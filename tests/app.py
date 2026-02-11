import streamlit as st
import os
import time
import uuid
from datetime import datetime, timedelta
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions, generate_container_sas, ContainerSasPermissions
from azure.ai.translation.document import DocumentTranslationClient, DocumentTranslationInput, TranslationTarget
from azure.core.credentials import AzureKeyCredential
import urllib.parse
import requests
import fitz # PyMuPDF for page count
import pandas as pd
import zipfile
import io

# Search Manager Import
from search_manager import AzureSearchManager

# Chat Manager Import  
from chat_manager_v2 import AzureOpenAIChatManager
from doc_intel_manager import DocumentIntelligenceManager
import excel_manager

# Authentication imports
from utils.auth_manager import AuthManager
from modules.login_page import render_login_page
from utils.chat_history_utils import load_history, save_history, get_session_title
import extra_streamlit_components as stx

# -----------------------------
# ì„¤ì • ë° ë¹„ë°€ ê´€ë¦¬
# -----------------------------
st.set_page_config(page_title="ì¸í…”ë¦¬ì „íŠ¸ ë‹¤íë¨¼íŠ¸", page_icon="ğŸ—ï¸", layout="wide")

# Custom CSS for larger tab labels and document list alignment
st.markdown("""
<style>
    /* Increase font size for tab labels */
    button[data-baseweb="tab"] {
        font-size: 20px !important;
    }
    button[data-baseweb="tab"] p {
        font-size: 20px !important;
        font-weight: 600 !important;
    }
    
    /* Document list - row alignment */
    [data-testid="stHorizontalBlock"] {
        display: flex !important;
        align-items: center !important;
        gap: 0.5rem !important;
        min-height: 42px !important;
    }
    
    /* Column layout - vertical centering */
    [data-testid="column"] {
        display: flex !important;
        flex-direction: column !important;
        justify-content: center !important;
    }
    
    /* All buttons - consistent height and sizing */
    .stButton button, .stLinkButton a {
        min-height: 38px !important;
        max-height: 38px !important;
        height: 38px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        padding: 0.25rem 0.75rem !important;
        white-space: nowrap !important;
        font-size: 1.1rem !important;
    }
    
    /* Popover button - same height */
    button[data-testid="baseButton-header"] {
        min-height: 38px !important;
        max-height: 38px !important;
        height: 38px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        padding: 0.25rem 0.75rem !important;
        font-size: 1.1rem !important;
    }
    
    /* Checkbox alignment */
    .stCheckbox {
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        min-height: 38px !important;
    }
    
    /* Markdown text alignment */
    .stMarkdown {
        display: flex !important;
        align-items: center !important;
        min-height: 38px !important;
    }
    
    /* Prevent wrapping in icon columns */
    [data-testid="column"] > div {
        white-space: nowrap !important;
    }
</style>
""", unsafe_allow_html=True)

def get_secret(key):
    if key in st.secrets:
        return st.secrets[key]
    return os.environ.get(key)

# í•„ìˆ˜ ìê²© ì¦ëª…
# 1. Storage
STORAGE_CONN_STR = get_secret("AZURE_STORAGE_CONNECTION_STRING")
CONTAINER_NAME = get_secret("AZURE_BLOB_CONTAINER_NAME") or "blob-leesunguk"

# 2. Translator
TRANSLATOR_KEY = get_secret("AZURE_TRANSLATOR_KEY")
TRANSLATOR_ENDPOINT = get_secret("AZURE_TRANSLATOR_ENDPOINT")

# 3. Search
SEARCH_ENDPOINT = get_secret("AZURE_SEARCH_ENDPOINT")
SEARCH_KEY = get_secret("AZURE_SEARCH_KEY")
SEARCH_INDEX_NAME = get_secret("AZURE_SEARCH_INDEX_NAME") or "pdf-search-index"
SEARCH_INDEXER_NAME = "pdf-indexer"
SEARCH_DATASOURCE_NAME = "blob-datasource"

# 4. Azure OpenAI
AZURE_OPENAI_ENDPOINT = get_secret("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_KEY = get_secret("AZURE_OPENAI_KEY")
AZURE_OPENAI_DEPLOYMENT = get_secret("AZURE_OPENAI_DEPLOYMENT") or get_secret("AZURE_OPENAI_DEPLOYMENT_NAME")
AZURE_OPENAI_API_VERSION = get_secret("AZURE_OPENAI_API_VERSION")

# 5. Document Intelligence
AZURE_DOC_INTEL_ENDPOINT = get_secret("AZURE_DOC_INTEL_ENDPOINT")
AZURE_DOC_INTEL_KEY = get_secret("AZURE_DOC_INTEL_KEY")

# -----------------------------
# Azure í´ë¼ì´ì–¸íŠ¸ í—¬í¼
# -----------------------------
def get_blob_service_client():
    if not STORAGE_CONN_STR:
        st.error("Azure Storage Connection Stringì´ ì„¤ì •ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
        st.stop()
    return BlobServiceClient.from_connection_string(STORAGE_CONN_STR)

def get_translation_client():
    if not TRANSLATOR_KEY or not TRANSLATOR_ENDPOINT:
        st.error("Azure Translator Key ë˜ëŠ” Endpointê°€ ì„¤ì •ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
        st.stop()
    return DocumentTranslationClient(TRANSLATOR_ENDPOINT, AzureKeyCredential(TRANSLATOR_KEY))

def get_search_manager():
    if not SEARCH_ENDPOINT or not SEARCH_KEY:
        st.error("Azure Search Endpoint ë˜ëŠ” Keyê°€ ì„¤ì •ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
        st.stop()
    return AzureSearchManager(SEARCH_ENDPOINT, SEARCH_KEY, SEARCH_INDEX_NAME)

def get_chat_manager():
    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_KEY:
        st.error("Azure OpenAI Endpoint ë˜ëŠ” Keyê°€ ì„¤ì •ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
        st.stop()
    return AzureOpenAIChatManager(
        AZURE_OPENAI_ENDPOINT, 
        AZURE_OPENAI_KEY, 
        AZURE_OPENAI_DEPLOYMENT, 
        AZURE_OPENAI_API_VERSION,
        get_search_manager(),
        STORAGE_CONN_STR,
        CONTAINER_NAME
    )

def get_doc_intel_manager():
    if not AZURE_DOC_INTEL_ENDPOINT or not AZURE_DOC_INTEL_KEY:
        st.error("Azure Document Intelligence Endpoint ë˜ëŠ” Keyê°€ ì„¤ì •ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
        st.stop()
    return DocumentIntelligenceManager(AZURE_DOC_INTEL_ENDPOINT, AZURE_DOC_INTEL_KEY)

def generate_sas_url(blob_service_client, container_name, blob_name=None, page=None, permission="r", expiry_hours=1, content_disposition=None):
    """
    Generates a SAS URL for a blob and wraps it in a web viewer (Google Docs/Office) if applicable.
    If blob_name is None, generates a Container SAS.
    """
    try:
        account_name = blob_service_client.account_name
        
        # Handle credential types
        if hasattr(blob_service_client.credential, 'account_key'):
            account_key = blob_service_client.credential.account_key
        else:
            account_key = blob_service_client.credential['account_key']
        
        start = datetime.utcnow() - timedelta(minutes=15)
        expiry = datetime.utcnow() + timedelta(hours=expiry_hours)
        
        if blob_name:
            # Clean blob name (remove page suffixes like " (p.1)")
            import re
            clean_name = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', blob_name).strip()
            
            # Determine content type
            import mimetypes
            content_type, _ = mimetypes.guess_type(clean_name)
            
            # Force PDF content type if extension matches (to ensure browser opens it)
            if clean_name.lower().endswith('.pdf'):
                content_type = "application/pdf"
                content_disposition = "inline"
            elif not content_type:
                content_type = "application/octet-stream"

            if content_disposition is None:
                content_disposition = "inline"

            sas_token = generate_blob_sas(
                account_name=account_name,
                container_name=container_name,
                blob_name=clean_name,
                account_key=account_key,
                permission=BlobSasPermissions(read=True),
                start=start,
                expiry=expiry,
                content_disposition=content_disposition,
                content_type=content_type
            )
            sas_url = f"https://{account_name}.blob.core.windows.net/{container_name}/{urllib.parse.quote(clean_name, safe='/')}?{sas_token}"
            
            lower_name = clean_name.lower()
            if lower_name.endswith(('.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls')):
                encoded_sas_url = urllib.parse.quote(sas_url)
                return f"https://view.officeapps.live.com/op/view.aspx?src={encoded_sas_url}"
            elif lower_name.endswith('.pdf'):
                # Use native browser viewer (better performance/reliability than Google Viewer)
                # encoded_sas_url = urllib.parse.quote(sas_url)
                # final_url = f"https://docs.google.com/viewer?url={encoded_sas_url}"
                
                # Direct SAS URL with content_disposition=inline opens in browser PDF viewer
                final_url = sas_url
                if page:
                    final_url += f"#page={page}"
                return final_url
            else:
                return sas_url
        else:
            # Container SAS
            sas_token = generate_container_sas(
                account_name=account_name,
                container_name=container_name,
                account_key=account_key,
                permission=ContainerSasPermissions(write=True, list=True, read=True, delete=True),
                start=start,
                expiry=expiry
            )
            return f"https://{account_name}.blob.core.windows.net/{container_name}?{sas_token}"
            
    except Exception as e:
        st.error(f"SAS URL ìƒì„± ì¤‘ ì˜¤ë¥˜ ë°œìƒ ({blob_name}): {e}")
        return "#"

# -----------------------------
# Progress Management (Resume Capability)
# -----------------------------
import json

TEMP_DIR = ".temp_analysis"
if not os.path.exists(TEMP_DIR):
    os.makedirs(TEMP_DIR)

def get_progress_file_path(safe_filename):
    return os.path.join(TEMP_DIR, f"{safe_filename}_progress.json")

def save_progress(safe_filename, page_chunks, total_pages):
    """Save intermediate analysis progress to disk"""
    try:
        filepath = get_progress_file_path(safe_filename)
        data = {
            "safe_filename": safe_filename,
            "total_pages": total_pages,
            "page_chunks": page_chunks,
            "last_updated": datetime.utcnow().isoformat()
        }
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        # print(f"DEBUG: Progress saved for {safe_filename} ({len(page_chunks)} chunks)")
    except Exception as e:
        print(f"Error saving progress: {e}")

def load_progress(safe_filename):
    """Load intermediate analysis progress from disk"""
    try:
        filepath = get_progress_file_path(safe_filename)
        if os.path.exists(filepath):
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
            print(f"DEBUG: Progress loaded for {safe_filename} ({len(data.get('page_chunks', []))} chunks)")
            return data
    except Exception as e:
        print(f"Error loading progress: {e}")
    return None

def delete_progress(safe_filename):
    """Delete progress file after successful completion"""
    try:
        filepath = get_progress_file_path(safe_filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            print(f"DEBUG: Progress file deleted for {safe_filename}")
    except Exception as e:
        print(f"Error deleting progress file: {e}")

# File Persistence for Resume
FILES_DIR = os.path.join(TEMP_DIR, "files")
if not os.path.exists(FILES_DIR):
    os.makedirs(FILES_DIR)

def save_uploaded_file_temp(uploaded_file, safe_filename):
    """Save uploaded file to temp dir for resume capability"""
    try:
        filepath = os.path.join(FILES_DIR, safe_filename)
        with open(filepath, "wb") as f:
            f.write(uploaded_file.getbuffer())
        return filepath
    except Exception as e:
        print(f"Error saving temp file: {e}")
        return None

def get_temp_file_path(safe_filename):
    return os.path.join(FILES_DIR, safe_filename)

def delete_temp_file(safe_filename):
    """Delete temp file after completion"""
    try:
        filepath = os.path.join(FILES_DIR, safe_filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            print(f"DEBUG: Temp file deleted for {safe_filename}")
    except Exception as e:
        print(f"Error deleting temp file: {e}")

class LocalFile:
    """Mock Streamlit UploadedFile for local files"""
    def __init__(self, path, name, type="application/pdf"):
        self.path = path
        self.name = name
        self.type = type
        self.size = os.path.getsize(path)
        self._file = open(path, "rb")

    def read(self, size=-1):
        return self._file.read(size)

    def seek(self, offset, whence=0):
        return self._file.seek(offset, whence)

    def tell(self):
        return self._file.tell()
        
    def getbuffer(self):
        # Return bytes
        self.seek(0)
        return self.read()

    def close(self):
        self._file.close()

def is_drm_protected(uploaded_file):
    """
    Check if the uploaded file is DRM protected or encrypted.
    Returns True if protected, False otherwise.
    """
    try:
        file_type = uploaded_file.name.split('.')[-1].lower()
        
        # 1. PDF Check
        if file_type == 'pdf':
            try:
                # Read file stream
                bytes_data = uploaded_file.getvalue()
                with fitz.open(stream=bytes_data, filetype="pdf") as doc:
                    if doc.is_encrypted:
                        return True
            except Exception as e:
                print(f"PDF DRM Check Error: {e}")
                # If we can't open it with fitz, it might be corrupted or heavily encrypted
                return True 

        # 2. Office Files (docx, pptx, xlsx) Check
        # Modern Office files are Zip archives. If they are encrypted/DRM'd, 
        # they often become OLE CF (Compound File) binaries or non-zip streams.
        elif file_type in ['docx', 'pptx', 'xlsx']:
            try:
                bytes_data = uploaded_file.getvalue()
                # Check if it is a valid zip file
                if not zipfile.is_zipfile(io.BytesIO(bytes_data)):
                    # Not a zip -> Likely Encrypted/DRM (OLE format)
                    return True
                
                # Optional: Try to open it to be sure
                with zipfile.ZipFile(io.BytesIO(bytes_data)) as zf:
                    # Check for standard OOXML structure (e.g., [Content_Types].xml)
                    if '[Content_Types].xml' not in zf.namelist():
                        return True
            except Exception as e:
                print(f"Office DRM Check Error: {e}")
                return True # Assume protected if we can't parse structure
                
        return False
    except Exception as e:
        print(f"General DRM Check Error: {e}")
        return False



# -----------------------------
# UI êµ¬ì„±
# -----------------------------


# ì§€ì› ì–¸ì–´ ëª©ë¡ ê°€ì ¸ì˜¤ê¸° (API)
@st.cache_data
def get_supported_languages():
    try:
        url = "https://api.cognitive.microsofttranslator.com/languages?api-version=3.0&scope=translation"
        # Accept-Language í—¤ë”ë¥¼ 'ko'ë¡œ ì„¤ì •í•˜ì—¬ ì–¸ì–´ ì´ë¦„ì„ í•œêµ­ì–´ë¡œ ë°›ìŒ
        headers = {"Accept-Language": "ko"}
        response = requests.get(url, headers=headers, timeout=5)
        response.raise_for_status()
        data = response.json()
        
        languages = {}
        for code, info in data['translation'].items():
            # "í•œêµ­ì–´ ì´ë¦„ (ì›ì–´ ì´ë¦„)" í˜•ì‹ìœ¼ë¡œ í‘œì‹œ (ì˜ˆ: ì˜ì–´ (English))
            label = f"{info['name']} ({info['nativeName']})"
            languages[label] = code
        return languages
    except requests.exceptions.SSLError:
        # ë¡œì»¬ í™˜ê²½(ì‚¬ë‚´ë§) ë“±ì—ì„œ SSL ì¸ì¦ì„œ ì˜¤ë¥˜ ë°œìƒ ì‹œ verify=Falseë¡œ ì¬ì‹œë„
        try:
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
            response = requests.get(url, headers=headers, verify=False, timeout=5)
            response.raise_for_status()
            data = response.json()
            languages = {}
            for code, info in data['translation'].items():
                label = f"{info['name']} ({info['nativeName']})"
                languages[label] = code
            return languages
        except Exception as e:
            print(f"SSL Bypass retry failed: {e}")
            # ì‹¤íŒ¨ ì‹œ ì•„ë˜ ê¸°ë³¸ ì–¸ì–´ ì œê³µìœ¼ë¡œ ë„˜ì–´ê°

    except Exception as e:
        print(f"ì–¸ì–´ ëª©ë¡ ê°€ì ¸ì˜¤ê¸° ì‹¤íŒ¨ (API): {e}")
        # UIì— ì—ëŸ¬ë¥¼ í‘œì‹œí•˜ì§€ ì•Šê³  ì½˜ì†”ì—ë§Œ ë‚¨ê¹€
    
    # ì‹¤íŒ¨ ì‹œ ê¸°ë³¸ ì–¸ì–´ ì œê³µ (í™•ì¥ëœ ëª©ë¡)
    return {
        "í•œêµ­ì–´ (Korean)": "ko", 
        "ì˜ì–´ (English)": "en",
        "ì¼ë³¸ì–´ (Japanese)": "ja",
        "ì¤‘êµ­ì–´ ê°„ì²´ (Chinese Simplified)": "zh-Hans",
        "ì¤‘êµ­ì–´ ë²ˆì²´ (Chinese Traditional)": "zh-Hant",
        "í”„ë‘ìŠ¤ì–´ (French)": "fr",
        "ë…ì¼ì–´ (German)": "de",
        "ìŠ¤í˜ì¸ì–´ (Spanish)": "es",
        "ëŸ¬ì‹œì•„ì–´ (Russian)": "ru",
        "ë² íŠ¸ë‚¨ì–´ (Vietnamese)": "vi"
    }

LANGUAGES = get_supported_languages()

# ì–¸ì–´ ì½”ë“œë³„ íŒŒì¼ ì ‘ë¯¸ì‚¬ ë§¤í•‘ (ê¸°ë³¸ì ìœ¼ë¡œ ëŒ€ë¬¸ì ì½”ë“œë¥¼ ì‚¬ìš©í•˜ë˜, ì¼ë¶€ ì»¤ìŠ¤í…€ ê°€ëŠ¥)
# ì—¬ê¸°ì„œëŠ” ìë™ ìƒì„± ë¡œì§ì„ ì‚¬ìš©í•˜ë¯€ë¡œ ë³„ë„ ë”•ì…”ë„ˆë¦¬ ë¶ˆí•„ìš”, 
# ë‹¤ë§Œ ì¤‘êµ­ì–´ ë“± íŠ¹ìˆ˜ ì¼€ì´ìŠ¤ë¥¼ ìœ„í•´ ë‚¨ê²¨ë‘˜ ìˆ˜ ìˆìŒ.
LANG_SUFFIX_OVERRIDE = {
    "zh-Hans": "CN",
    "zh-Hant": "TW",
}

# Initialize session state for page navigation
if "page" not in st.session_state:
    st.session_state.page = "í™ˆ"

def change_page(page_name):
    st.session_state.page = page_name

# Initialize AuthManager
auth_manager = AuthManager(STORAGE_CONN_STR)

# Initialize Cookie Manager
# Initialize Cookie Manager
cookie_manager = stx.CookieManager(key="auth_cookie_manager")

# Initialize login state
if 'is_logged_in' not in st.session_state:
    st.session_state.is_logged_in = False

# Check for existing session cookie (Auto-login)
# Check for existing session cookie (Auto-login)
if not st.session_state.is_logged_in and not st.session_state.get('just_logged_out', False):
    try:
        # Improved robust cookie retrieval (Retry mechanism)
        auth_email = None
        
        # Method 1: Direct get with retries (wait for component to sync)
        # extra_streamlit_components sometimes needs a moment to load cookies from frontend
        for i in range(5):
            auth_email = cookie_manager.get(cookie="auth_email")
            if auth_email:
                break
            time.sleep(0.1) 
            
        # Method 2: Fallback to get_all() if direct get failed
        if not auth_email:
            cookies = cookie_manager.get_all()
            if cookies and isinstance(cookies, dict):
                auth_email = cookies.get("auth_email")
        
        if auth_email:
            # Validate email exists in auth_manager
            user = auth_manager.get_user_by_email(auth_email)
            if user:
                st.session_state.is_logged_in = True
                st.session_state.user_info = user
                st.toast(f"ìë™ ë¡œê·¸ì¸ë˜ì—ˆìŠµë‹ˆë‹¤: {user.get('name')}")
    except Exception as e:
        print(f"Cookie check failed: {e}")

# Check if user is logged in
if not st.session_state.is_logged_in:
    render_login_page(auth_manager, cookie_manager)
    st.stop()

# User is logged in - get their info
user_info = st.session_state.get('user_info', {})
user_role = user_info.get('role', 'guest')
user_perms = user_info.get('permissions', [])

def get_user_folder_name(user_info):
    """Get sanitized user folder name"""
    if not user_info:
        return "guest"
    # Use name but fallback to ID if empty
    name = user_info.get('name', user_info.get('id', 'guest'))
    return name.strip()

user_folder = get_user_folder_name(user_info)

# Define role-based menu permissions (Fallback / Admin)
ALL_MENUS = ["í™ˆ", "ë²ˆì—­í•˜ê¸°", "íŒŒì¼ ë³´ê´€í•¨", "ë¬¼ì–´ë³´ë©´ ë‹µí•˜ëŠ” ë¬¸ì„œ AI", "ë„ë©´/ìŠ¤í™ ë¹„êµ", "ì—‘ì…€ë°ì´í„° ìë™ì¶”ì¶œ", "ì‚¬ì§„ëŒ€ì§€ ìë™ì‘ì„±", "ì‘ì—…ê³„íš ë° íˆ¬ì…ë¹„ ìë™ì‘ì„±", "ê´€ë¦¬ì ì„¤ì •", "ì‚¬ìš©ì ì„¤ì •", "ë””ë²„ê·¸ (Debug)"]
GUEST_MENUS = ["í™ˆ", "ì‚¬ìš©ì ì„¤ì •"]

if user_role == 'admin':
    available_menus = ALL_MENUS
else:
    # Use assigned permissions, ensuring mandatory menus are present
    available_menus = user_perms if user_perms else GUEST_MENUS
    
    # Filter out button permissions (btn:download, btn:edit, etc.) from menu list
    available_menus = [menu for menu in available_menus if not menu.startswith('btn:')]
    
    # Map old menu names to new names (Migration fix)
    available_menus = [
        "ë„ë©´/ìŠ¤í™ ë¹„êµ" if menu == "ë„ë©´/ìŠ¤í™ ë¶„ì„" else 
        "ë¬¼ì–´ë³´ë©´ ë‹µí•˜ëŠ” ë¬¸ì„œ AI" if menu in ["ê²€ìƒ‰ & AI ì±„íŒ…", "ë¬¸ì„œ ì—…ë¡œë“œ & AI ì±„íŒ…"] else menu 
        for menu in available_menus
    ]
    # Ensure "í™ˆ" and "ì‚¬ìš©ì ì„¤ì •" are always available
    if "í™ˆ" not in available_menus:
        available_menus.insert(0, "í™ˆ")
    if "ì‚¬ìš©ì ì„¤ì •" not in available_menus:
        available_menus.append("ì‚¬ìš©ì ì„¤ì •")
    
    # Remove "ê´€ë¦¬ì ì„¤ì •" if somehow present for non-admins
    if "ê´€ë¦¬ì ì„¤ì •" in available_menus:
        available_menus.remove("ê´€ë¦¬ì ì„¤ì •")

with st.sidebar:
    # User profile
    st.markdown(f"### ğŸ‘¤ {user_info.get('name', 'User')}")
    st.caption(f"**{user_info.get('email', '')}**")
    st.caption(f"ê¶Œí•œ: {user_role.upper()}")
    
    # Debug: Show permissions and menus
    # st.caption(f"Perms: {user_perms}")
    # st.caption(f"Menus: {available_menus}")
    
    # --- Persistent Error Display ---
    if "drm_error_message" in st.session_state and st.session_state.drm_error_message:
        st.error(st.session_state.drm_error_message)
        # Clear it after showing
        del st.session_state.drm_error_message
    
    if st.button("ğŸšª ë¡œê·¸ì•„ì›ƒ", key="logout_btn", use_container_width=True):
        st.session_state.is_logged_in = False
        st.session_state.user_info = None
        st.session_state.just_logged_out = True # Prevent immediate auto-login
        # Delete cookie
        cookie_manager.delete("auth_email")
        st.rerun()
    
    st.divider()
    
    st.header("ë©”ë‰´")
    # Filter menu based on user role
    menu = st.radio("ì´ë™", available_menus, key="page")
    
    st.divider()
    
    if menu == "ë²ˆì—­í•˜ê¸°":
        st.header("ì„¤ì •")
        # í•œêµ­ì–´ë¥¼ ê¸°ë³¸ê°’ìœ¼ë¡œ ì°¾ê¸°
        default_index = 0
        lang_labels = list(LANGUAGES.keys())
        for i, label in enumerate(lang_labels):
            if "Korean" in label or "í•œêµ­ì–´" in label:
                default_index = i
                break
                
        target_lang_label = st.selectbox("ëª©í‘œ ì–¸ì–´ ì„ íƒ", lang_labels, index=default_index)
        target_lang_code = LANGUAGES[target_lang_label]
        st.info(f"ì„ íƒëœ ëª©í‘œ ì–¸ì–´: {target_lang_code}")

    # ìê²© ì¦ëª… ìƒíƒœ í™•ì¸
    if STORAGE_CONN_STR and TRANSLATOR_KEY and SEARCH_KEY:
        st.success("âœ… Azure ìê²© ì¦ëª… í™•ì¸ë¨")
    else:
        st.warning("âš ï¸ ì¼ë¶€ Azure ìê²© ì¦ëª…ì´ ëˆ„ë½ë˜ì—ˆìŠµë‹ˆë‹¤.")

# Common Header for non-Home pages - Removed to allow custom placement
# if menu != "í™ˆ":
#     st.title(menu)


if menu == "í™ˆ":
    # Use the new home_chat module with function calling support
    from home_chat import render_home_chat
    chat_manager = get_chat_manager()
    render_home_chat(chat_manager)
    
if menu == "ë²ˆì—­í•˜ê¸°":
    _, col_main, _ = st.columns([0.1, 0.8, 0.1])
    with col_main:
        if "translate_uploader_key" not in st.session_state:
            st.session_state.translate_uploader_key = 0

        uploaded_file = st.file_uploader("ë²ˆì—­í•  ë¬¸ì„œ ì—…ë¡œë“œ (PPTX, PDF, DOCX, XLSX ë“±)", type=["pptx", "pdf", "docx", "xlsx"], key=f"translate_{st.session_state.translate_uploader_key}")

        # ì´ì „ ë²ˆì—­ ê²°ê³¼ê°€ ìˆìœ¼ë©´ í‘œì‹œ
        if "last_translation_result" in st.session_state:
            result = st.session_state.last_translation_result
            st.success("âœ… ë²ˆì—­ì´ ì™„ë£Œë˜ì—ˆìŠµë‹ˆë‹¤!")
            st.markdown(f"[{result['file_name']} ë‹¤ìš´ë¡œë“œ]({result['url']})", unsafe_allow_html=True)
            
            # ê²°ê³¼ë¥¼ ì§€ìš°ê³  ì‹¶ì„ ìˆ˜ ìˆìœ¼ë¯€ë¡œ ë‹«ê¸° ë²„íŠ¼ ì œê³µ (ì„ íƒ ì‚¬í•­)
            if st.button("ê²°ê³¼ ë‹«ê¸°"):
                del st.session_state.last_translation_result
                st.rerun()

        if uploaded_file:
            if is_drm_protected(uploaded_file):
                st.session_state.drm_error_message = "â›” DRMìœ¼ë¡œ ë³´í˜¸ëœ íŒŒì¼(ì•”í˜¸í™”ëœ íŒŒì¼)ì€ ë²ˆì—­í•  ìˆ˜ ì—†ìŠµë‹ˆë‹¤. íŒŒì¼ ëª©ë¡ì—ì„œ ì œê±°ë˜ì—ˆìŠµë‹ˆë‹¤."
                st.session_state.translate_uploader_key += 1
                st.rerun()

        if st.button("ë²ˆì—­ ì‹œì‘", type="primary", disabled=not uploaded_file):
            if not uploaded_file:
                st.error("íŒŒì¼ì„ ì—…ë¡œë“œí•´ì£¼ì„¸ìš”.")
            else:
                with st.spinner("Azure Blobì— íŒŒì¼ ì—…ë¡œë“œ ì¤‘..."):
                    try:
                        blob_service_client = get_blob_service_client()
                        container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                        
                        # ì»¨í…Œì´ë„ˆ ì ‘ê·¼ ê¶Œí•œ í™•ì¸
                        try:
                            if not container_client.exists():
                                container_client.create_container()
                        except Exception as e:
                            if "AuthenticationFailed" in str(e):
                                st.error("ğŸš¨ ì¸ì¦ ì‹¤íŒ¨: Azure Storage Keyê°€ ì˜¬ë°”ë¥´ì§€ ì•ŠìŠµë‹ˆë‹¤. Secrets ì„¤ì •ì„ í™•ì¸í•´ì£¼ì„¸ìš”.")
                                st.stop()
                            else:
                                raise e

                        # íŒŒì¼ëª… ìœ ë‹ˆí¬í•˜ê²Œ ì²˜ë¦¬ (UUID ì œê±°, ë®ì–´ì“°ê¸° í—ˆìš©)
                        # file_uuid = str(uuid.uuid4())[:8] 
                        original_filename = uploaded_file.name
                        input_blob_name = f"{user_folder}/original/{original_filename}"
                        
                        # ì—…ë¡œë“œ
                        blob_client = container_client.get_blob_client(input_blob_name)
                        blob_client.upload_blob(uploaded_file, overwrite=True)
                        
                        st.success("ì—…ë¡œë“œ ì™„ë£Œ! ë²ˆì—­ ìš”ì²­ ì¤‘...")
                        
                        # SAS ìƒì„±
                        source_url = generate_sas_url(blob_service_client, CONTAINER_NAME, input_blob_name)
                        
                        # Target URL ì„¤ì •
                        target_base_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}"
                        # Target URLì€ ì»¨í…Œì´ë„ˆ ë˜ëŠ” í´ë” ê²½ë¡œì—¬ì•¼ í•¨ (íŒŒì¼ ê²½ë¡œ ë¶ˆê°€)
                        # ì‚¬ìš©ìë³„ translated í´ë”ë¡œ ì„¤ì •
                        # URL ì¸ì½”ë”© í•„ìš”
                        encoded_user_folder = urllib.parse.quote(user_folder)
                        target_output_url = f"{target_base_url}/{encoded_user_folder}/translated/?{generate_sas_url(blob_service_client, CONTAINER_NAME).split('?')[1]}"
                        
                    except Exception as e:
                        st.error(f"ì—…ë¡œë“œ/SAS ìƒì„± ì‹¤íŒ¨: {e}")
                        st.stop()

                with st.spinner("ë²ˆì—­ ì‘ì—… ìš”ì²­ ë° ëŒ€ê¸° ì¤‘..."):
                    try:
                        client = get_translation_client()
                        
                        poller = client.begin_translation(
                            inputs=[
                                DocumentTranslationInput(
                                    source_url=source_url,
                                    storage_type="File",
                                    targets=[
                                        TranslationTarget(
                                            target_url=target_output_url,
                                            language=target_lang_code
                                        )
                                    ]
                                )
                            ]
                        )
                        
                        result = poller.result()
                        
                        for doc in result:
                            if doc.status == "Succeeded":
                                st.success(f"ë²ˆì—­ ì™„ë£Œ! (ìƒíƒœ: {doc.status})")
                            else:
                                st.error(f"ë¬¸ì„œ ë²ˆì—­ ì‹¤íŒ¨! (ìƒíƒœ: {doc.status})")
                                if doc.error:
                                    st.error(f"ì—ëŸ¬ ì½”ë“œ: {doc.error.code}, ë©”ì‹œì§€: {doc.error.message}")
                        
                        # ê²°ê³¼ íŒŒì¼ ì°¾ê¸°
                        time.sleep(2)
                        # UUID í´ë”ê°€ ì—†ìœ¼ë¯€ë¡œ translated í´ë” ì „ì²´ì—ì„œ í•´ë‹¹ íŒŒì¼ëª… ê²€ìƒ‰
                        output_prefix_search = f"{user_folder}/translated/"
                        output_blobs = list(container_client.list_blobs(name_starts_with=output_prefix_search))
                        
                        # ë°©ê¸ˆ ë²ˆì—­ëœ íŒŒì¼ ì°¾ê¸° (íŒŒì¼ëª… ë§¤ì¹­)
                        # Azure ë²ˆì—­ì€ ì›ë³¸ íŒŒì¼ëª…ì„ ìœ ì§€í•˜ê±°ë‚˜ ì–¸ì–´ ì½”ë“œë¥¼ ë¶™ì„
                        target_blobs = []
                        for blob in output_blobs:
                            if original_filename in blob.name:
                                target_blobs.append(blob)
                        
                        if not target_blobs:
                            st.warning(f"ê²°ê³¼ íŒŒì¼ì„ ì°¾ëŠ” ì¤‘ì…ë‹ˆë‹¤... (ê²½ë¡œ: {output_prefix_search})")
                            # Fallback: list all to debug
                            # all_output = list(container_client.list_blobs(name_starts_with=output_prefix_search))
                            # debug_msg = "\n".join([b.name for b in all_output[:10]])
                            # st.error(f"ê²°ê³¼ íŒŒì¼ì„ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤.\ní˜„ì¬ í´ë” íŒŒì¼ ëª©ë¡:\n{debug_msg}")
                        else:
                            st.subheader("ë‹¤ìš´ë¡œë“œ")
                            for blob in target_blobs:
                                blob_name = blob.name
                                file_name = blob_name.split("/")[-1]
                                
                                # íŒŒì¼ëª…ì— ì–¸ì–´ ì ‘ë¯¸ì‚¬ ì¶”ê°€ (Rename)
                                suffix = LANG_SUFFIX_OVERRIDE.get(target_lang_code, target_lang_code.upper())
                                name_part, ext_part = os.path.splitext(file_name)
                                
                                # ì´ë¯¸ ì ‘ë¯¸ì‚¬ê°€ ìˆëŠ”ì§€ í™•ì¸ (í˜¹ì‹œ ëª¨ë¥¼ ì¤‘ë³µ ë°©ì§€)
                                if not name_part.endswith(f"_{suffix}"):
                                    new_file_name = f"{name_part}_{suffix}{ext_part}"
                                    new_blob_name = f"{user_folder}/translated/{new_file_name}"
                                    
                                    try:
                                        # Rename: Copy to new name -> Delete old
                                        source_blob = container_client.get_blob_client(blob_name)
                                        dest_blob = container_client.get_blob_client(new_blob_name)
                                        
                                        source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob_name)
                                        dest_blob.start_copy_from_url(source_sas)
                                        
                                        # Wait for copy
                                        for _ in range(10):
                                            props = dest_blob.get_blob_properties()
                                            if props.copy.status == "success":
                                                break
                                            time.sleep(0.2)
                                            
                                        source_blob.delete_blob()
                                        
                                        # Update variables for download link
                                        blob_name = new_blob_name
                                        file_name = new_file_name
                                        st.toast(f"íŒŒì¼ëª… ë³€ê²½ë¨: {file_name}")
                                        
                                    except Exception as e:
                                        st.warning(f"íŒŒì¼ëª… ë³€ê²½ ì‹¤íŒ¨ (ê¸°ë³¸ ì´ë¦„ìœ¼ë¡œ ìœ ì§€): {e}")

                                # PPTX í°íŠ¸ ë³€ê²½ (Times New Roman)
                                if file_name.lower().endswith(".pptx"):
                                    try:
                                        from pptx import Presentation
                                        
                                        # ì„ì‹œ íŒŒì¼ë¡œ ë‹¤ìš´ë¡œë“œ
                                        temp_pptx = f"temp_{original_filename}"
                                        blob_client_temp = container_client.get_blob_client(blob_name)
                                        with open(temp_pptx, "wb") as f:
                                            data = blob_client_temp.download_blob().readall()
                                            f.write(data)
                                        
                                        # í°íŠ¸ ë³€ê²½ ë¡œì§
                                        prs = Presentation(temp_pptx)
                                        font_name = "Times New Roman"
                                        
                                        def change_font(shapes):
                                            for shape in shapes:
                                                if shape.has_text_frame:
                                                    for paragraph in shape.text_frame.paragraphs:
                                                        for run in paragraph.runs:
                                                            run.font.name = font_name
                                                
                                                if shape.has_table:
                                                    for row in shape.table.rows:
                                                        for cell in row.cells:
                                                            if cell.text_frame:
                                                                for paragraph in cell.text_frame.paragraphs:
                                                                    for run in paragraph.runs:
                                                                        run.font.name = font_name
                                                
                                                if shape.shape_type == 6: # Group
                                                    change_font(shape.shapes)

                                        for slide in prs.slides:
                                            change_font(slide.shapes)
                                        
                                        prs.save(temp_pptx)
                                        
                                        # ë‹¤ì‹œ ì—…ë¡œë“œ (ë®ì–´ì“°ê¸°)
                                        with open(temp_pptx, "rb") as f:
                                            blob_client_temp.upload_blob(f, overwrite=True)
                                        
                                        os.remove(temp_pptx)
                                        st.toast("PPTX í°íŠ¸ ë³€ê²½ ì™„ë£Œ (Times New Roman)")
                                        
                                    except Exception as e:
                                        st.warning(f"PPTX í°íŠ¸ ë³€ê²½ ì‹¤íŒ¨: {e}")

                                download_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob_name)
                                st.markdown(f"[{file_name} ë‹¤ìš´ë¡œë“œ]({download_sas})", unsafe_allow_html=True)
                                
                                # ê²°ê³¼ ì„¸ì…˜ì— ì €ì¥
                                st.session_state.last_translation_result = {
                                    "file_name": file_name,
                                    "url": download_sas
                                }
                                
                        # ì„±ê³µì ìœ¼ë¡œ ì™„ë£Œë˜ë©´ ì—…ë¡œë” ì´ˆê¸°í™” (í‚¤ ë³€ê²½)
                        st.session_state.translate_uploader_key += 1
                        time.sleep(1) # ì ì‹œ ëŒ€ê¸°
                        st.rerun()
                                
                    except Exception as e:
                        st.error(f"ë²ˆì—­ ìš”ì²­ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")

elif menu == "íŒŒì¼ ë³´ê´€í•¨":
    _, col_main, _ = st.columns([0.1, 0.8, 0.1])
    with col_main:
        # st.subheader("ğŸ“‚ í´ë¼ìš°ë“œ íŒŒì¼ ë³´ê´€í•¨") - Removed to avoid duplication
        
        st.divider()
        
        if st.button("ğŸ”„ ëª©ë¡ ìƒˆë¡œê³ ì¹¨"):
            st.rerun()
            
        try:
            blob_service_client = get_blob_service_client()
            container_client = blob_service_client.get_container_client(CONTAINER_NAME)
            
            # íƒ­ìœ¼ë¡œ Input/Output êµ¬ë¶„
            tab1, tab2 = st.tabs(["ì›ë³¸ ë¬¸ì„œ (Input)", "ë²ˆì—­ëœ ë¬¸ì„œ (Output)"])
            
            def render_file_list(prefixes, tab_name):
                all_blobs = []
                for prefix in prefixes:
                    blobs = list(container_client.list_blobs(name_starts_with=prefix))
                    all_blobs.extend(blobs)
                
                # ì¤‘ë³µ ì œê±° (í˜¹ì‹œ ëª¨ë¥¼ ê²½ìš° ëŒ€ë¹„)
                unique_blobs = {b.name: b for b in all_blobs}.values()
                blobs = list(unique_blobs)
                blobs.sort(key=lambda x: x.creation_time, reverse=True)
                
                if not blobs:
                    st.info(f"{tab_name}ì— íŒŒì¼ì´ ì—†ìŠµë‹ˆë‹¤.")
                    return

                for i, blob in enumerate(blobs):
                    file_name = blob.name.split("/")[-1]
                    creation_time = blob.creation_time.strftime('%Y-%m-%d %H:%M')
                    
                    # í´ë” ê²½ë¡œ í‘œì‹œ (ê´€ë¦¬ì í¸ì˜)
                    folder_path = "/".join(blob.name.split("/")[:-1])
                    
                    with st.container():
                        col1, col2, col3 = st.columns([6, 2, 2])
                        
                        with col1:
                            sas_url = generate_sas_url(blob_service_client, CONTAINER_NAME, blob.name)
                            st.markdown(f"**[{file_name}]({sas_url})**")
                            st.caption(f"ğŸ“‚ {folder_path} | ğŸ“… {creation_time} | ğŸ“¦ {blob.size / 1024:.1f} KB")
                        
                        with col2:
                            # ìˆ˜ì • (ì´ë¦„ ë³€ê²½)
                            with st.popover("ìˆ˜ì •"):
                                new_name = st.text_input("ìƒˆ íŒŒì¼ëª…", value=file_name, key=f"rename_{i}_{blob.name}")
                                if st.button("ì´ë¦„ ë³€ê²½", key=f"btn_rename_{i}_{blob.name}"):
                                    try:
                                        # ìƒˆ ê²½ë¡œ ìƒì„± (ê¸°ì¡´ í´ë” êµ¬ì¡° ìœ ì§€)
                                        path_parts = blob.name.split("/")
                                        folder = "/".join(path_parts[:-1])
                                        new_blob_name = f"{folder}/{new_name}"
                                        
                                        # ë³µì‚¬ (Renameì€ Copy + Delete)
                                        source_blob = container_client.get_blob_client(blob.name)
                                        dest_blob = container_client.get_blob_client(new_blob_name)
                                        
                                        # SAS URL for Copy Source
                                        source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob.name)
                                        
                                        dest_blob.start_copy_from_url(source_sas)
                                        
                                        # ë³µì‚¬ ì™„ë£Œ ëŒ€ê¸° (ê°„ë‹¨í•œ í´ë§)
                                        for _ in range(10):
                                            props = dest_blob.get_blob_properties()
                                            if props.copy.status == "success":
                                                break
                                            time.sleep(0.5)
                                        
                                        # ì›ë³¸ ì‚­ì œ
                                        source_blob.delete_blob()
                                        st.success("ì´ë¦„ ë³€ê²½ ì™„ë£Œ!")
                                        time.sleep(1)
                                        st.rerun()
                                    except Exception as e:
                                        st.error(f"ì´ë¦„ ë³€ê²½ ì‹¤íŒ¨: {e}")

                        with col3:
                            # ì‚­ì œ
                            if st.button("ì‚­ì œ", key=f"del_{prefix}_{i}", type="secondary"):
                                try:
                                    container_client.delete_blob(blob.name)
                                    st.success("ì‚­ì œë˜ì—ˆìŠµë‹ˆë‹¤.")
                                    time.sleep(1)
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"ì‚­ì œ ì‹¤íŒ¨: {e}")
                        
                        st.divider()

            with tab1:
                input_prefixes = [f"{user_folder}/documents/"]
                if user_role == 'admin':
                    input_prefixes.extend(["input/", "gulflng/"])
                render_file_list(input_prefixes, "ë‚´ ë¬¸ì„œ (Documents)")
                
            with tab2:
                output_prefixes = [f"{user_folder}/translated/"]
                if user_role == 'admin':
                    output_prefixes.extend(["output/"])
                render_file_list(output_prefixes, "ë²ˆì—­ëœ ë¬¸ì„œ")
                    
        except Exception as e:
            st.error(f"íŒŒì¼ ëª©ë¡ì„ ë¶ˆëŸ¬ì˜¤ëŠ” ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")

elif menu == "ë¬¼ì–´ë³´ë©´ ë‹µí•˜ëŠ” ë¬¸ì„œ AI":
    from utils.chat_history_utils import load_history, save_history, get_session_title
    SEARCH_HISTORY_FILE = "search_chat_history.json"

    # Initialize Session State for Search History
    if "search_chat_history_data" not in st.session_state:
        st.session_state.search_chat_history_data = load_history(SEARCH_HISTORY_FILE)
    
    if "current_search_session_id" not in st.session_state:
        new_id = str(uuid.uuid4())
        st.session_state.current_search_session_id = new_id
        st.session_state.search_chat_history_data[new_id] = {
            "title": "ìƒˆë¡œìš´ ëŒ€í™”",
            "timestamp": datetime.now().isoformat(),
            "messages": []
        }
        st.session_state.chat_messages = [] # This maps to the current session messages

    # Layout: Spacer L (25%) | Main Content (50%) | Spacer R (10%) | History Sidebar (15%)
    col_spacer_l, col_main, col_spacer_r, col_history = st.columns([0.25, 0.5, 0.1, 0.15])
    
    # Custom CSS for Sidebar Styling (Same as Home)
    st.markdown("""
    <style>
    /* Target the fourth column (History Sidebar) */
    [data-testid="stHorizontalBlock"] > [data-testid="column"]:nth-of-type(4) {
        background-color: #1E1E1E;
        border-left: 1px solid #333;
        padding: 1rem;
        border-radius: 10px;
    }
    [data-testid="column"]:nth-of-type(4) button {
        text-align: left;
    }
    </style>
    """, unsafe_allow_html=True)

    # --- Right Sidebar (History) ---
    with col_history:
        st.markdown("### ì±„íŒ… ê¸°ë¡")
        
        if st.button("â• ìƒˆ ì±„íŒ…", key="new_search_chat", use_container_width=True):
            new_id = str(uuid.uuid4())
            st.session_state.current_search_session_id = new_id
            st.session_state.search_chat_history_data[new_id] = {
                "title": "ìƒˆë¡œìš´ ëŒ€í™”",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.chat_messages = []
            st.rerun()
            
        st.markdown("---")
        
        sorted_sessions = sorted(
            st.session_state.search_chat_history_data.items(),
            key=lambda x: x[1].get("timestamp", ""),
            reverse=True
        )
        
        for session_id, session_data in sorted_sessions:
            title = session_data.get("title", "ëŒ€í™”")
            if session_id == st.session_state.current_search_session_id:
                if st.button(f"ğŸ“‚ {title}", key=f"search_hist_{session_id}", use_container_width=True, type="primary"):
                    pass
            else:
                if st.button(f"ğŸ“„ {title}", key=f"search_hist_{session_id}", use_container_width=True):
                    st.session_state.current_search_session_id = session_id
                    st.session_state.chat_messages = session_data.get("messages", [])
                    st.rerun()
        
        if st.button("ğŸ—‘ï¸ ê¸°ë¡ ì‚­ì œ", key="del_search_hist", use_container_width=True):
            st.session_state.search_chat_history_data = {}
            save_history(SEARCH_HISTORY_FILE, {})
            new_id = str(uuid.uuid4())
            st.session_state.current_search_session_id = new_id
            st.session_state.search_chat_history_data[new_id] = {
                "title": "ìƒˆë¡œìš´ ëŒ€í™”",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.chat_messages = []
            st.rerun()

    with col_main:
        st.title("ë¬¼ì–´ë³´ë©´ ë‹µí•˜ëŠ” ë¬¸ì„œ AI")
        # Tabs for Search and Chat to preserve state
        tab1, tab2, tab3 = st.tabs(["ğŸ“¤ ë¬¸ì„œ ë“±ë¡", "ğŸ” í‚¤ì›Œë“œ ê²€ìƒ‰", "ğŸ¤– AI ì§ˆì˜ì‘ë‹µ"])
        
        with tab1:
            # File Uploader (Simplified)
            if "doc_search_uploader_key" not in st.session_state:
                st.session_state.doc_search_uploader_key = 0
                
            doc_upload = st.file_uploader("ë¬¸ì„œë¥¼ ë“±ë¡í•˜ë©´ ê²€ìƒ‰ê³¼ ì§ˆì˜ì‘ë‹µì´ ê°€ëŠ¥í•©ë‹ˆë‹¤.", type=['pdf', 'docx', 'txt', 'pptx'], key=f"doc_search_upload_{st.session_state.doc_search_uploader_key}")
            
            if doc_upload:
                if is_drm_protected(doc_upload):
                    st.session_state.drm_error_message = "â›” DRMìœ¼ë¡œ ë³´í˜¸ëœ íŒŒì¼(ì•”í˜¸í™”ëœ íŒŒì¼)ì€ ì—…ë¡œë“œí•  ìˆ˜ ì—†ìŠµë‹ˆë‹¤. ë³´ì•ˆì„ í•´ì œí•œ í›„ ë‹¤ì‹œ ì‹œë„í•´ì£¼ì„¸ìš”."
                    st.session_state.doc_search_uploader_key += 1
                    st.rerun()

            if doc_upload and st.button("ì—…ë¡œë“œ", key="btn_doc_upload"):
                try:
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    
                    # Upload to {user_folder}/documents/ (Flat structure)
                    blob_name = f"{user_folder}/documents/{doc_upload.name}"
                    blob_client = container_client.get_blob_client(blob_name)
                    blob_client.upload_blob(doc_upload, overwrite=True)
                    st.success(f"'{doc_upload.name}' ì—…ë¡œë“œ ì™„ë£Œ! (ì¸ë±ì‹±ì— ì‹œê°„ì´ ê±¸ë¦´ ìˆ˜ ìˆìŠµë‹ˆë‹¤)")
                except Exception as e:
                    st.error(f"ì—…ë¡œë“œ ì‹¤íŒ¨: {e}")
            
            st.divider()
            
            # Indexed Document List
            st.markdown("### ğŸ—‚ï¸ ë“±ë¡ ë¬¸ì„œ ëª©ë¡")
            
            try:
                search_manager = get_search_manager()
                
                # Construct prefix URL for filtering
                account_name = get_blob_service_client().account_name
                encoded_user_folder = urllib.parse.quote(user_folder)
                prefix_url = f"https://{account_name}.blob.core.windows.net/{CONTAINER_NAME}/{encoded_user_folder}/"
                
                # Filter logic
                if user_role == 'admin':
                    filter_expr = None
                else:
                    # Workaround: Use range query instead of startswith if startswith is not supported
                    # prefix_url ends with '/' (ASCII 47). Next char is '0' (ASCII 48).
                    # So we want path >= prefix_url AND path < prefix_url_with_next_char
                    # Actually, let's just use the next char logic safely.
                    # Or just try startswith again with debug? No, let's try the range.
                    # prefix_url = .../
                    # upper_bound = ...0
                    upper_bound = prefix_url[:-1] + '0'
                    filter_expr = f"metadata_storage_path ge '{prefix_url}' and metadata_storage_path lt '{upper_bound}'"
                
                # Debug
                # st.write(f"Debug Filter: {filter_expr}")
                
                # Search all documents (*)
                results = search_manager.search("*", filter_expr=filter_expr, top=1000)
                
                # Filter out .json files first
                filtered_results = []
                for res in results:
                    file_name = res.get('metadata_storage_name', 'Unknown')
                    if not file_name.lower().endswith('.json'):
                        filtered_results.append(res)
                
                if not filtered_results:
                    st.info("ì¸ë±ì‹±ëœ ë¬¸ì„œê°€ ì—†ìŠµë‹ˆë‹¤.")
                else:
                    st.write(f"ì´ {len(filtered_results)}ê°œ ë¬¸ì„œê°€ ë“±ë¡ë˜ì–´ ìˆìŠµë‹ˆë‹¤. (ê²€ìƒ‰ê³¼ ì§ˆì˜ê°€ ê°€ëŠ¥í•©ë‹ˆë‹¤).")
                    
                    # Display as a table
                    doc_data = []
                    for res in filtered_results:
                        file_name = res.get('metadata_storage_name', 'Unknown')
                        size = res.get('metadata_storage_size', 0)
                        last_modified = res.get('metadata_storage_last_modified', '')
                        path = res.get('metadata_storage_path', '')
                        
                        # Convert size to MB
                        size_mb = f"{int(size) / (1024 * 1024):.2f} MB"
                        
                        # Format date
                        try:
                            dt = datetime.fromisoformat(last_modified.replace('Z', '+00:00'))
                            date_str = dt.strftime("%Y-%m-%d %H:%M")
                        except:
                            date_str = last_modified
                            
                        doc_data.append({
                            "Name": file_name,
                            "Size": size_mb,
                            "Last Modified": date_str,
                            "path": path # Hidden for logic
                        })
                    
                    # Use dataframe for better display
                    import pandas as pd
                    df = pd.DataFrame(doc_data)
                    
                    # Display table with selection (optional, maybe just list)
                    # For now, just a clean dataframe display
                    st.dataframe(
                        df[["Name", "Size", "Last Modified"]],
                        use_container_width=True,
                        hide_index=True
                    )
            except Exception as e:
                st.error(f"ë¬¸ì„œ ëª©ë¡ì„ ë¶ˆëŸ¬ì˜¤ëŠ” ì¤‘ ì˜¤ë¥˜ê°€ ë°œìƒí–ˆìŠµë‹ˆë‹¤: {e}")
                if 'filter_expr' in locals():
                    st.code(filter_expr)
        
        with tab2:
            # -----------------------------
            # ê²€ìƒ‰ ì˜µì…˜
            # -----------------------------
            with st.expander("âš™ï¸ ê³ ê¸‰ ê²€ìƒ‰ ì˜µì…˜ (RAG ì„¤ì •)", expanded=False):
                c1, c2 = st.columns(2)
                with c1:
                    search_use_semantic = st.checkbox("ì‹œë§¨í‹± ë­ì»¤ ì‚¬ìš©", value=True, key="search_use_semantic", help="ì˜ë¯¸ ê¸°ë°˜ ê²€ìƒ‰ì„ ì‚¬ìš©í•˜ì—¬ ì •í™•ë„ë¥¼ ë†’ì…ë‹ˆë‹¤.")
                with c2:
                    search_mode_opt = st.radio("ê²€ìƒ‰ ëª¨ë“œ", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="search_mode_opt", help="any: í‚¤ì›Œë“œ ì¤‘ í•˜ë‚˜ë¼ë„ í¬í•¨ë˜ë©´ ê²€ìƒ‰ (ì¶”ì²œ)")
                    search_mode = "all" if "all" in search_mode_opt else "any"

            # Display Chat History (Shared with AI Chat)
            for msg in st.session_state.chat_messages:
                with st.chat_message(msg["role"]):
                    if msg["role"] == "user":
                        st.markdown(msg["content"])
                    else:
                        # Assistant message (Results)
                        if "results" in msg:
                            results = msg["results"]
                            if not results:
                                st.info("ê²€ìƒ‰ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤.")
                            else:
                                st.success(f"ì´ {len(results)}ê°œì˜ ë¬¸ì„œë¥¼ ì°¾ì•˜ìŠµë‹ˆë‹¤.")
                                for result in results:
                                    with st.container():
                                        file_name = result.get('metadata_storage_name', 'Unknown File')
                                        path = result.get('metadata_storage_path', '')
                                        
                                        # Highlights
                                        highlights = result.get('@search.highlights')
                                        if highlights:
                                            snippets = []
                                            if 'content' in highlights:
                                                snippets.extend(highlights['content'])
                                            if 'content_exact' in highlights:
                                                snippets.extend(highlights['content_exact'])
                                            unique_snippets = list(set(snippets))[:3]
                                            content_snippet = " ... ".join(unique_snippets)
                                        else:
                                            content_snippet = result.get('content', '')[:300] + "..."
                                        
                                        # Text Cleaning Logic (Restored & Improved)
                                        import re
                                        def clean_text(text):
                                            # 1. Escape Markdown special characters except HTML tags we want to keep (like <mark>)
                                            text = text.replace('~', '\\~')
                                            
                                            # 2. Handle HTML Tables (convert to Markdown-ish for display)
                                            text = re.sub(r'</td>', ' | ', text, flags=re.IGNORECASE)
                                            text = re.sub(r'</th>', ' | ', text, flags=re.IGNORECASE)
                                            text = re.sub(r'</tr>', '\n', text, flags=re.IGNORECASE)
                                            
                                            # 3. If text contains pipes (|), it might be a table. Preserve structure.
                                            if "|" in text:
                                                # Remove other HTML tags except <mark>
                                                text = re.sub(r'<(?!/?mark\b)[^>]+>', '', text)
                                                text = re.sub(r'^\s*(\|[\s\|]*)+\s*$', '', text, flags=re.MULTILINE)
                                                text = re.sub(r'\n\s*\n', '\n', text)
                                                return text.strip()

                                            # 4. Remove other HTML tags except <mark>
                                            text = re.sub(r'<(?!/?mark\b)[^>]+>', '', text)
                                            
                                            # 5. Replace single newlines with space
                                            cleaned = re.sub(r'(?<!\.)\n(?!\n)', ' ', text)
                                            cleaned = re.sub(r' +', ' ', cleaned)
                                            return cleaned.strip()
                                            
                                        st.markdown(f"### ğŸ“„ {file_name}")
                                        st.markdown(f"> {clean_text(content_snippet)}", unsafe_allow_html=True)
                                        
                                        # Generate SAS link
                                        try:
                                            blob_service_client = get_blob_service_client()
                                            from urllib.parse import unquote
                                            
                                            if "https://direct_fetch/" in path:
                                                blob_path = unquote(path.replace("https://direct_fetch/", "").split('#')[0])
                                            elif CONTAINER_NAME in path:
                                                blob_path = unquote(path.split(f"/{CONTAINER_NAME}/")[1].split('#')[0])
                                            else:
                                                blob_path = path
                                            
                                            import mimetypes
                                            content_type, _ = mimetypes.guess_type(file_name)
                                            
                                            sas_token = generate_blob_sas(
                                                account_name=blob_service_client.account_name,
                                                container_name=CONTAINER_NAME,
                                                blob_name=blob_path,
                                                account_key=blob_service_client.credential.account_key,
                                                permission=BlobSasPermissions(read=True),
                                                expiry=datetime.utcnow() + timedelta(hours=1),
                                                content_disposition="inline",
                                                content_type=content_type
                                            )
                                            sas_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_path)}?{sas_token}"
                                            
                                            lower_name = file_name.lower()
                                            if lower_name.endswith(('.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls')):
                                                final_url = f"https://view.officeapps.live.com/op/view.aspx?src={urllib.parse.quote(sas_url)}"
                                                link_text = "ğŸ“„ ì›¹ì—ì„œ ë³´ê¸° (Office Viewer)"
                                            elif lower_name.endswith('.pdf'):
                                                final_url = f"https://docs.google.com/viewer?url={urllib.parse.quote(sas_url)}"
                                                link_text = "ğŸ“„ ì›¹ì—ì„œ ë³´ê¸° (PDF Viewer)"
                                            else:
                                                final_url = sas_url
                                                link_text = "ğŸ“„ ë¬¸ì„œ ì—´ê¸° (ìƒˆ íƒ­)"
                                                
                                            st.markdown(f'<a href="{final_url}" target="_blank">{link_text}</a>', unsafe_allow_html=True)
                                        except Exception as e:
                                            st.caption(f"ë§í¬ ìƒì„± ì‹¤íŒ¨: {e}")
                                        st.divider()
                        else:
                            st.markdown(msg["content"])

            # Chat Input for Search
            if query := st.chat_input("ê²€ìƒ‰í•  í‚¤ì›Œë“œë¥¼ ì…ë ¥í•˜ì„¸ìš”...", key="keyword_chat_input"):
                st.session_state.chat_messages.append({"role": "user", "content": query})
                
                with st.spinner("ê²€ìƒ‰ ì¤‘..."):
                    try:
                        search_manager = get_search_manager()
                        account_name = get_blob_service_client().account_name
                        encoded_user_folder = urllib.parse.quote(user_folder)
                        prefix_url = f"https://{account_name}.blob.core.windows.net/{CONTAINER_NAME}/{encoded_user_folder}/"
                        
                        if user_role == 'admin':
                            filter_expr = None
                        else:
                            upper_bound = prefix_url[:-1] + '0'
                            filter_expr = f"metadata_storage_path ge '{prefix_url}' and metadata_storage_path lt '{upper_bound}'"
                        
                        results = search_manager.search(query, filter_expr=filter_expr, use_semantic_ranker=search_use_semantic, search_mode=search_mode)
                        
                        # Filter out .json files
                        filtered_results = [res for res in results if not res.get('metadata_storage_name', '').lower().endswith('.json')]
                        
                        st.session_state.chat_messages.append({
                            "role": "assistant",
                            "content": f"'{query}'ì— ëŒ€í•œ ê²€ìƒ‰ ê²°ê³¼ì…ë‹ˆë‹¤.",
                            "results": filtered_results
                        })
                        
                        # --- Auto-Save History ---
                        current_id = st.session_state.current_search_session_id
                        current_title = st.session_state.search_chat_history_data[current_id]["title"]
                        if current_title == "ìƒˆë¡œìš´ ëŒ€í™”" and len(st.session_state.chat_messages) > 0:
                            new_title = get_session_title(st.session_state.chat_messages)
                            st.session_state.search_chat_history_data[current_id]["title"] = new_title
                        
                        st.session_state.search_chat_history_data[current_id]["messages"] = st.session_state.chat_messages
                        st.session_state.search_chat_history_data[current_id]["timestamp"] = datetime.now().isoformat()
                        save_history(SEARCH_HISTORY_FILE, st.session_state.search_chat_history_data)
                        st.rerun()
                    except Exception as e:
                        st.error(f"ê²€ìƒ‰ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")

        with tab3:
            # -----------------------------
            # ê²€ìƒ‰ ì˜µì…˜ (Chat Tab) - Moved to Top
            # -----------------------------
            with st.expander("âš™ï¸ ê³ ê¸‰ ê²€ìƒ‰ ì˜µì…˜ (RAG ì„¤ì •)", expanded=False):
                c1, c2 = st.columns(2)
                with c1:
                    chat_use_semantic = st.checkbox("ì‹œë§¨í‹± ë­ì»¤ ì‚¬ìš©", value=True, key="chat_use_semantic", help="ì˜ë¯¸ ê¸°ë°˜ ê²€ìƒ‰ì„ ì‚¬ìš©í•˜ì—¬ ì •í™•ë„ë¥¼ ë†’ì…ë‹ˆë‹¤.")
                with c2:
                    chat_search_mode_opt = st.radio("ê²€ìƒ‰ ëª¨ë“œ", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="chat_search_mode", help="any: í‚¤ì›Œë“œ ì¤‘ í•˜ë‚˜ë¼ë„ í¬í•¨ë˜ë©´ ê²€ìƒ‰ (ì¶”ì²œ)")
                    chat_search_mode = "all" if "all" in chat_search_mode_opt else "any"
            
            # Initialize chat history in session state
            if "chat_messages" not in st.session_state:
                st.session_state.chat_messages = []
            
            # Display chat messages
            for message in st.session_state.chat_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
                    
                    # Display citations if present
                    if "citations" in message and message["citations"]:
                        st.markdown("---")
                        st.caption("ğŸ“š **ì°¸ì¡° ë¬¸ì„œ:**")
                        for i, citation in enumerate(message["citations"], 1):
                            filepath = citation.get('filepath', 'Unknown')
                            # Use pre-generated final_url if available, otherwise generate one
                            display_url = citation.get('final_url')
                            if not display_url:
                                try:
                                    blob_service_client = get_blob_service_client()
                                    display_url = generate_sas_url(
                                        blob_service_client, 
                                        CONTAINER_NAME, 
                                        filepath, 
                                        page=citation.get('page')
                                    )
                                except:
                                    display_url = "#"
                            
                            st.markdown(f"{i}. [{filepath}]({display_url})")
            
            # -----------------------------
            # ê²€ìƒ‰ ì˜µì…˜ (Chat Tab) - Moved to Top
            # -----------------------------
            # st.write("")
            # with st.expander("âš™ï¸ ê³ ê¸‰ ê²€ìƒ‰ ì˜µì…˜ (RAG ì„¤ì •)", expanded=False):
            #     c1, c2 = st.columns(2)
            #     with c1:
            #         chat_use_semantic = st.checkbox("ì‹œë§¨í‹± ë­ì»¤ ì‚¬ìš©", value=False, key="chat_use_semantic", help="ì˜ë¯¸ ê¸°ë°˜ ê²€ìƒ‰ì„ ì‚¬ìš©í•˜ì—¬ ì •í™•ë„ë¥¼ ë†’ì…ë‹ˆë‹¤.")
            #     with c2:
            #         chat_search_mode_opt = st.radio("ê²€ìƒ‰ ëª¨ë“œ", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="chat_search_mode", help="any: í‚¤ì›Œë“œ ì¤‘ í•˜ë‚˜ë¼ë„ í¬í•¨ë˜ë©´ ê²€ìƒ‰ (ì¶”ì²œ)")
            #         chat_search_mode = "all" if "all" in chat_search_mode_opt else "any"

            # Chat input
            if prompt := st.chat_input("ì§ˆë¬¸ì„ ì…ë ¥í•˜ì„¸ìš” (ì˜ˆ: 10-P-101Aì˜ ì‚¬ì–‘ì€ ë¬´ì—‡ì¸ê°€ìš”?)", key="search_chat_input"):
                # Add user message to chat history
                st.session_state.chat_messages.append({"role": "user", "content": prompt})
                
                # Display user message
                with st.chat_message("user"):
                    st.markdown(prompt)
                
                # Get AI response
                with st.chat_message("assistant"):
                    with st.spinner("ë‹µë³€ ìƒì„± ì¤‘..."):
                        try:
                            chat_manager = get_chat_manager()
                            
                            # Prepare conversation history (exclude citations from history)
                            conversation_history = [
                                {"role": msg["role"], "content": msg["content"]}
                                for msg in st.session_state.chat_messages[:-1]  # Exclude the just-added user message
                            ]
                            
                            # Pass the selected search options to the chat manager
                            response_text, citations, context, final_filter, search_results = chat_manager.get_chat_response(
                                prompt, 
                                conversation_history, 
                                search_mode=chat_search_mode, 
                                use_semantic_ranker=chat_use_semantic,
                                filter_expr=None,
                                user_folder=user_folder, # Pass Name-based folder (matches Blob/Index path)
                                is_admin=(user_role == 'admin')
                            )
                            
                            # ---------------------------------------------------------
                            # CRITICAL: Linkify Inline Citations & Escape Tildes
                            # ---------------------------------------------------------
                            
                            # 1. Pre-generate Web Viewer URLs for all citations
                            citation_links = {}
                            processed_citations = []
                            
                            if citations:
                                for cit in citations:
                                    filepath = cit.get('filepath', 'Unknown')
                                    # CRITICAL: Clean filepath from page suffixes like " (p.1)" or " (p.1) (p.1)"
                                    import re
                                    clean_filepath = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', filepath).strip()
                                    
                                    page = cit.get('page')
                                    url = cit.get('url', '')
                                    
                                    # Generate Web Viewer URL
                                    try:
                                        blob_service_client = get_blob_service_client()
                                        final_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            clean_filepath, 
                                            page=page
                                        )
                                    except Exception as e:
                                        st.error(f"URL ìƒì„± ì‹¤íŒ¨ ({clean_filepath}): {e}")
                                        final_url = "#"
                                    
                                    cit['final_url'] = final_url
                                    processed_citations.append(cit)
                                    
                                    filename = os.path.basename(filepath)
                                    if page:
                                        citation_links[(filename, str(page))] = final_url
                            
                            # 2. Replace text citations with Markdown links (Support both [] and ())
                            if response_text:
                                import re
                                # Pattern to match [filename: p.1] or (filename: p.1)
                                # Improved: Allow parentheses in filenames (common in EPC drawings)
                                # CRITICAL FIX: Exclude pipe (|) to prevent crossing table boundaries
                                pattern = r'[\[\(]([^\[\]|]+?:\s*p\.?\s*(\d+))[\]\)]'
                                
                                def replace_citation(match):
                                    content = match.group(1).strip()
                                    
                                    # Remove "ë¬¸ì„œëª…:" prefix if present (Common in Korean LLM outputs)
                                    content = re.sub(r'^ë¬¸ì„œëª…\s*:\s*', '', content)
                                    
                                    # Split by last colon to separate filename and page
                                    if ':' in content:
                                        fname = content.rsplit(':', 1)[0].strip()
                                        p_num = match.group(2)
                                    else:
                                        return match.group(0)
                                        
                                    original_text = match.group(0)
                                    
                                    # Try to find matching citation with fuzzy logic
                                    target_url = None
                                    
                                    # 1. Clean LLM filename for comparison
                                    clean_llm = re.sub(r'\.pdf$', '', fname.lower().strip())
                                    
                                    for (k_fname, k_page), url in citation_links.items():
                                        # 2. Clean known filename (remove .pdf and (p.N) suffixes)
                                        clean_known = re.sub(r'\.pdf$', '', k_fname.lower().strip())
                                        clean_known = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', clean_known).strip()
                                        
                                        # CRITICAL FIX: Skip empty filenames to prevent false positive matches
                                        if not clean_known:
                                            continue

                                        # 3. Match page and filename (fuzzy)
                                        if str(k_page) == str(p_num):
                                            # Exact match or containment (handling "ë¬¸ì„œëª…: " residue if regex failed)
                                            if clean_llm == clean_known or clean_llm in clean_known or clean_known in clean_llm:
                                                target_url = url
                                                # CRITICAL FIX: Capture the matched filename to replace text
                                                matched_filename = k_fname
                                                break
                                    
                                    if target_url:
                                        # Use Markdown link for table compatibility
                                        # Escape parentheses in URL to avoid breaking Markdown link
                                        safe_url = target_url.replace('(', '%28').replace(')', '%29')
                                        
                                        # CRITICAL FIX: Replace original text (e.g. "Same Document") with actual filename
                                        # Reconstruct text: (Filename: p.N)
                                        if matched_filename:
                                            new_text = f"({matched_filename}: p.{p_num})"
                                            return f"**[{new_text}]({safe_url})**"
                                        
                                        return f"**[{original_text}]({safe_url})**"
                                    
                                    return original_text

                                # DEBUG: Show raw response before linkification
                                st.code(response_text, language="markdown")
                                
                                response_text = re.sub(pattern, replace_citation, response_text)

                                # 3. Escape tildes
                                response_text = response_text.replace('~', '\\~')
                        
                            # Display response
                            st.markdown(response_text, unsafe_allow_html=True)
                            
                            # Display citations
                            if processed_citations:
                                st.markdown("---")
                                st.caption("ğŸ“š **ì°¸ì¡° ë¬¸ì„œ:**")
                                for i, citation in enumerate(processed_citations, 1):
                                    filepath = citation.get('filepath', 'Unknown')
                                    filename = os.path.basename(filepath)
                                    display_url = citation.get('final_url', '#')
                                    
                                    link_text = "ë¬¸ì„œ ë³´ê¸°"
                                    if "docs.google.com" in display_url: link_text = "PDF Viewer"
                                    elif "view.officeapps" in display_url: link_text = "Office Viewer"
                                    
                                    st.markdown(f"{i}. [{filename}]({display_url}) - {link_text}")
                            
                            # Debug: Show Citation Links (Hidden by default)
                            # with st.expander("ğŸ” ë§í¬ ë””ë²„ê¹… (Debug Links)", expanded=False):
                            #     st.write("Citation Links Keys:", list(citation_links.keys()))
                            #     st.write("Processed Citations:", processed_citations)
                            
                            # Add assistant response to chat history
                            st.session_state.chat_messages.append({
                                "role": "assistant",
                                "content": response_text,
                                "citations": citations,
                                "context": context,
                                "debug_filter": final_filter
                            })
                            
                            # Debug: Show Context
                            with st.expander("ğŸ” ê²€ìƒ‰ëœ ì»¨í…ìŠ¤íŠ¸ í™•ì¸ (Debug Context)", expanded=False):
                                if final_filter:
                                    st.caption(f"**OData Filter:** `{final_filter}`")
                                st.text_area("LLMì—ê²Œ ì „ë‹¬ëœ ì›ë¬¸ ë°ì´í„°", value=context, height=300)
                            
                            # --- Auto-Save History ---
                            current_id = st.session_state.current_search_session_id
                            current_title = st.session_state.search_chat_history_data[current_id]["title"]
                            if current_title == "ìƒˆë¡œìš´ ëŒ€í™”" and len(st.session_state.chat_messages) > 0:
                                new_title = get_session_title(st.session_state.chat_messages)
                                st.session_state.search_chat_history_data[current_id]["title"] = new_title
                            
                            st.session_state.search_chat_history_data[current_id]["messages"] = st.session_state.chat_messages
                            st.session_state.search_chat_history_data[current_id]["timestamp"] = datetime.now().isoformat()
                            save_history(SEARCH_HISTORY_FILE, st.session_state.search_chat_history_data)
                            st.rerun()

                        except Exception as e:
                            st.error(f"ì˜¤ë¥˜ê°€ ë°œìƒí–ˆìŠµë‹ˆë‹¤: {str(e)}")
        

elif menu == "ë„ë©´/ìŠ¤í™ ë¹„êµ":
    DRAWING_HISTORY_FILE = "drawing_chat_history.json"
    
    # Initialize Session State for Drawing History
    if "drawing_chat_history_data" not in st.session_state:
        st.session_state.drawing_chat_history_data = load_history(DRAWING_HISTORY_FILE)
    
    if "current_drawing_session_id" not in st.session_state:
        new_id = str(uuid.uuid4())
        st.session_state.current_drawing_session_id = new_id
        st.session_state.drawing_chat_history_data[new_id] = {
            "title": "ìƒˆë¡œìš´ ëŒ€í™”",
            "timestamp": datetime.now().isoformat(),
            "messages": []
        }
        st.session_state.rag_chat_messages = []

    # Layout: Spacer L (25%) | Main Content (50%) | Spacer R (10%) | History Sidebar (15%)
    col_spacer_l, col_main, col_spacer_r, col_history = st.columns([0.25, 0.5, 0.1, 0.15])
    
    # Custom CSS for Sidebar Styling (Same as Home)
    st.markdown("""
    <style>
    /* Target the fourth column (History Sidebar) */
    [data-testid="stHorizontalBlock"] > [data-testid="column"]:nth-of-type(4) {
        background-color: #1E1E1E;
        border-left: 1px solid #333;
        padding: 1rem;
        border-radius: 10px;
    }
    [data-testid="column"]:nth-of-type(4) button {
        text-align: left;
    }
    </style>
    """, unsafe_allow_html=True)

    # --- Right Sidebar (History) ---
    with col_history:
        st.markdown("### ì±„íŒ… ê¸°ë¡")
        
        if st.button("â• ìƒˆ ì±„íŒ…", key="new_drawing_chat", use_container_width=True):
            new_id = str(uuid.uuid4())
            st.session_state.current_drawing_session_id = new_id
            st.session_state.drawing_chat_history_data[new_id] = {
                "title": "ìƒˆë¡œìš´ ëŒ€í™”",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.rag_chat_messages = []
            st.rerun()
            
        st.markdown("---")
        
        sorted_sessions = sorted(
            st.session_state.drawing_chat_history_data.items(),
            key=lambda x: x[1].get("timestamp", ""),
            reverse=True
        )
        
        for session_id, session_data in sorted_sessions:
            title = session_data.get("title", "ëŒ€í™”")
            if session_id == st.session_state.current_drawing_session_id:
                if st.button(f"ğŸ“‚ {title}", key=f"drawing_hist_{session_id}", use_container_width=True, type="primary"):
                    pass
            else:
                if st.button(f"ğŸ“„ {title}", key=f"drawing_hist_{session_id}", use_container_width=True):
                    st.session_state.current_drawing_session_id = session_id
                    st.session_state.rag_chat_messages = session_data.get("messages", [])
                    st.rerun()
        
        if st.button("ğŸ—‘ï¸ ê¸°ë¡ ì‚­ì œ", key="del_drawing_hist", use_container_width=True):
            st.session_state.drawing_chat_history_data = {}
            save_history(DRAWING_HISTORY_FILE, {})
            new_id = str(uuid.uuid4())
            st.session_state.current_drawing_session_id = new_id
            st.session_state.drawing_chat_history_data[new_id] = {
                "title": "ìƒˆë¡œìš´ ëŒ€í™”",
                "timestamp": datetime.now().isoformat(),
                "messages": []
            }
            st.session_state.rag_chat_messages = []
            st.rerun()

    with col_main:
        st.title("ë„ë©´/ìŠ¤í™ ë¹„êµ")
    
        tab1, tab2 = st.tabs(["ğŸ“¤ ë¬¸ì„œ ì—…ë¡œë“œ", "ğŸ’¬ AIë¶„ì„"])
    
        with tab1:
        
            if "drawing_uploader_key" not in st.session_state:
                st.session_state.drawing_uploader_key = 0
            
            # High Resolution OCR Toggle
            use_high_res = st.toggle("ê³ í•´ìƒë„ OCR ì ìš© (ë„ë©´ ë¯¸ì„¸ ê¸€ì ì¶”ì¶œìš©)", value=False, help="ë³µì¡í•œ ë„ë©´ì˜ ì‘ì€ ê¸€ì”¨ë¥¼ ë” ì •í™•í•˜ê²Œ ì½ìŠµë‹ˆë‹¤. ë¶„ì„ ì‹œê°„ì´ ë” ì˜¤ë˜ ê±¸ë¦´ ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
        

        
            # --- RESUME UI SECTION ---
            # Check for interrupted sessions
            import glob
            resumable_files = []
            if os.path.exists(FILES_DIR):
                for filepath in glob.glob(os.path.join(FILES_DIR, "*")):
                    filename = os.path.basename(filepath)
                    # Check if progress exists
                    if load_progress(filename):
                        resumable_files.append(filename)
            
            files_to_process = []
            
            if resumable_files:
                st.warning(f"âš ï¸ ì¤‘ë‹¨ëœ ë¶„ì„ ì‘ì—…ì´ {len(resumable_files)}ê±´ ë°œê²¬ë˜ì—ˆìŠµë‹ˆë‹¤.")
                for r_file in resumable_files:
                    progress = load_progress(r_file)
                    processed = progress.get('processed_pages', 0)
                    total = progress.get('total_pages', '?')
                    last_updated = progress.get('last_updated', 'Unknown')
                    
                    # Format timestamp
                    try:
                        dt = datetime.fromisoformat(last_updated)
                        time_str = dt.strftime("%Y-%m-%d %H:%M:%S")
                    except:
                        time_str = last_updated
                    
                    with st.expander(f"ğŸ“„ {r_file} ({processed}/{total} í˜ì´ì§€ ì™„ë£Œ) - {time_str}", expanded=True):
                        col_res1, col_res2 = st.columns(2)
                        if col_res1.button(f"â–¶ï¸ ì´ì–´ì„œ ë¶„ì„í•˜ê¸° (Resume)", key=f"resume_{r_file}"):
                            # Create mock file object
                            local_path = get_temp_file_path(r_file)
                            if os.path.exists(local_path):
                                mock_file = LocalFile(local_path, r_file)
                                files_to_process.append(mock_file)
                                st.session_state.is_resuming = True
                            else:
                                st.error("ì›ë³¸ ì„ì‹œ íŒŒì¼ì„ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤.")
                        
                        if col_res2.button(f"ğŸ—‘ï¸ ì·¨ì†Œ ë° ì‚­ì œ (Discard)", key=f"discard_{r_file}"):
                            delete_progress(r_file)
                            delete_temp_file(r_file)
                            st.rerun()

            uploaded_files = st.file_uploader("PDF ë„ë©´, ìŠ¤í™, ì‚¬ì–‘ì„œ ë“±ì„ ì—…ë¡œë“œí•˜ì„¸ìš”", accept_multiple_files=True, type=['pdf', 'png', 'jpg', 'jpeg', 'tiff', 'bmp'], key=f"drawing_{st.session_state.drawing_uploader_key}")
        
            if uploaded_files or files_to_process:
                if "analysis_status" not in st.session_state:
                    st.session_state.analysis_status = {}
                
                # If resuming, we skip the "Start" button check or auto-click it
                start_analysis = False
                if files_to_process:
                    start_analysis = True
                    target_files = files_to_process
                elif uploaded_files:
                    # Immediate DRM Check
                    drm_files = [f.name for f in uploaded_files if is_drm_protected(f)]
                    if drm_files:
                        st.session_state.drm_error_message = f"â›” ë‹¤ìŒ íŒŒì¼ë“¤ì€ DRMìœ¼ë¡œ ë³´í˜¸ë˜ì–´ ìˆì–´ ì—…ë¡œë“œí•  ìˆ˜ ì—†ìŠµë‹ˆë‹¤: {', '.join(drm_files)}. íŒŒì¼ ëª©ë¡ì´ ì´ˆê¸°í™”ë˜ì—ˆìŠµë‹ˆë‹¤."
                        st.session_state.drawing_uploader_key += 1
                        st.rerun()

                    if st.button("ì—…ë¡œë“œ ë° ë¶„ì„ ì‹œì‘"):
                        start_analysis = True
                        target_files = uploaded_files
                
                if start_analysis:
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    doc_intel_manager = get_doc_intel_manager()
                    search_manager = get_search_manager()
                    
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    total_files = len(target_files)
                    
                    for idx, file in enumerate(target_files):
                        try:
                            # Normalize filename to NFC (to match search query logic)
                            import unicodedata
                            safe_filename = unicodedata.normalize('NFC', file.name)
                            
                            # Save to temp dir for resume capability (only if it's a fresh upload)
                            if not isinstance(file, LocalFile):
                                save_uploaded_file_temp(file, safe_filename)
                            
                            # Initialize status
                            st.session_state.analysis_status[safe_filename] = {
                                "status": "Extracting",
                                "total_pages": 0,
                                "processed_pages": 0,
                                "chunks": {},
                                "error": None
                            }
                            
                            status_text.text(f"ì²˜ë¦¬ ì¤‘ ({idx+1}/{total_files}): {safe_filename}")
                            
                            blob_path = f"{user_folder}/drawings/{safe_filename}"
                            # 2. Upload to Azure Blob Storage
                            status_text.text(f"ì—…ë¡œë“œ ì¤‘ ({idx+1}/{total_files}): {file.name}...")
                            blob_client = blob_service_client.get_blob_client(container=CONTAINER_NAME, blob=blob_path)
                            
                            # CRITICAL: Reset file pointer to ensure full upload
                            file.seek(0)
                            blob_client.upload_blob(file, overwrite=True)
                            
                            # Verify upload size
                            props = blob_client.get_blob_properties()
                            if props.size != file.size:
                                st.error(f"âš ï¸ íŒŒì¼ ì—…ë¡œë“œ í¬ê¸° ë¶ˆì¼ì¹˜! (ì›ë³¸: {file.size}, ì—…ë¡œë“œë¨: {props.size})")
                            else:
                                print(f"DEBUG: Upload verified. Size: {props.size} bytes")

                            # Generate SAS Token for Document Intelligence access
                            sas_token = generate_blob_sas(
                                account_name=blob_service_client.account_name,
                                container_name=CONTAINER_NAME,
                                blob_name=blob_path,
                                account_key=blob_service_client.credential.account_key,
                                permission=BlobSasPermissions(read=True),
                                expiry=datetime.utcnow() + timedelta(hours=1)
                            )
                            blob_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_path)}?{sas_token}"
                            
                            # 3. Analyze with Document Intelligence (Chunked)
                            file.seek(0)
                            pdf_data = file.read()
                            doc = fitz.open(stream=pdf_data, filetype="pdf")
                            total_pages = doc.page_count
                            file.seek(0)
                            
                            status_text.text(f"ë¶„ì„ ì¤€ë¹„ ì¤‘ ({idx+1}/{total_files}): {file.name} (ì´ {total_pages} í˜ì´ì§€)")
                            
                            st.session_state.analysis_status[safe_filename]["total_pages"] = total_pages
                            
                            chunk_size = 50
                            page_chunks = []
                            
                            # RESUME LOGIC: Check for existing progress
                            existing_progress = load_progress(safe_filename)
                            processed_ranges = set()
                            
                            if existing_progress and existing_progress.get('total_pages') == total_pages:
                                st.info(f"ğŸ”„ ì´ì „ ë¶„ì„ ì§„í–‰ ìƒí™©ì„ ë°œê²¬í–ˆìŠµë‹ˆë‹¤. ({len(existing_progress.get('page_chunks', []))} í˜ì´ì§€ ì™„ë£Œë¨) ì´ì–´ì„œ ì§„í–‰í•©ë‹ˆë‹¤.")
                                page_chunks = existing_progress.get('page_chunks', [])
                                
                                # Mark processed ranges based on loaded chunks
                                for chunk in page_chunks:
                                    # We need to reconstruct which ranges are done. 
                                    # Since we don't store ranges explicitly in chunks, we can infer or just skip logic below.
                                    # Better approach: Calculate the range this chunk belongs to and mark it.
                                    p_num = chunk['page_number']
                                    # Calculate start page of the chunk this page belongs to
                                    range_start = ((p_num - 1) // chunk_size) * chunk_size + 1
                                    range_end = min(range_start + chunk_size - 1, total_pages)
                                    processed_ranges.add(f"{range_start}-{range_end}")
                                
                                st.session_state.analysis_status[safe_filename]["processed_pages"] = len(page_chunks)
                            
                            for start_page in range(1, total_pages + 1, chunk_size):
                                end_page = min(start_page + chunk_size - 1, total_pages)
                                page_range = f"{start_page}-{end_page}"
                                
                                # Skip if already processed
                                if page_range in processed_ranges:
                                    st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Ready"
                                    # status_text.text(f"ìŠ¤í‚µ ì¤‘ ({idx+1}/{total_files}): {file.name} - í˜ì´ì§€ {page_range} (ì´ë¯¸ ì™„ë£Œë¨)")
                                    continue
                                
                                st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Extracting"
                                status_text.text(f"ë¶„ì„ ì¤‘ ({idx+1}/{total_files}): {file.name} - í˜ì´ì§€ {page_range} ë¶„ì„ ì¤‘...")
                                
                                # Retry logic for each chunk
                                max_retries = 3
                                for retry in range(max_retries):
                                    try:
                                        chunks = doc_intel_manager.analyze_document(blob_url, page_range=page_range, high_res=use_high_res)
                                        page_chunks.extend(chunks)
                                        
                                        # Save progress immediately
                                        save_progress(safe_filename, page_chunks, total_pages)
                                        
                                        st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Ready"
                                        st.session_state.analysis_status[safe_filename]["processed_pages"] += len(chunks)
                                        break
                                    except Exception as e:
                                        if retry == max_retries - 1:
                                            st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Failed"
                                            st.session_state.analysis_status[safe_filename]["error"] = str(e)
                                            raise e
                                        
                                        # Transient error - show friendly message
                                        wait_time = 5 * (retry + 1)
                                        status_text.text(f"â³ ì¼ì‹œì  ì§€ì—°ìœ¼ë¡œ ì¬ì—°ê²° ì¤‘ ({retry+1}/{max_retries}): {file.name} - í˜ì´ì§€ {page_range} (ì•½ {wait_time}ì´ˆ ëŒ€ê¸°)...")
                                        time.sleep(wait_time)
                            
                            # 4. Indexing
                            st.session_state.analysis_status[safe_filename]["status"] = "Indexing"
                            
                            if len(page_chunks) == 0:
                                st.warning(f"âš ï¸ ê²½ê³ : '{file.name}'ì—ì„œ í˜ì´ì§€ë¥¼ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤.")
                            
                            documents_to_index = []
                            for page_chunk in page_chunks:
                                # Create document object for each page
                                # ID must be unique and URL safe. Include page number in ID.
                                import base64
                                page_id_str = f"{blob_path}_page_{page_chunk['page_number']}"
                                doc_id = base64.urlsafe_b64encode(page_id_str.encode('utf-8')).decode('utf-8')
                                
                                document = {
                                    "id": doc_id,
                                    "content": page_chunk['content'],
                                    "content_exact": page_chunk['content'],
                                    "metadata_storage_name": f"{safe_filename} (p.{page_chunk['page_number']})",
                                    "metadata_storage_path": f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{blob_path}#page={page_chunk['page_number']}",
                                    "metadata_storage_last_modified": datetime.utcnow().isoformat() + "Z",
                                    "metadata_storage_size": file.size,
                                    "metadata_storage_content_type": file.type,
                                    "project": "drawings_analysis",  # Tag for filtering
                                    "title": page_chunk.get('ë„ë©´ëª…(TITLE)', ''),  # Drawing title
                                    "drawing_no": page_chunk.get('ë„ë©´ë²ˆí˜¸(DWG. NO.)', ''),  # Drawing number
                                    "page_number": page_chunk['page_number'],  # Page number for filtering
                                    "filename": safe_filename  # Filename for search
                                }
                                documents_to_index.append(document)
                            
                            # Batch upload all pages (50 pages at a time to avoid request size limits)
                            indexing_success = True
                            if documents_to_index:
                                batch_size = 50
                                for i in range(0, len(documents_to_index), batch_size):
                                    batch = documents_to_index[i:i + batch_size]
                                    status_text.text(f"ì¸ë±ì‹± ì¤‘ ({idx+1}/{total_files}): {safe_filename} - ë°°ì¹˜ ì „ì†¡ ì¤‘ ({i//batch_size + 1}/{(len(documents_to_index)-1)//batch_size + 1})")
                                    success, msg = search_manager.upload_documents(batch)
                                    if not success:
                                        st.error(f"ì¸ë±ì‹± ì‹¤íŒ¨ ({file.name}, ë°°ì¹˜ {i//batch_size + 1}): {msg}")
                                        indexing_success = False
                                        break
                                
                                # 5. Save Analysis JSON to Blob Storage (Dual Retrieval Strategy)
                                # Only save if indexing was successful
                                if indexing_success:
                                    status_text.text(f"ë¶„ì„ ê²°ê³¼ ì €ì¥ ì¤‘ ({idx+1}/{total_files}): {safe_filename}...")
                                    search_manager.upload_analysis_json(container_client, user_folder, safe_filename, page_chunks)
                                else:
                                    st.warning(f"âš ï¸ ì¸ë±ì‹± ì‹¤íŒ¨ë¡œ ì¸í•´ '{safe_filename}'ì˜ ë¶„ì„ ê²°ê³¼(JSON)ë¥¼ ì €ì¥í•˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
                                    # Delete the original file to prevent orphans
                                    try:
                                        st.warning(f"ğŸ§¹ ì¸ë±ì‹± ì‹¤íŒ¨ë¡œ ì¸í•´ ì›ë³¸ íŒŒì¼ '{safe_filename}'ì„ ì‚­ì œí•©ë‹ˆë‹¤.")
                                        blob_client.delete_blob()
                                        st.info("ì›ë³¸ íŒŒì¼ ì‚­ì œ ì™„ë£Œ.")
                                    except Exception as e:
                                        st.error(f"ì›ë³¸ íŒŒì¼ ì‚­ì œ ì‹¤íŒ¨: {e}")
                            
                            # Success cleanup
                            if indexing_success:
                                delete_progress(safe_filename)
                                delete_temp_file(safe_filename)
                            
                            st.session_state.analysis_status[safe_filename]["status"] = "Ready"
                            progress_bar.progress((idx + 1) / total_files)
                            
                        except Exception as e:
                            st.error(f"ì˜¤ë¥˜ ë°œìƒ ({file.name}): {str(e)}")
                    
                    status_text.text("ëª¨ë“  ì‘ì—…ì´ ì™„ë£Œë˜ì—ˆìŠµë‹ˆë‹¤!")
                    st.success("ì—…ë¡œë“œ, ë¶„ì„ ë° ì¸ë±ì‹±ì´ ì™„ë£Œë˜ì—ˆìŠµë‹ˆë‹¤.")
                    
                    # ì„±ê³µì ìœ¼ë¡œ ì™„ë£Œë˜ë©´ ì—…ë¡œë” ì´ˆê¸°í™”
                    st.session_state.drawing_uploader_key += 1
                    time.sleep(2)
                    st.rerun()

            # ğŸ“Š ë¶„ì„ ëª¨ë‹ˆí„°ë§ ëŒ€ì‹œë³´ë“œ
            if "analysis_status" in st.session_state and st.session_state.analysis_status:
                st.divider()
                st.markdown("#### ğŸ“Š ë¶„ì„ ëª¨ë‹ˆí„°ë§ ëŒ€ì‹œë³´ë“œ")
                for filename, info in st.session_state.analysis_status.items():
                    status_color = "green" if info['status'] == "Ready" else "orange" if info['status'] != "Failed" else "red"
                    with st.expander(f":{status_color}[{filename}] - {info['status']}", expanded=(info['status'] != "Ready")):
                        col1, col2 = st.columns([3, 1])
                        with col1:
                            st.write(f"**ì „ì²´ ìƒíƒœ:** {info['status']}")
                            progress = info['processed_pages'] / info['total_pages'] if info['total_pages'] > 0 else 0
                            st.progress(progress)
                            st.write(f"**ì§„í–‰ë„:** {info['processed_pages']} / {info['total_pages']} í˜ì´ì§€ ì™„ë£Œ")
                    
                        if info['error']:
                            st.error(f"**ìµœê·¼ ì˜¤ë¥˜:** {info['error']}")
                    
                        # ì„¸ë¶€ ì²­í¬ ìƒíƒœ
                        if info['chunks']:
                            st.markdown("---")
                            st.caption("ğŸ§© **í˜ì´ì§€ ì²­í¬ë³„ ìƒíƒœ**")
                            chunk_cols = st.columns(4)
                            for i, (chunk_range, chunk_status) in enumerate(info['chunks'].items()):
                                with chunk_cols[i % 4]:
                                    if chunk_status == "Ready":
                                        st.success(f"âœ… {chunk_range}")
                                    elif chunk_status == "Failed":
                                        st.error(f"âŒ {chunk_range}")
                                        # ì¬ì‹œë„ ë²„íŠ¼ (ê°„ì†Œí™”ëœ êµ¬í˜„)
                                        if st.button("ğŸ”„", key=f"retry_{filename}_{chunk_range}", help=f"{chunk_range} ì¬ì‹œë„"):
                                            st.info("ì¬ì‹œë„ëŠ” 'ì—…ë¡œë“œ ë° ë¶„ì„ ì‹œì‘'ì„ ë‹¤ì‹œ ëˆŒëŸ¬ì£¼ì„¸ìš” (ë©±ë“±ì„± ë³´ì¥)")
                                    else:
                                        st.info(f"â³ {chunk_range}")

        with tab2:

        
            # Display analyzed documents
            st.markdown("#### ğŸ“‹ ë¶„ì„ëœ ë¬¸ì„œ ëª©ë¡")
            try:
                blob_service_client = get_blob_service_client()
                container_client = blob_service_client.get_container_client(CONTAINER_NAME)
            
                # List files in user's drawings folder + Admin access to root drawings
                blobs = []
                # User folder
                blobs.extend(list(container_client.list_blobs(name_starts_with=f"{user_folder}/drawings/")))
            
                if user_role == 'admin':
                    # Admin root folder
                    blobs.extend(list(container_client.list_blobs(name_starts_with="drawings/")))
            
                # Deduplicate
                unique_blobs = {b.name: b for b in blobs}.values()
            
                blob_list = []
                available_filenames = []
                for blob in unique_blobs:
                    if not blob.name.endswith('/'):  # Skip folder markers
                        filename = blob.name.split('/')[-1]
                        blob_list.append({
                            'name': filename,
                            'full_name': blob.name,
                            'size': blob.size,
                            'modified': blob.last_modified
                        })
                        available_filenames.append(filename)
            
                # Sort by modified date (most recent first)
                blob_list.sort(key=lambda x: x['modified'], reverse=True)
            
                selected_filenames = []
            
                if blob_list:
                    st.info(f"ì´ {len(blob_list)}ê°œì˜ ë¬¸ì„œê°€ ë¶„ì„ë˜ì–´ ìˆìŠµë‹ˆë‹¤. ë¶„ì„í•  ë¬¸ì„œë¥¼ ì„ íƒí•˜ì„¸ìš”.")
                
                    # Add "Select All" checkbox
                    def toggle_all():
                        new_state = st.session_state.select_all_files
                        # Update state for ALL files in the list, not just existing keys
                        for b in blob_list:
                            st.session_state[f"chk_{b['name']}"] = new_state

                    select_all = st.checkbox("ì „ì²´ ì„ íƒ", value=False, key="select_all_files", on_change=toggle_all)
                
                    # Display as expandable list
                    with st.expander("ğŸ“„ ë¬¸ì„œ ëª©ë¡ ë° ì„ íƒ", expanded=True):
                        for idx, blob_info in enumerate(blob_list, 1):
                            # Improved column layout for better zoom stability
                            # col0: checkbox (3%), col1: filename (59%), col2: 3 icons (27%), col3: delete+JSON (11%)
                            col0, col1, col2, col3 = st.columns([0.3, 5.9, 2.7, 1.1])
                            
                            with col0:
                                # Checkbox for selection
                                chk_key = f"chk_{blob_info['name']}"
                                if chk_key not in st.session_state:
                                    st.session_state[chk_key] = False
                                
                                is_selected = st.checkbox(f"select_{idx}", key=chk_key, label_visibility="collapsed")
                                if is_selected:
                                    selected_filenames.append(blob_info['name'])
                        
                            with col1:
                                size_mb = blob_info['size'] / (1024 * 1024)
                                st.markdown(f"**{blob_info['name']}** ({size_mb:.2f} MB)")
                        
                            with col2:
                                # 3 action icons in a row
                                icon_c1, icon_c2, icon_c3 = st.columns(3)
                            
                                with icon_c1:
                                    # Download Button
                                    try:
                                        sas_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            blob_info['full_name'], 
                                            content_disposition="attachment"
                                        )
                                        st.link_button("ğŸ“¥", sas_url, help="ë‹¤ìš´ë¡œë“œ", use_container_width=True)
                                    except Exception as e:
                                        st.error(f"Err: {e}")

                                with icon_c2:
                                    # 2. Rename Button (Popover)
                                    with st.popover("âœï¸", use_container_width=True):
                                        new_name_input = st.text_input("ìƒˆ íŒŒì¼ëª…", value=blob_info['name'], key=f"ren_{blob_info['name']}")
                                        if st.button("ì´ë¦„ ë³€ê²½", key=f"btn_ren_{blob_info['name']}"):
                                            if new_name_input != blob_info['name']:
                                                try:
                                                    with st.spinner("ì´ë¦„ ë³€ê²½ ë° ì¸ë±ìŠ¤ ì—…ë°ì´íŠ¸ ì¤‘..."):
                                                        # A. Rename Blob
                                                        old_blob_name = blob_info['full_name']
                                                        folder_path = old_blob_name.rsplit('/', 1)[0]
                                                        new_blob_name = f"{folder_path}/{new_name_input}"
                                                    
                                                        source_blob = container_client.get_blob_client(old_blob_name)
                                                        dest_blob = container_client.get_blob_client(new_blob_name)
                                                    
                                                        # Copy
                                                        source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, old_blob_name)
                                                        dest_blob.start_copy_from_url(source_sas)
                                                    
                                                        # Wait for copy
                                                        for _ in range(20):
                                                            props = dest_blob.get_blob_properties()
                                                            if props.copy.status == "success":
                                                                break
                                                            time.sleep(0.2)
                                                    
                                                        # B. Update Search Index (Preserve OCR Data)
                                                        search_manager = get_search_manager()
                                                        import unicodedata
                                                        safe_old_filename = unicodedata.normalize('NFC', blob_info['name'])
                                                        safe_new_filename = unicodedata.normalize('NFC', new_name_input)
                                                    
                                                        # Find old docs
                                                        results = search_manager.search_client.search(
                                                            search_text="*",
                                                            filter=f"project eq 'drawings_analysis'",
                                                            select=["id", "content", "content_exact", "metadata_storage_name", "metadata_storage_path", "metadata_storage_size", "metadata_storage_content_type"]
                                                        )
                                                    
                                                        docs_to_upload = []
                                                        ids_to_delete = []
                                                    
                                                        for doc in results:
                                                            # Check if this doc belongs to the file (by name prefix)
                                                            # Name format: "{filename} (p.{page})"
                                                            if doc['metadata_storage_name'].startswith(safe_old_filename):
                                                                # Create new doc
                                                                page_suffix = doc['metadata_storage_name'].split(safe_old_filename)[-1] # e.g. " (p.1)"
                                                            
                                                                # New ID
                                                                import base64
                                                                # Extract page number from suffix or path if possible, or just reconstruct
                                                                # Path format: .../filename#page=N
                                                                try:
                                                                    page_num = doc['metadata_storage_path'].split('#page=')[-1]
                                                                    new_page_id_str = f"{new_blob_name}_page_{page_num}"
                                                                    new_doc_id = base64.urlsafe_b64encode(new_page_id_str.encode('utf-8')).decode('utf-8')
                                                                
                                                                    new_doc = {
                                                                        "id": new_doc_id,
                                                                        "content": doc['content'],
                                                                        "content_exact": doc.get('content_exact', doc['content']),
                                                                        "metadata_storage_name": f"{safe_new_filename}{page_suffix}",
                                                                        "metadata_storage_path": f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{new_blob_name}#page={page_num}",
                                                                        "metadata_storage_last_modified": datetime.utcnow().isoformat() + "Z",
                                                                        "metadata_storage_size": doc['metadata_storage_size'],
                                                                        "metadata_storage_content_type": doc['metadata_storage_content_type'],
                                                                        "project": "drawings_analysis"
                                                                    }
                                                                    docs_to_upload.append(new_doc)
                                                                    ids_to_delete.append({"id": doc['id']})
                                                                except:
                                                                    pass
                                                                    
                                                        if docs_to_upload:
                                                            search_manager.upload_documents(docs_to_upload)
                                                        if ids_to_delete:
                                                            search_manager.search_client.delete_documents(documents=ids_to_delete)

                                                        # C. Delete old blob
                                                        source_blob.delete_blob()
                                                    
                                                        st.success("ì´ë¦„ ë³€ê²½ ì™„ë£Œ!")
                                                        time.sleep(1)
                                                        st.rerun()
                                                    
                                                except Exception as e:
                                                    st.error(f"ë³€ê²½ ì‹¤íŒ¨: {e}")

                                with icon_c3:
                                    # 3. Re-analyze Button
                                    if st.button("ğŸ”„", key=f"reanalyze_{blob_info['name']}", help="ì¬ë¶„ì„ (ì¸ë±ìŠ¤ ë³µêµ¬)", use_container_width=True):
                                        try:
                                            with st.spinner("ì¬ë¶„ì„ ì‹œì‘... (íŒŒì¼ ë‹¤ìš´ë¡œë“œ ì¤‘)"):
                                                # A. Download Blob to memory
                                                blob_client = container_client.get_blob_client(blob_info['full_name'])
                                                download_stream = blob_client.download_blob()
                                                pdf_data = download_stream.readall()
                                            
                                                # B. Count Pages
                                                import fitz
                                                doc = fitz.open(stream=pdf_data, filetype="pdf")
                                                total_pages = doc.page_count
                                            
                                                # C. Initialize Status
                                                if "analysis_status" not in st.session_state:
                                                    st.session_state.analysis_status = {}
                                            
                                                safe_filename = blob_info['name']
                                                st.session_state.analysis_status[safe_filename] = {
                                                    "status": "Extracting",
                                                    "total_pages": total_pages,
                                                    "processed_pages": 0,
                                                    "chunks": {},
                                                    "error": None
                                                }
                                            
                                                # D. Analyze Chunks
                                                doc_intel_manager = get_doc_intel_manager()
                                                search_manager = get_search_manager()
                                                blob_service_client = get_blob_service_client()
                                            
                                                # Generate SAS for Analysis
                                                sas_token = generate_blob_sas(
                                                    account_name=blob_service_client.account_name,
                                                    container_name=CONTAINER_NAME,
                                                    blob_name=blob_info['full_name'],
                                                    account_key=blob_service_client.credential.account_key,
                                                    permission=BlobSasPermissions(read=True),
                                                    expiry=datetime.utcnow() + timedelta(hours=1)
                                                )
                                                # Use relative path for URL construction if needed, but full_name is usually relative to container if listed from container_client?
                                                # container_client.list_blobs returns name relative to container.
                                                blob_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_info['full_name'])}?{sas_token}"
                                            
                                                chunk_size = 50
                                                page_chunks = []
                                            
                                                progress_bar = st.progress(0)
                                                status_text = st.empty()
                                            
                                                for start_page in range(1, total_pages + 1, chunk_size):
                                                    end_page = min(start_page + chunk_size - 1, total_pages)
                                                    page_range = f"{start_page}-{end_page}"
                                                
                                                    st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Extracting"
                                                    status_text.text(f"ì¬ë¶„ì„ ì¤‘: {safe_filename} ({page_range})...")
                                                
                                                    # Retry logic
                                                    max_retries = 3
                                                    for retry in range(max_retries):
                                                        try:
                                                            # Use default high_res=False for re-analysis
                                                            chunks = doc_intel_manager.analyze_document(blob_url, page_range=page_range, high_res=False)
                                                            page_chunks.extend(chunks)
                                                            st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Ready"
                                                            st.session_state.analysis_status[safe_filename]["processed_pages"] += len(chunks)
                                                            break
                                                        except Exception as e:
                                                            if retry == max_retries - 1:
                                                                st.session_state.analysis_status[safe_filename]["chunks"][page_range] = "Failed"
                                                                st.session_state.analysis_status[safe_filename]["error"] = str(e)
                                                                raise e
                                                            time.sleep(5 * (retry + 1))
                                            
                                                # E. Indexing
                                                st.session_state.analysis_status[safe_filename]["status"] = "Indexing"
                                                status_text.text("ì¸ë±ì‹± ì¤‘...")
                                            
                                                documents_to_index = []
                                                for page_chunk in page_chunks:
                                                    import base64
                                                    # Use full_name (path in container) for ID generation to match upload logic
                                                    page_id_str = f"{blob_info['full_name']}_page_{page_chunk['page_number']}"
                                                    doc_id = base64.urlsafe_b64encode(page_id_str.encode('utf-8')).decode('utf-8')
                                                
                                                    document = {
                                                        "id": doc_id,
                                                        "content": page_chunk['content'],
                                                        "content_exact": page_chunk['content'],
                                                        "metadata_storage_name": f"{safe_filename} (p.{page_chunk['page_number']})",
                                                        "metadata_storage_path": f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{blob_info['full_name']}#page={page_chunk['page_number']}",
                                                        "metadata_storage_last_modified": datetime.utcnow().isoformat() + "Z",
                                                        "metadata_storage_size": blob_info['size'],
                                                        "metadata_storage_content_type": "application/pdf",
                                                        "project": "drawings_analysis",
                                                        "title": page_chunk.get('ë„ë©´ëª…(TITLE)', ''),  # Drawing title
                                                        "drawing_no": page_chunk.get('ë„ë©´ë²ˆí˜¸(DWG. NO.)', ''),  # Drawing number
                                                        "page_number": page_chunk['page_number'],  # Page number for filtering
                                                        "filename": safe_filename  # Filename for search
                                                    }
                                                    documents_to_index.append(document)
                                            
                                                if documents_to_index:
                                                    batch_size = 50
                                                    for i in range(0, len(documents_to_index), batch_size):
                                                        batch = documents_to_index[i:i + batch_size]
                                                        success, msg = search_manager.upload_documents(batch)
                                                        if not success:
                                                            st.error(f"âŒ ì¸ë±ìŠ¤ ì—…ë¡œë“œ ì‹¤íŒ¨ (ë°°ì¹˜ {i//batch_size + 1}): {msg}")
                                                            raise Exception(f"Index upload failed: {msg}")
                                                
                                                    # Save JSON only if upload succeeded
                                                    search_manager.upload_analysis_json(container_client, user_folder, safe_filename, page_chunks)
                                            
                                                st.session_state.analysis_status[safe_filename]["status"] = "Ready"
                                                st.success("ì¬ë¶„ì„ ì™„ë£Œ! ì´ì œ ê²€ìƒ‰ì´ ê°€ëŠ¥í•©ë‹ˆë‹¤.")
                                                time.sleep(1)
                                                st.rerun()

                                        except Exception as e:
                                            st.error(f"ì¬ë¶„ì„ ì‹¤íŒ¨: {e}")

                                # 3. JSON (Admin only)
                                if user_role == 'admin':
                                    json_key = f"json_data_{blob_info['name']}"
                                
                                    if json_key not in st.session_state:
                                        if st.button("JSON", key=f"gen_json_{blob_info['name']}"):
                                            with st.spinner("..."):
                                                search_manager = get_search_manager()
                                                # Dual Retrieval Strategy: Try Blob first
                                                docs = search_manager.get_document_json_from_blob(container_client, user_folder, blob_info['name'])
                                            
                                                # Fallback to AI Search if Blob JSON not found (for older files)
                                                if not docs:
                                                    st.info("Blob JSONì„ ì°¾ì„ ìˆ˜ ì—†ì–´ AI Searchì—ì„œ ê²€ìƒ‰í•©ë‹ˆë‹¤...")
                                                    docs = search_manager.get_document_json(blob_info['name'])
                                                
                                                if docs:
                                                    import json
                                                    json_str = json.dumps(docs, ensure_ascii=False, indent=2)
                                                    st.session_state[json_key] = json_str
                                                    st.rerun()
                                                else:
                                                    st.error(f"No Data found for '{blob_info['name']}'")
                                                    # Try one more time without project filter to see if it exists at all
                                                    safe_name = blob_info['name'].replace("'", "''")
                                                    debug_docs = search_manager.search_client.search(
                                                        search_text="*",
                                                        filter=f"search.ismatch('\"{safe_name}*\"', 'metadata_storage_name')",
                                                        select=["metadata_storage_name", "project"],
                                                        top=5
                                                    )
                                                    debug_list = list(debug_docs)
                                                    if debug_list:
                                                        st.warning(f"Found {len(debug_list)} docs without correct project tag. Example: {debug_list[0].get('metadata_storage_name')} (Project: {debug_list[0].get('project')})")
                                                    else:
                                                        st.error("Document not found in index at all.")
                                    else:
                                        # Show download button
                                        json_data = st.session_state[json_key]
                                        st.download_button(
                                            label="ğŸ’¾",
                                            data=json_data,
                                            file_name=f"{blob_info['name']}.json",
                                            mime="application/json",
                                            key=f"dl_json_{blob_info['name']}"
                                        )

                            with col3:
                                if st.button("ğŸ—‘ï¸ ì‚­ì œ", key=f"del_{blob_info['name']}"):
                                    try:
                                        # 1. Delete from Blob Storage (Use full_name)
                                        blob_client = container_client.get_blob_client(blob_info['full_name'])
                                        blob_client.delete_blob()
                                    
                                        # 2. Delete from Search Index
                                        search_manager = get_search_manager()
                                    
                                        # Find docs to delete
                                        import unicodedata
                                        safe_filename = unicodedata.normalize('NFC', blob_info['name'])
                                    
                                        # Clean up index (Find ALL pages)
                                        ids_to_delete = []
                                        while True:
                                            results = search_manager.search_client.search(
                                                search_text="*",
                                                filter=f"project eq 'drawings_analysis'",
                                                select=["id", "metadata_storage_name"],
                                                top=1000
                                            )
                                        
                                            batch_ids = []
                                            for doc in results:
                                                # Use NFC normalization for comparison
                                                doc_name = unicodedata.normalize('NFC', doc['metadata_storage_name'])
                                                if doc_name.startswith(safe_filename):
                                                    batch_ids.append({"id": doc['id']})
                                        
                                            if not batch_ids:
                                                break
                                            
                                            search_manager.search_client.delete_documents(documents=batch_ids)
                                            ids_to_delete.extend(batch_ids)
                                        
                                            # If we found less than 1000, we might be done, but to be safe we continue 
                                            # until a search returns no matches for our file.
                                            # Actually, if we delete them, the next search will return different docs.
                                            # So we continue until no more docs match.
                                            if len(batch_ids) == 0:
                                                break
                                    
                                        # Clear JSON state if exists
                                        json_key = f"json_data_{blob_info['name']}"
                                        if json_key in st.session_state:
                                            del st.session_state[json_key]

                                        st.success(f"{blob_info['name']} ì‚­ì œ ì™„ë£Œ")
                                        st.rerun()
                                    except Exception as e:
                                        st.error(f"ì‚­ì œ ì‹¤íŒ¨: {e}")
                else:
                    st.warning("ë¶„ì„ëœ ë¬¸ì„œê°€ ì—†ìŠµë‹ˆë‹¤. 'ë¬¸ì„œ ì—…ë¡œë“œ ë° ë¶„ì„' íƒ­ì—ì„œ ë¬¸ì„œë¥¼ ì—…ë¡œë“œí•˜ì„¸ìš”.")
            except Exception as e:
                st.error(f"ë¬¸ì„œ ëª©ë¡ì„ ë¶ˆëŸ¬ì˜¤ëŠ” ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")
        
            st.divider()
        
            # DEBUG: Show selected files
            if user_role == 'admin':
                # st.write(f"DEBUG: Selected Files ({len(selected_filenames)}): {selected_filenames}")
                pass

            # -----------------------------
            # Advanced Search Options (RAG) - Added to match AI Q&A
            # -----------------------------
            with st.expander("âš™ï¸ ê³ ê¸‰ ê²€ìƒ‰ ì˜µì…˜ (RAG ì„¤ì •)", expanded=False):
                c1, c2 = st.columns(2)
                with c1:
                    rag_use_semantic = st.checkbox("ì‹œë§¨í‹± ë­ì»¤ ì‚¬ìš©", value=True, key="rag_use_semantic", help="ì˜ë¯¸ ê¸°ë°˜ ê²€ìƒ‰ì„ ì‚¬ìš©í•˜ì—¬ ì •í™•ë„ë¥¼ ë†’ì…ë‹ˆë‹¤.")
                with c2:
                    rag_search_mode_opt = st.radio("ê²€ìƒ‰ ëª¨ë“œ", ["all (AND)", "any (OR)"], index=1, horizontal=True, key="rag_search_mode", help="any: í‚¤ì›Œë“œ ì¤‘ í•˜ë‚˜ë¼ë„ í¬í•¨ë˜ë©´ ê²€ìƒ‰ (ì¶”ì²œ)")
                    rag_search_mode = "all" if "all" in rag_search_mode_opt else "any"
                
                # Output Format Toggle
                st.write("")
                st.markdown("**ë‹µë³€ í˜•ì‹ (Output Format)**")
                output_format = st.radio(
                    "ë‹µë³€ í˜•ì‹ì„ ì„ íƒí•˜ì„¸ìš”", 
                    ["Table (í‘œ)", "Text (í…ìŠ¤íŠ¸)"], 
                    index=1, 
                    horizontal=True, 
                    label_visibility="collapsed",
                    key="rag_output_format"
                )
        
            # Chat Interface (Similar to main chat but focused)
            if "rag_chat_messages" not in st.session_state:
                st.session_state.rag_chat_messages = []
            
            for message in st.session_state.rag_chat_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
                    if "citations" in message and message["citations"]:
                        st.markdown("---")
                        st.caption("ğŸ“š **ì°¸ì¡° ë¬¸ì„œ:**")
                        for i, citation in enumerate(message["citations"], 1):
                            filepath = citation.get('filepath', 'Unknown')
                            # Use pre-generated final_url if available, otherwise generate one
                            display_url = citation.get('final_url')
                            if not display_url:
                                try:
                                    blob_service_client = get_blob_service_client()
                                    display_url = generate_sas_url(
                                        blob_service_client, 
                                        CONTAINER_NAME, 
                                        filepath, 
                                        page=citation.get('page')
                                    )
                                except:
                                    display_url = "#"
                        
                            st.markdown(f"{i}. [{filepath}]({display_url})")

            if prompt := st.chat_input("ë„ë©´ì´ë‚˜ ìŠ¤í™ì— ëŒ€í•´ ì§ˆë¬¸í•˜ì„¸ìš”..."):
                st.session_state.rag_chat_messages.append({"role": "user", "content": prompt})
                with st.chat_message("user"):
                    st.markdown(prompt)
            
                with st.chat_message("assistant"):
                    with st.spinner("ë¶„ì„ ì¤‘..."):
                        try:
                            chat_manager = get_chat_manager()
                        
                            conversation_history = [
                                {"role": msg["role"], "content": msg["content"]}
                                for msg in st.session_state.rag_chat_messages[:-1]
                            ]
                        
                            # Use 'any' search mode for better recall (find documents even with partial keyword match)
                            # This is important because technical drawings may have specific terms
                            # Filter to only search documents from the drawings folder
                            # Pass selected_filenames for specific file filtering
                            # If selected_filenames is empty (user deselected all), we should probably warn or search nothing.
                            # But for now let's pass it. If empty, the chat manager might search nothing or all depending on logic.
                            # Actually, let's default to all if none selected? No, user explicitly deselected.
                            # Let's pass the list as is.
                        
                            # Note: selected_filenames is defined in the outer scope of the tab
                            current_files = selected_filenames
                        
                            # Construct robust filter expression
                            # Include fallback for documents that might have lost their project tag but are in the drawings folder
                            base_filter = "(project eq 'drawings_analysis' or search.ismatch('/drawings/', 'metadata_storage_path'))"
                        
                            # Note: We used to filter by path here, but OData encoding issues caused 0 results.
                            # Now we pass user_folder to chat_manager for Python-side filtering.

                            # Append Output Format Instruction
                            final_prompt = prompt
                            if "Table" in output_format:
                                final_prompt += "\n\n[OUTPUT INSTRUCTION]: Please summarize the comparison in a **Markdown Table**."
                            else:
                                final_prompt += "\n\n[OUTPUT INSTRUCTION]: Please summarize the comparison in **Structured Markdown Text**. Do NOT use a table."

                            response_text, citations, context, final_filter, search_results = chat_manager.get_chat_response(
                                final_prompt, 
                                conversation_history,
                                search_mode=rag_search_mode,
                                use_semantic_ranker=rag_use_semantic,
                                filter_expr=base_filter,
                                available_files=current_files,
                                user_folder=user_folder,
                                is_admin=(user_role == 'admin')
                            )

                            # ---------------------------------------------------------
                            # CRITICAL: Linkify Inline Citations & Escape Tildes
                            # ---------------------------------------------------------
                            
                            # 1. Pre-generate Web Viewer URLs for all citations
                            citation_links = {}
                            processed_citations = [] # Store processed citations with URLs for the bottom list
                            
                            if citations:
                                for cit in citations:
                                    filepath = cit.get('filepath', 'Unknown')
                                    # CRITICAL: Clean filepath from page suffixes like " (p.1)" or " (p.1) (p.1)"
                                    import re
                                    clean_filepath = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', filepath).strip()
                                    
                                    page = cit.get('page')
                                    url = cit.get('url', '')
                                    
                                    # Generate Web Viewer URL
                                    try:
                                        blob_service_client = get_blob_service_client()
                                        final_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            clean_filepath, 
                                            page=page
                                        )
                                    except Exception as e:
                                        st.error(f"URL ìƒì„± ì‹¤íŒ¨ ({clean_filepath}): {e}")
                                        final_url = "#"
                                    
                                    cit['final_url'] = final_url
                                    processed_citations.append(cit)
                                    
                                    filename = os.path.basename(filepath)
                                    if page:
                                        citation_links[(filename, str(page))] = final_url
                            
                            # 2. Replace text citations with Markdown links (Support both [] and ())
                            if response_text:
                                import re
                                
                                # Reference dictionary to store URLs
                                link_references = {}
                                
                                def replace_citation(match):
                                    content = match.group(1).strip()
                                    
                                    # Remove "ë¬¸ì„œëª…:" prefix if present (Common in Korean LLM outputs)
                                    content = re.sub(r'^ë¬¸ì„œëª…\s*:\s*', '', content)

                                    # Split by last colon to separate filename and page
                                    if ':' in content:
                                        fname = content.rsplit(':', 1)[0].strip()
                                        p_num = match.group(2)
                                    else:
                                        return match.group(0)
                                        
                                    original_text = match.group(0)
                                    
                                    # Try to find matching citation with fuzzy logic
                                    target_url = None
                                    matched_filename = None
                                    
                                    # 1. Clean LLM filename for comparison
                                    clean_llm = re.sub(r'\.pdf$', '', fname.lower().strip())
                                    
                                    for (k_fname, k_page), url in citation_links.items():
                                        # 2. Clean known filename (remove .pdf and (p.N) suffixes)
                                        clean_known = re.sub(r'\.pdf$', '', k_fname.lower().strip())
                                        clean_known = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', clean_known).strip()
                                        
                                        # CRITICAL FIX: Skip empty filenames to prevent false positive matches
                                        if not clean_known:
                                            continue

                                        # 3. Match page and filename (fuzzy)
                                        if str(k_page) == str(p_num):
                                            # Exact match or containment (handling "ë¬¸ì„œëª…: " residue if regex failed)
                                            if clean_llm == clean_known or clean_llm in clean_known or clean_known in clean_llm:
                                                target_url = url
                                                # CRITICAL FIX: Capture the matched filename to replace text
                                                matched_filename = k_fname
                                                break
                                    
                                    if target_url:
                                        # Use Markdown link for table compatibility
                                        # Escape parentheses in URL to avoid breaking Markdown link
                                        safe_url = target_url.replace('(', '%28').replace(')', '%29')
                                        
                                        # Use Inline Link with Icon and Hover Title for compact display
                                        # Escape parentheses in URL
                                        safe_url = target_url.replace('(', '%28').replace(')', '%29')
                                        
                                        hover_text = original_text
                                        if matched_filename:
                                            # Remove .pdf extension for cleaner display
                                            display_filename = re.sub(r'\.pdf$', '', matched_filename, flags=re.IGNORECASE)
                                            # Use hyphen instead of parentheses to avoid Markdown link issues
                                            hover_text = f"{display_filename} - p.{p_num}"
                                        
                                        # Return Icon Link
                                        return f"[ğŸ”—]({safe_url} \"{hover_text}\")"
                                    
                                    # Fallback: Check if original_text contains a URL (e.g. from LLM output)
                                    # Match (http...) or (blob...)
                                    url_match = re.search(r'\((https?://[^)]+|blob:[^)]+)\)', original_text)
                                    if url_match:
                                        fallback_url = url_match.group(1)
                                        safe_url = fallback_url.replace('(', '%28').replace(')', '%29')
                                        
                                        # Construct hover text from captured groups
                                        # fname is from group(1), p_num from group(2) of the main regex
                                        display_filename = re.sub(r'\.pdf$', '', fname.strip(), flags=re.IGNORECASE)
                                        hover_text = f"{display_filename} - p.{p_num}"
                                        
                                        return f"[ğŸ”—]({safe_url} \"{hover_text}\")"

                                    return original_text

                                # DEBUG: Show raw response before linkification
                                # st.code(response_text, language="markdown")
                                
                                # Pass 1: Strict match (Double Brackets [[...]]) - NEW STANDARD
                                # Matches [[File: p.1]] or [[File - p.1]]
                                # Also consumes trailing URL if present (e.g. [[...]] (url))
                                pattern_double = r'\[\[(.*?)[:\-]\s*p\.?\s*(\d+)\]\](?:\s*\((?:https?|blob):[^)]+\))?'
                                response_text = re.sub(pattern_double, replace_citation, response_text)

                                # Pass 2: Strict match (Closed parentheses/brackets) - LEGACY FALLBACK
                                # Updated to consume trailing URL if present (e.g. [Title: p.1](url)) to avoid double links
                                # Also updated to allow parentheses in filename (e.g. (File(v1): p.1))
                                # Updated to allow hyphen separator (e.g. [File - p.1])
                                pattern_strict = r'[\[\(]([^\[\]]+?[:\-]\s*p\.?\s*(\d+))[\]\)](?:\s*\((?:https?|blob):[^)]+\))?'
                                response_text = re.sub(pattern_strict, replace_citation, response_text)
                                
                                # Pass 3: Truncated match (Open parentheses at end of string) - Fallback
                                # Only matches if it's at the very end of the string or followed by newline
                                pattern_truncated = r'[\[\(]([^\[\]]+?[:\-]\s*p\.?\s*(\d+))(?:\s*\((?:https?|blob):[^)]+\))?$'
                                response_text = re.sub(pattern_truncated, replace_citation, response_text)

                                # Append Reference Definitions to the end of the response
                                if link_references:
                                    response_text += "\n\n"
                                    for ref_id, url in link_references.items():
                                        response_text += f"[{ref_id}]: {url}\n"

                                # 3. Escape tildes (AFTER linkification to avoid breaking links if they contained tildes, though unlikely in URLs)
                                response_text = response_text.replace('~', '\\~')
                        
                            st.markdown(response_text, unsafe_allow_html=True)
                            
                            # Display Google-like search results (Snippets + Links)
                            if search_results:
                                with st.expander("ğŸ” ê²€ìƒ‰ ê²°ê³¼ ë° ìŠ¤ë‹ˆí« (ìƒìœ„ í›„ë³´)", expanded=True):
                                    for i, res in enumerate(search_results[:5]): # Show top 5 for clarity
                                        res_name = res.get('metadata_storage_name', 'Unknown')
                                        res_path = res.get('metadata_storage_path', '')
                                    
                                        # Extract snippet from highlights
                                        highlights = res.get('@search.highlights', {})
                                        snippet = highlights.get('content', [""])[0] if highlights else ""
                                        if not snippet:
                                            snippet = res.get('content', '')[:200] + "..."
                                    
                                        # Generate SAS link for the result
                                        try:
                                            # Extract blob path from metadata_storage_path
                                            from urllib.parse import unquote
                                            import re
                                        
                                            if "https://direct_fetch/" in res_path:
                                                # Handle custom direct fetch scheme
                                                path_without_scheme = res_path.replace("https://direct_fetch/", "")
                                                blob_path_part = path_without_scheme.split('#')[0]
                                                blob_path_part = unquote(blob_path_part)
                                            elif CONTAINER_NAME in res_path:
                                                # Handle standard Azure Blob URL
                                                blob_path_part = res_path.split(f"/{CONTAINER_NAME}/")[1].split('#')[0]
                                                blob_path_part = unquote(blob_path_part)
                                            else:
                                                # Fallback or relative path
                                                blob_path_part = res_path
                                        
                                            # CRITICAL FIX: Strip " (p.N)" suffix if present in the path
                                            # This happens if the indexer appended it to the path
                                            blob_path_part = re.sub(r'\s*\(p\.\d+\)$', '', blob_path_part)
                                            
                                            # Generate SAS Token
                                            sas_token = generate_blob_sas(
                                                account_name=blob_service_client.account_name,
                                                container_name=CONTAINER_NAME,
                                                blob_name=blob_path_part,
                                                account_key=blob_service_client.credential.account_key,
                                                permission=BlobSasPermissions(read=True),
                                                expiry=datetime.utcnow() + timedelta(hours=1),
                                                content_disposition="inline",
                                                content_type="application/pdf" # Default to PDF for viewer hint
                                            )
                                            sas_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}/{urllib.parse.quote(blob_path_part, safe='/')}?{sas_token}"

                                            # Use Office Online Viewer for Office files ONLY
                                            # PDF files use direct SAS URL (browser viewer) for better page linking
                                            lower_name = res_name.lower()
                                            if lower_name.endswith(('.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls')):
                                                encoded_sas_url = urllib.parse.quote(sas_url)
                                                final_url = f"https://view.officeapps.live.com/op/view.aspx?src={encoded_sas_url}"
                                                link_text = "ğŸ“„ ì›¹ì—ì„œ ë³´ê¸° (Office Viewer)"
                                            elif lower_name.endswith('.pdf'):
                                                # Direct SAS URL for PDF (No Google Viewer)
                                                final_url = sas_url
                                                link_text = "ğŸ“„ ë¬¸ì„œ ì—´ê¸° (ìƒˆ íƒ­)"
                                            else:
                                                final_url = sas_url
                                                link_text = "ğŸ“„ ë¬¸ì„œ ì—´ê¸° (ìƒˆ íƒ­)"

                                        except:
                                            final_url = "#"
                                            link_text = "ë§í¬ ìƒì„± ì‹¤íŒ¨"

                                        st.markdown(f"**{i+1}. {res_name}**")
                                        st.write(f"_{snippet}_")
                                        if final_url != "#":
                                            st.markdown(f"[{link_text}]({final_url})")
                                        st.divider()

                            if processed_citations:
                                st.markdown("---")
                                st.caption("ğŸ“š **ì°¸ì¡° ë¬¸ì„œ (í˜ì´ì§€ë³„ ë§í¬):**")
                                
                                # Group citations by filename
                                from collections import defaultdict
                                pages_by_file = defaultdict(set)
                                
                                for cit in processed_citations:
                                    fp = cit.get('filepath', 'Unknown')
                                    pg = cit.get('page')
                                    
                                    # Clean filepath
                                    clean_fp = re.sub(r'\s*\(\s*p\.?\s*\d+\s*\)', '', fp).strip()
                                    
                                    if pg:
                                        try:
                                            pg_int = int(pg)
                                            pages_by_file[clean_fp].add(pg_int)
                                        except:
                                            pass
                                    else:
                                        # Ensure file is listed even if no specific page
                                        if clean_fp not in pages_by_file:
                                            pages_by_file[clean_fp] = set()

                                # Display grouped citations
                                for i, (fp, pages) in enumerate(sorted(pages_by_file.items()), 1):
                                    filename = os.path.basename(fp)
                                    
                                    # Generate Doc URL (Page 1)
                                    try:
                                        blob_service_client = get_blob_service_client()
                                        doc_url = generate_sas_url(
                                            blob_service_client, 
                                            CONTAINER_NAME, 
                                            fp, 
                                            page=1
                                        )
                                    except:
                                        doc_url = "#"
                                    
                                    # Base line: Document Title
                                    line = f"**{i}. [{filename}]({doc_url})**"
                                    
                                    # Append Page Links
                                    sorted_pages = sorted(pages)
                                    if sorted_pages:
                                        page_links = []
                                        for p in sorted_pages:
                                            try:
                                                p_url = generate_sas_url(
                                                    blob_service_client, 
                                                    CONTAINER_NAME, 
                                                    fp, 
                                                    page=p
                                                )
                                                page_links.append(f"[p.{p}]({p_url})")
                                            except:
                                                pass
                                        
                                        if page_links:
                                            line += " â€” " + " Â· ".join(page_links)
                                    
                                    st.markdown(line)
                            
                            # Debug: Show Citation Links (Hidden by default)
                            # with st.expander("ğŸ” ë§í¬ ë””ë²„ê¹… (Debug Links)", expanded=False):
                            #     st.write("Citation Links Keys:", list(citation_links.keys()))
                            #     st.write("Processed Citations:", processed_citations)
                        
                            # Debug: Show Context
                            with st.expander("ğŸ” ê²€ìƒ‰ëœ ì»¨í…ìŠ¤íŠ¸ í™•ì¸ (Debug Context)", expanded=False):
                                if final_filter:
                                    st.caption(f"**OData Filter:** `{final_filter}`")
                                if search_results:
                                    st.caption(f"**Search Results:** {len(search_results)} chunks found")
                                st.text_area("LLMì—ê²Œ ì „ë‹¬ëœ ì›ë¬¸ ë°ì´í„°", value=context, height=300)

                            st.session_state.rag_chat_messages.append({
                                "role": "assistant",
                                "content": response_text,
                                "citations": citations,
                                "context": context
                            })
                            
                            # --- Auto-Save History ---
                            current_id = st.session_state.current_drawing_session_id
                            current_title = st.session_state.drawing_chat_history_data[current_id]["title"]
                            if current_title == "ìƒˆë¡œìš´ ëŒ€í™”" and len(st.session_state.rag_chat_messages) > 0:
                                new_title = get_session_title(st.session_state.rag_chat_messages)
                                st.session_state.drawing_chat_history_data[current_id]["title"] = new_title
                            
                            st.session_state.drawing_chat_history_data[current_id]["messages"] = st.session_state.rag_chat_messages
                            st.session_state.drawing_chat_history_data[current_id]["timestamp"] = datetime.now().isoformat()
                            save_history(DRAWING_HISTORY_FILE, st.session_state.drawing_chat_history_data)
                            
                            st.rerun()


                        except Exception as e:
                            st.error(f"ì˜¤ë¥˜: {e}")
                            import traceback
                            st.code(traceback.format_exc())
elif menu == "ë””ë²„ê·¸ (Debug)":
    st.title("ğŸ•µï¸â€â™‚ï¸ RAG Deep Diagnostic Tool (Integrated)")
    
    # Check if admin
    if user_role != 'admin':
        st.error("Admin access required.")
        st.stop()

    search_manager = get_search_manager()
    blob_service_client = get_blob_service_client()
    container_client = blob_service_client.get_container_client(CONTAINER_NAME)

    # Fetch list of files for selection (Filter for drawings only)
    blob_list = []
    try:
        blobs = container_client.list_blobs()
        for b in blobs:
            # Filter: Must be a file (not folder) AND must be in a 'drawings' folder
            if not b.name.endswith('/') and '/drawings/' in b.name:
                blob_list.append(b.name)
    except Exception as e:
        st.error(f"Failed to list blobs: {e}")
    
    blob_list.sort(key=lambda x: x.split('/')[-1]) # Sort by filename
    
    target_blob = st.selectbox("Select Target File", blob_list)
    
    # Extract filename for search
    if target_blob:
        filename = target_blob.split('/')[-1]
        st.caption(f"Selected Filename for Search: `{filename}`")
    else:
        filename = st.text_input("Target Filename", value="ì œ4ê¶Œ ë„ë©´(ì²­ì£¼).pdf")

    if st.button("Run Diagnostics"):
        st.divider()
        
        # 1. Index Inspection
        st.subheader("1. Index Inspection")
        
        # Search for ALL pages (including chunks like "filename (p.1)")
        try:
            # Search Strategy: Look for documents where metadata_storage_name starts with the filename
            # This will catch both the main file and all page chunks
            import unicodedata
            norm_filename = unicodedata.normalize('NFC', filename)
            
            # Use text search for the filename and then filter client-side
            results = search_manager.search_client.search(
                search_text=f"\"{filename}\"",
                search_mode="all",
                select=["id", "metadata_storage_name", "metadata_storage_path", "project", "content"],
                top=1000  # Increase to capture all pages
            )
            
            # Filter to get documents that start with our filename (including page chunks)
            results = [
                doc for doc in results 
                if unicodedata.normalize('NFC', doc.get('metadata_storage_name', '')).startswith(norm_filename)
            ]
            
        except Exception as e:
            st.warning(f"Search failed ({str(e)}). This might indicate an indexing issue.")
            results = []
        
        st.write(f"Found **{len(results)}** documents in index.")
        
        # Show breakdown by type
        if results:
            main_docs = [d for d in results if '(p.' not in d.get('metadata_storage_name', '')]
            page_docs = [d for d in results if '(p.' in d.get('metadata_storage_name', '')]
            st.caption(f"ğŸ“„ Main file: {len(main_docs)} | ğŸ“‘ Page chunks: {len(page_docs)}")
        
        if results:
            # Analyze First Result
            first = results[0]
            st.json({
                "First Doc ID": first['id'],
                "Name": first['metadata_storage_name'],
                "Path": first['metadata_storage_path'],
                "Project": first['project']
            })

            # 2. Blob Verification
            st.subheader("2. Blob Verification")
            path = first['metadata_storage_path']
            blob_path = None
            
            if "https://direct_fetch/" in path:
                st.warning("âš ï¸ Using 'direct_fetch' scheme. This is a virtual path.")
                blob_path = path.replace("https://direct_fetch/", "").split('#')[0]
            elif CONTAINER_NAME in path:
                try:
                    blob_path = path.split(f"/{CONTAINER_NAME}/")[1].split('#')[0]
                    blob_path = urllib.parse.unquote(blob_path)
                except:
                    pass
            
            if blob_path:
                st.write(f"**Extracted Blob Path:** `{blob_path}`")
                blob_client = container_client.get_blob_client(blob_path)
                if blob_client.exists():
                    st.success("âœ… Blob exists in storage.")
                else:
                    st.error("âŒ Blob DOES NOT exist at this path!")
                    
                    # Search for it
                    st.write("Searching for file in container...")
                    found_blobs = list(container_client.list_blobs(name_starts_with=os.path.dirname(blob_path)))
                    if found_blobs:
                        st.write("Found similar blobs:")
                        for b in found_blobs:
                            st.code(b.name)
                    else:
                        st.warning("No similar blobs found.")
            else:
                st.error("Could not extract blob path from metadata.")

            # 3. List Page Check
            st.subheader("3. List Page Check")
            list_keywords = ["PIPING AND INSTRUMENT DIAGRAM FOR LIST", "DRAWING LIST", "ë„ë©´ ëª©ë¡"]
            found_list = False
            
            for doc in results:
                # Handle None content safely
                content = doc.get('content')
                if content is None:
                    content = ""
                    st.warning(f"âš ï¸ Document '{doc['metadata_storage_name']}' has NO CONTENT (NULL).")
                
                content_upper = content.upper()
                if any(k in content_upper for k in list_keywords):
                    st.success(f"âœ… Found List Page! Name: `{doc['metadata_storage_name']}`")
                    st.text_area("Content Preview", content[:500], height=150)
                    found_list = True
                    break
            
            if not found_list:
                st.error("âŒ List Page NOT found in the top 50 results.")
                st.write("Top 5 Results Content Snippets:")
                for i, doc in enumerate(results[:5]):
                    content_preview = (doc.get('content') or "")[:100]
                    st.text(f"{i+1}. {doc['metadata_storage_name']}: {content_preview}...")

            # 4. Cleanup Tool
            st.divider()
            st.subheader("4. Index Cleanup")
            st.warning("If this document is corrupt (No Content / No Project), you can delete it here.")
            
            if st.button(f"ğŸ—‘ï¸ Delete ALL {len(results)} found documents from Index"):
                try:
                    # Collect IDs
                    ids_to_delete = [{"id": doc['id']} for doc in results]
                    search_manager.search_client.delete_documents(documents=ids_to_delete)
                    st.success(f"Successfully deleted {len(results)} documents.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Delete failed: {e}")

        else:
            st.error("No documents found in index matching this filename.")
            
            # Debug: List what IS in the index (drawings only)
            st.divider()
            st.subheader("ğŸ•µï¸ Index Content Peek (Top 20 Drawings)")
            try:
                # Get top 20 docs from drawings_analysis project
                peek_results = search_manager.search_client.search(
                    search_text="*",
                    filter="project eq 'drawings_analysis'",
                    select=["metadata_storage_name", "project", "metadata_storage_last_modified"],
                    top=20
                )
                peek_list = list(peek_results)
                if peek_list:
                    st.write(f"Index contains at least {len(peek_list)} documents. Here are the top 20:")
                    peek_data = []
                    for d in peek_list:
                        peek_data.append({
                            "Name": d.get('metadata_storage_name'),
                            "Project": d.get('project'),
                            "Modified": d.get('metadata_storage_last_modified')
                        })
                    st.table(peek_data)
                else:
                    st.error("âš ï¸ The Index appears to be COMPLETELY EMPTY.")
            except Exception as e:
                st.error(f"Failed to peek index: {e}")

    # -----------------------------
    # ë””ë²„ê¹… ë„êµ¬ (Debug Tools)
    # -----------------------------
    if user_role == 'admin':
        with st.expander("ğŸ› ï¸ ì¸ë±ìŠ¤ ë° ê²€ìƒ‰ ì§„ë‹¨ (Debug Tools)", expanded=False):
            st.warning("ì´ ë„êµ¬ëŠ” ê²€ìƒ‰ ë¬¸ì œë¥¼ ì§„ë‹¨í•˜ê¸° ìœ„í•œ ê²ƒì…ë‹ˆë‹¤.")
            
            # Secret Inspector
            st.write("### ğŸ” ìê²© ì¦ëª… í™•ì¸ (Secret Inspector)")
            def mask_secret(s):
                if not s: return "Not Set"
                if len(s) <= 8: return "*" * len(s)
                return s[:4] + "*" * (len(s)-8) + s[-4:]
            
            secrets_to_check = {
                "AZURE_STORAGE_CONNECTION_STRING": STORAGE_CONN_STR,
                "AZURE_BLOB_CONTAINER_NAME": CONTAINER_NAME,
                "AZURE_OPENAI_ENDPOINT": AZURE_OPENAI_ENDPOINT,
                "AZURE_OPENAI_KEY": AZURE_OPENAI_KEY,
                "AZURE_SEARCH_ENDPOINT": SEARCH_ENDPOINT,
                "AZURE_SEARCH_KEY": SEARCH_KEY,
                "AZURE_TRANSLATOR_KEY": TRANSLATOR_KEY,
                "AZURE_DOC_INTEL_ENDPOINT": AZURE_DOC_INTEL_ENDPOINT,
                "AZURE_DOC_INTEL_KEY": AZURE_DOC_INTEL_KEY
            }
            
            import pandas as pd
            secret_data = []
            for k, v in secrets_to_check.items():
                secret_data.append({"Secret Key": k, "Status": "âœ… Loaded" if v else "âŒ Missing", "Value (Masked)": mask_secret(v)})
            
            st.table(pd.DataFrame(secret_data))
            
            st.write("---")
            
            if st.button("ğŸ” ì¸ë±ìŠ¤ ìƒíƒœ ë° ê²€ìƒ‰ í…ŒìŠ¤íŠ¸ ì‹¤í–‰"):
                try:
                    search_manager = get_search_manager()
                    client = search_manager.search_client
                    
                    st.write("### 1. ì¸ë±ìŠ¤ ë¬¸ì„œ í™•ì¸ (project='drawings_analysis')")
                    results = client.search(search_text="*", filter="project eq 'drawings_analysis'", select=["id", "metadata_storage_name", "project"], top=20)
                    
                    docs = list(results)
                    st.write(f"Found {len(docs)} docs with project='drawings_analysis'")
                    
                    if docs:
                        for doc in docs:
                            st.code(f"ID: {doc['id']}\nName: {doc['metadata_storage_name']}\nProject: {doc['project']}")
                    
                    st.write("---")
                    st.write("### 1-B. ì¸ë±ìŠ¤ ë¬¸ì„œ í™•ì¸ (ì „ì²´ - í•„í„° ì—†ìŒ)")
                    results_all = client.search(search_text="*", select=["id", "metadata_storage_name", "project"], top=20)
                    docs_all = list(results_all)
                    st.write(f"Found {len(docs_all)} docs in total (top 20)")
                    for doc in docs_all:
                        proj = doc.get('project', 'None')
                        st.code(f"Name: {doc['metadata_storage_name']}\nProject: {proj}")
                    
                    st.write("---")
                    st.write("### 2. í‚¤ì›Œë“œ ê²€ìƒ‰ í…ŒìŠ¤íŠ¸ ('foundation loading data')")
                    search_results = client.search(search_text="foundation loading data", filter="project eq 'drawings_analysis'", top=5, select=["metadata_storage_name", "content"])
                    search_docs = list(search_results)
                    
                    st.write(f"ê²€ìƒ‰ ê²°ê³¼: {len(search_docs)}ê°œ")
                    for doc in search_docs:
                        st.text(f"Match: {doc['metadata_storage_name']}")
                        st.caption(f"Content: {doc['content'][:200]}...")
                    
                    st.write("---")
                    st.write("### 3. ì™€ì¼ë“œì¹´ë“œ ê²€ìƒ‰ í…ŒìŠ¤íŠ¸ ('*')")
                    wild_results = client.search(search_text="*", filter="project eq 'drawings_analysis'", top=5, select=["metadata_storage_name", "content"])
                    wild_docs = list(wild_results)
                    
                    st.write(f"ê²€ìƒ‰ ê²°ê³¼: {len(wild_docs)}ê°œ")
                    for doc in wild_docs:
                        st.text(f"Match: {doc['metadata_storage_name']}")
                        st.caption(f"Content: {doc['content'][:200]}...")
                        
                except Exception as e:
                    st.error(f"ì§„ë‹¨ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {str(e)}")
                    st.code(str(e))
            
            st.write("---")
            st.write("### ğŸ” ì¸ë±ìŠ¤ ë°ì´í„° í™•ì¸")
            if st.button("ğŸ“‘ ì¸ë±ìŠ¤ëœ ëª¨ë“  íŒŒì¼ëª… ë³´ê¸°"):
                with st.spinner("ì¸ë±ìŠ¤ ì¡°íšŒ ì¤‘..."):
                    try:
                        search_manager = get_search_manager()
                        # Get all docs (limit to top 1000 to be safe)
                        results = search_manager.search("*", select=["metadata_storage_name"], top=1000)
                        indexed_files = set()
                        for res in results:
                            # Remove page suffix (p.N) to get base filename
                            name = res['metadata_storage_name']
                            base_name = name.split(' (p.')[0]
                            indexed_files.add(base_name)
                        
                        st.write(f"ì´ {len(indexed_files)}ê°œì˜ íŒŒì¼ì´ ì¸ë±ìŠ¤ì—ì„œ ë°œê²¬ë˜ì—ˆìŠµë‹ˆë‹¤.")
                        st.dataframe(list(indexed_files), use_container_width=True)
                    except Exception as e:
                        st.error(f"ì¡°íšŒ ì‹¤íŒ¨: {e}")

            st.write("---")
            st.write("### ğŸ§ª ì‚¬ìš©ì ì§€ì • ê²€ìƒ‰ í…ŒìŠ¤íŠ¸")
            debug_query = st.text_input("ê²€ìƒ‰ì–´ ì…ë ¥ (ì˜ˆ: filter element)", key="debug_query")
            if st.button("ê²€ìƒ‰ í…ŒìŠ¤íŠ¸ ì‹¤í–‰", key="run_debug_search"):
                if debug_query:
                    try:
                        search_manager = get_search_manager()
                        client = search_manager.search_client
                        
                        st.write(f"Query: '{debug_query}'")
                        # Use 'any' search mode to match behavior
                        results = client.search(
                            search_text=debug_query, 
                            filter="project eq 'drawings_analysis'", 
                            search_mode="any",
                            select=["metadata_storage_name", "content"],
                            top=10
                        )
                        docs = list(results)
                        st.write(f"ê²€ìƒ‰ ê²°ê³¼: {len(docs)}ê°œ")
                        
                        if docs:
                            for doc in docs:
                                st.text(f"Match: {doc['metadata_storage_name']}")
                                st.caption(f"Content: {doc['content'][:200]}...")
                        else:
                            st.warning("ê²€ìƒ‰ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤.")
                    except Exception as e:
                        st.error(f"ê²€ìƒ‰ ì˜¤ë¥˜: {e}")

            st.write("---")
            st.write("### âš ï¸ ì¸ë±ìŠ¤ ì´ˆê¸°í™”")
            if st.button("ğŸ—‘ï¸ ëª¨ë“  ë„ë©´ ë°ì´í„° ì‚­ì œ (Index & Blob)", type="primary"):
                try:
                    # 1. Delete all blobs in any drawings/, json/ folder (Global reset)
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    
                    # List all blobs and filter for drawings or json
                    blobs = container_client.list_blobs()
                    deleted_blobs = 0
                    for blob in blobs:
                        if '/drawings/' in blob.name or blob.name.startswith('drawings/') or '/json/' in blob.name or blob.name.startswith('json/'):
                            container_client.delete_blob(blob.name)
                            deleted_blobs += 1
                    
                    # 2. Delete all docs in index with project='drawings_analysis'
                    search_manager = get_search_manager()
                    
                    deleted_total = 0
                    while True:
                        results = search_manager.search_client.search(
                            search_text="*",
                            filter="project eq 'drawings_analysis'",
                            select=["id"],
                            top=1000
                        )
                        ids_to_delete = [{"id": doc['id']} for doc in results]
                        if not ids_to_delete:
                            break
                            
                        search_manager.search_client.delete_documents(documents=ids_to_delete)
                        deleted_total += len(ids_to_delete)
                        if len(ids_to_delete) < 1000:
                            break
                    
                    st.success(f"ëª¨ë“  ë„ë©´ ë°ì´í„°ê°€ ì‚­ì œë˜ì—ˆìŠµë‹ˆë‹¤. (Blob ì‚­ì œ ì™„ë£Œ, Index {deleted_total}ê°œ ì‚­ì œ ì™„ë£Œ) ì´ì œ íŒŒì¼ì„ ë‹¤ì‹œ ì—…ë¡œë“œí•˜ì„¸ìš”.")
                    st.rerun()
                except Exception as e:
                    st.error(f"ì´ˆê¸°í™” ì‹¤íŒ¨: {e}")

            if st.button("ğŸ§¹ 'í˜ì´ì§€ ë²ˆí˜¸ ì—†ëŠ”' ì¤‘ë³µ ë°ì´í„° ì •ë¦¬ (ê¶Œì¥)", help="ì¸ë±ìŠ¤ì—ì„œ (p.N) í˜•ì‹ì´ ì•„ë‹Œ ì˜ëª»ëœ ë°ì´í„°ë¥¼ ì°¾ì•„ ì‚­ì œí•©ë‹ˆë‹¤."):
                try:
                    search_manager = get_search_manager()
                    results = search_manager.search_client.search(
                        search_text="*",
                        filter="project eq 'drawings_analysis'",
                        select=["id", "metadata_storage_name"],
                        top=1000
                    )
                    
                    ids_to_delete = []
                    count = 0
                    for doc in results:
                        name = doc['metadata_storage_name']
                        # Delete if it doesn't contain "(p." (standard page suffix)
                        if "(p." not in name:
                            ids_to_delete.append({"id": doc['id']})
                            count += 1
                    
                    if ids_to_delete:
                        search_manager.search_client.delete_documents(documents=ids_to_delete)
                        st.success(f"ì •ë¦¬ ì™„ë£Œ! {count}ê°œì˜ ì¤‘ë³µ/ì˜ëª»ëœ ë¬¸ì„œë¥¼ ì‚­ì œí–ˆìŠµë‹ˆë‹¤.")
                        st.rerun()
                    else:
                        st.info("ì‚­ì œí•  ì˜ëª»ëœ ë°ì´í„°ê°€ ì—†ìŠµë‹ˆë‹¤. ì¸ë±ìŠ¤ê°€ ê¹¨ë—í•©ë‹ˆë‹¤! âœ¨")
                        
                except Exception as e:
                    st.error(f"ì •ë¦¬ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")

            if st.button("ğŸ·ï¸ ëˆ„ë½ëœ 'drawings_analysis' íƒœê·¸ ë³µêµ¬", help="ë“œë¡œì‰ í´ë”ì— ìˆì§€ë§Œ í”„ë¡œì íŠ¸ íƒœê·¸ê°€ ì—†ëŠ” ë¬¸ì„œë¥¼ ì°¾ì•„ íƒœê·¸ë¥¼ ì¶”ê°€í•©ë‹ˆë‹¤."):
                try:
                    search_manager = get_search_manager()
                    # Search for all docs with missing project tag and filter path in Python
                    results = search_manager.search_client.search(
                        search_text="*",
                        filter="(project eq null)",
                        select=["id", "metadata_storage_name", "metadata_storage_path", "content", "content_exact", "metadata_storage_last_modified", "metadata_storage_size", "metadata_storage_content_type"],
                        top=10000 # Increase to cover all docs
                    )
                    
                    docs_to_fix = []
                    for doc in results:
                        # Filter by path in Python
                        if '/drawings/' in doc.get('metadata_storage_path', ''):
                            doc['project'] = 'drawings_analysis'
                            docs_to_fix.append(doc)
                    
                    if docs_to_fix:
                        success, msg = search_manager.upload_documents(docs_to_fix)
                        if success:
                            st.success(f"ë³µêµ¬ ì™„ë£Œ! {len(docs_to_fix)}ê°œì˜ ë¬¸ì„œì— 'drawings_analysis' íƒœê·¸ë¥¼ ì¶”ê°€í–ˆìŠµë‹ˆë‹¤.")
                            st.rerun()
                        else:
                            st.error(f"ë³µêµ¬ ì‹¤íŒ¨: {msg}")
                    else:
                        st.info("íƒœê·¸ë¥¼ ë³µêµ¬í•  ë¬¸ì„œê°€ ì—†ìŠµë‹ˆë‹¤.")
                except Exception as e:
                    st.error(f"íƒœê·¸ ë³µêµ¬ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")

            if st.button("ğŸ“Š ì¸ë±ìŠ¤ í†µê³„ í™•ì¸", help="í”„ë¡œì íŠ¸ë³„ ë¬¸ì„œ ê°œìˆ˜ë¥¼ í™•ì¸í•©ë‹ˆë‹¤."):
                try:
                    search_manager = get_search_manager()
                    
                    # Count drawings_analysis
                    drawings_res = search_manager.search_client.search(
                        search_text="*",
                        filter="project eq 'drawings_analysis'",
                        include_total_count=True,
                        top=0
                    )
                    drawings_count = drawings_res.get_count()
                    
                    # Count others (likely standard indexed)
                    others_res = search_manager.search_client.search(
                        search_text="*",
                        filter="project eq null",
                        include_total_count=True,
                        top=0
                    )
                    others_count = others_res.get_count()
                    
                    st.write(f"**ë„ë©´ ë¶„ì„ ë°ì´í„° (drawings_analysis):** {drawings_count}ê°œ")
                    st.write(f"**ì¼ë°˜ ë¬¸ì„œ ë°ì´í„° (Standard Indexer):** {others_count}ê°œ")
                    
                    # Check Standard Indexer Status
                    st.divider()
                    st.write("**í‘œì¤€ ì¸ë±ì„œ (Standard Indexer) ìƒíƒœ í™•ì¸:**")
                    # Try common indexer names
                    for idx_name in ["pdf-indexer", "indexer-all", "indexer-drawings"]:
                        try:
                            status = search_manager.indexer_client.get_indexer_status(idx_name)
                            last_res = status.last_result
                            if last_res:
                                st.write(f"- `{idx_name}`: {last_res.status} (ì„±ê³µ: {last_res.item_count}, ì‹¤íŒ¨: {last_res.failed_item_count})")
                                if last_res.failed_item_count > 0:
                                    with st.expander(f"âŒ {idx_name} ì—ëŸ¬ ìƒì„¸ ë³´ê¸°"):
                                        for err in last_res.errors[:5]:
                                            st.error(f"ë¬¸ì„œ: {err.key}\nì—ëŸ¬: {err.message}")
                            else:
                                st.write(f"- `{idx_name}`: ì‹¤í–‰ ê¸°ë¡ ì—†ìŒ")
                        except:
                            pass

                    if drawings_count == 0 and others_count > 0:
                        st.warning("ë„ë©´ ë°ì´í„°ê°€ í•˜ë‚˜ë„ ì—†ìŠµë‹ˆë‹¤. ì¸ë±ì‹± ê³¼ì •ì— ë¬¸ì œê°€ ìˆì„ ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
                except Exception as e:
                    st.error(f"í†µê³„ í™•ì¸ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")

            with st.expander("ğŸ” ì¸ë±ìŠ¤ ìƒì„¸ ì§„ë‹¨ ë„êµ¬", expanded=False):
                st.caption("ì¸ë±ìŠ¤ì— ì €ì¥ëœ ì‹¤ì œ íŒŒì¼ëª…ê³¼ íƒœê·¸ë¥¼ ì§ì ‘ í™•ì¸í•©ë‹ˆë‹¤.")
                
                # Add search input for specific file diagnosis (Outside button for persistence)
                diag_query = st.text_input("ì§„ë‹¨í•  íŒŒì¼ëª… ê²€ìƒ‰ (ì¼ë¶€ë§Œ ì…ë ¥ ê°€ëŠ¥)", value="", key="diag_query")
                diag_path_filter = st.checkbox("'/drawings/' ê²½ë¡œë§Œ ë³´ê¸°", value=True, key="diag_path_filter")
                
                if st.button("ğŸ“‹ ì§„ë‹¨ ì‹¤í–‰ (ìµœê·¼ 100ê°œ)"):
                    try:
                        search_manager = get_search_manager()
                        
                        # Use a more inclusive search for diagnosis
                        # If query is provided, use it as search_text. If not, use *
                        results = search_manager.search_client.search(
                            search_text=diag_query if diag_query else "*",
                            select=["metadata_storage_name", "project", "metadata_storage_path"],
                            top=1000 # Increase for better diagnosis
                        )
                        
                        dump_data = []
                        for doc in results:
                            name = doc.get('metadata_storage_name', '')
                            path = doc.get('metadata_storage_path', '')
                            
                            if diag_path_filter and '/drawings/' not in path:
                                continue
                                
                            dump_data.append({
                                "Name": name,
                                "Project": doc.get('project'),
                                "Path": path
                            })
                        
                        if dump_data:
                            st.write(f"ê²€ìƒ‰ ê²°ê³¼: {len(dump_data)}ê°œì˜ ë¬¸ì„œ ë°œê²¬")
                            st.table(dump_data)
                        else:
                            st.warning("ê²€ìƒ‰ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤. íŒŒì¼ëª…ì´ ì¸ë±ìŠ¤ì— ì¡´ì¬í•˜ì§€ ì•Šê±°ë‚˜ í•„í„°ì— ê±¸ëŸ¬ì¡Œì„ ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
                            
                        # Extra check: Search by path only if query failed
                        if diag_query and not dump_data:
                            st.info(f"'{diag_query}'ë¡œ ê²€ìƒ‰ëœ ê²°ê³¼ê°€ ì—†ì–´ ê²½ë¡œ ê¸°ë°˜ìœ¼ë¡œ ë‹¤ì‹œ ì°¾ìŠµë‹ˆë‹¤...")
                            # Use startswith on metadata_storage_path (SimpleField/Filterable)
                            # We don't know the full prefix, but we can try to find anything in drawings
                            path_results = search_manager.search_client.search(
                                search_text="*",
                                filter="startswith(metadata_storage_path, 'https://')", # Broad filter
                                select=["metadata_storage_name", "project", "metadata_storage_path"],
                                top=5000 # Increase to cover more docs
                            )
                            # Filter for '/drawings/' in Python for maximum reliability
                            path_data = [
                                {"Name": d['metadata_storage_name'], "Project": d['project'], "Path": d['metadata_storage_path']} 
                                for d in path_results 
                                if '/drawings/' in d.get('metadata_storage_path', '')
                            ]
                            if path_data:
                                st.write("'/drawings/' ê²½ë¡œì—ì„œ ë°œê²¬ëœ íŒŒì¼ë“¤ (ìµœê·¼ 100ê°œ ì¤‘):")
                                st.table(path_data[:20])
                            else:
                                st.error("'/drawings/' ê²½ë¡œì—ì„œ ë¬¸ì„œë¥¼ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤. ì¸ë±ì„œê°€ í•´ë‹¹ í´ë”ë¥¼ ìŠ¤ìº”í•˜ì§€ ì•Šì•˜ì„ ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
                                
                    except Exception as e:
                        st.error(f"ì§„ë‹¨ ì¤‘ ì˜¤ë¥˜ ë°œìƒ: {e}")
                        results_list = list(results)
                        st.info(f"ê²€ìƒ‰ëœ ì²­í¬(Chunk) ìˆ˜: {len(results_list)}ê°œ")
                        
                        total_chars = 0
                        for i, doc in enumerate(results_list):
                            content = doc.get('content', '')
                            char_count = len(content)
                            total_chars += char_count
                            
                            with st.expander(f"Chunk {i+1}: {doc.get('metadata_storage_name')} ({char_count}ì)"):
                                st.code(content[:1000] + ("..." if len(content) > 1000 else ""))
                        
                        st.divider()
                        st.metric("ì´ ê¸€ì ìˆ˜ (Total Characters)", f"{total_chars:,}")
                        est_tokens = int(total_chars / 4)
                        st.metric("ì˜ˆìƒ í† í° ìˆ˜ (Estimated Tokens)", f"{est_tokens:,}")
                        
                        if est_tokens > 5000:
                            st.warning(f"âš ï¸ í† í° ìˆ˜ê°€ ë§ìŠµë‹ˆë‹¤ ({est_tokens} > 5000). AI ë‹µë³€ ìƒì„± ì‹œ 'Token Limit Exceeded' ì˜¤ë¥˜ê°€ ë°œìƒí•  ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
                        else:
                            st.success(f"âœ… í† í° ìˆ˜ê°€ ì ì ˆí•©ë‹ˆë‹¤ ({est_tokens}).")
                            
                    except Exception as e:
                        st.error(f"ë¶„ì„ ì‹¤íŒ¨: {e}")
            
            if st.button("ğŸ—‘ï¸ ëŒ€í™” ì´ˆê¸°í™”", key="clear_rag_chat"):
                st.session_state.rag_chat_messages = []
                st.rerun()

    st.markdown("---")

if menu == "ì—‘ì…€ë°ì´í„° ìë™ì¶”ì¶œ":
    # Integrated Excel Tool
    excel_manager.render_excel_tool()

if menu == "ì‚¬ì§„ëŒ€ì§€ ìë™ì‘ì„±":
    st.caption("ê±´ì„¤ í˜„ì¥ ì‚¬ì§„ì„ ì—…ë¡œë“œí•˜ì—¬ Excel ì‚¬ì§„ëŒ€ì§€ë¥¼ ìë™ìœ¼ë¡œ ìƒì„±í•©ë‹ˆë‹¤.")
    
    # Embed Photo Log app via iframe
    st.components.v1.iframe(
        src="https://photo-log-a0215.web.app/",
        height=800,
        scrolling=True
    )

if menu == "ì‘ì—…ê³„íš ë° íˆ¬ì…ë¹„ ìë™ì‘ì„±":
    st.caption("ì‘ì—… ê³„íšì„ ìˆ˜ë¦½í•˜ê³  íˆ¬ì…ë¹„ë¥¼ ìë™ìœ¼ë¡œ ì‚°ì¶œí•©ë‹ˆë‹¤.")
    
    # Embed Work Schedule app via iframe
    st.components.v1.iframe(
        src="https://workschedule-7b1cf.web.app/",
        height=800,
        scrolling=True
    )

if menu == "ê´€ë¦¬ì ì„¤ì •":
    # st.subheader("âš™ï¸ ê´€ë¦¬ì ì„¤ì •") - Removed to avoid duplication
    st.info("Azure AI Search ë¦¬ì†ŒìŠ¤ë¥¼ ì´ˆê¸°í™”í•˜ê±°ë‚˜ ìƒíƒœë¥¼ í™•ì¸í•©ë‹ˆë‹¤.")
    
    # ì¸ë±ì‹± ëŒ€ìƒ í´ë” ì„¤ì •
    # í´ë” ëª©ë¡ ê°€ì ¸ì˜¤ê¸°
    folder_options = ["(ì „ì²´)"]
    try:
        blob_service_client = get_blob_service_client()
        container_client = blob_service_client.get_container_client(CONTAINER_NAME)
        # walk_blobsë¥¼ ì‚¬ìš©í•˜ì—¬ ìµœìƒìœ„ í´ë”ë§Œ ì¡°íšŒ
        for blob in container_client.walk_blobs(delimiter='/'):
            if blob.name.endswith('/'):
                folder_options.append(blob.name.strip('/'))
    except Exception as e:
        st.warning(f"í´ë” ëª©ë¡ì„ ê°€ì ¸ì˜¤ì§€ ëª»í–ˆìŠµë‹ˆë‹¤: {e}")
        folder_options.append("GULFLNG") # Fallback

    # ê¸°ë³¸ê°’ ì„¤ì • (GULFLNGê°€ ìˆìœ¼ë©´ ê·¸ê±¸ë¡œ, ì—†ìœ¼ë©´ ì „ì²´)
    default_idx = 0
    if "GULFLNG" in folder_options:
        default_idx = folder_options.index("GULFLNG")

    selected_folder = st.selectbox(
        "ì¸ë±ì‹± ëŒ€ìƒ í´ë” ì„ íƒ", 
        folder_options, 
        index=default_idx,
        help="ì¸ë±ì‹±í•  í”„ë¡œì íŠ¸ í´ë”ë¥¼ ì„ íƒí•˜ì„¸ìš”."
    )
    
    
    # '(ì „ì²´)' ì„ íƒ ì‹œ Noneìœ¼ë¡œ ì²˜ë¦¬
    target_folder = None if selected_folder == "(ì „ì²´)" else selected_folder
    
    st.info("ğŸ’¡ **í´ë”ë³„ ì¸ë±ì‹±**: ê° í´ë”ëŠ” ë…ë¦½ì ìœ¼ë¡œ ì¸ë±ì‹±ë©ë‹ˆë‹¤. ë‹¤ë¥¸ í´ë”ì˜ ë°ì´í„°ì— ì˜í–¥ì„ ì£¼ì§€ ì•ŠìŠµë‹ˆë‹¤.")
    
    # ------------------------------------------------------------------
    # ì¸ë±ìŠ¤ ìŠ¤í‚¤ë§ˆ ì—…ë°ì´íŠ¸ ë²„íŠ¼
    # ------------------------------------------------------------------
    st.subheader("ğŸ“ ì¸ë±ìŠ¤ ìŠ¤í‚¤ë§ˆ ê´€ë¦¬")
    st.markdown("""
    ë„ë©´ ë©”íƒ€ë°ì´í„° í•„ë“œ(`title`, `drawing_no`)ë¥¼ ì¸ë±ìŠ¤ì— ì¶”ê°€í•©ë‹ˆë‹¤.
    
    **ì£¼ì˜:** Azure SearchëŠ” ê¸°ì¡´ í•„ë“œì˜ íƒ€ì…ì„ ë³€ê²½í•  ìˆ˜ ì—†ìŠµë‹ˆë‹¤.  
    ìŠ¤í‚¤ë§ˆ ì¶©ëŒ ì‹œ ì¸ë±ìŠ¤ë¥¼ ì‚­ì œí•˜ê³  ì¬ìƒì„±í•´ì•¼ í•©ë‹ˆë‹¤.
    """)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("##### ğŸ”„ ìŠ¤í‚¤ë§ˆ ì—…ë°ì´íŠ¸ (ê¸°ì¡´ ë°ì´í„° ìœ ì§€)")
        if st.button("ìŠ¤í‚¤ë§ˆ ì—…ë°ì´íŠ¸", help="ìƒˆ í•„ë“œë§Œ ì¶”ê°€í•©ë‹ˆë‹¤. ê¸°ì¡´ í•„ë“œ ë³€ê²½ ì‹œ ì—ëŸ¬ê°€ ë°œìƒí•  ìˆ˜ ìˆìŠµë‹ˆë‹¤."):
            with st.spinner("ì¸ë±ìŠ¤ ìŠ¤í‚¤ë§ˆ ì—…ë°ì´íŠ¸ ì¤‘..."):
                manager = get_search_manager()
                success, msg = manager.create_index()
                
                if success:
                    st.success(f"âœ… {msg}")
                    st.info("""
                    **ë‹¤ìŒ ë‹¨ê³„:**
                    1. ê¸°ì¡´ íŒŒì¼ì˜ `ğŸ”„` ë²„íŠ¼ì„ í´ë¦­í•´ ì¬ë¶„ì„í•˜ê±°ë‚˜
                    2. ìƒˆ íŒŒì¼ì„ ì—…ë¡œë“œí•˜ì„¸ìš”
                    
                    **ìƒˆë¡œ ì¶”ê°€ëœ í•„ë“œ:**
                    - `title` (ë„ë©´ëª…)
                    - `drawing_no` (ë„ë©´ë²ˆí˜¸)
                    """)
                else:
                    st.error(f"âŒ ìŠ¤í‚¤ë§ˆ ì—…ë°ì´íŠ¸ ì‹¤íŒ¨: {msg}")
                    st.warning("âš ï¸ í•„ë“œ ì¶©ëŒì´ ë°œìƒí–ˆë‹¤ë©´ ì˜¤ë¥¸ìª½ì˜ 'ì¸ë±ìŠ¤ ì‚­ì œ í›„ ì¬ìƒì„±'ì„ ì‚¬ìš©í•˜ì„¸ìš”.")
    
    with col2:
        st.markdown("##### ğŸ—‘ï¸ ì‚­ì œ í›„ ì¬ìƒì„± (ëª¨ë“  ë°ì´í„° ì‚­ì œ)")
        st.warning("âš ï¸ **ê²½ê³ :** ëª¨ë“  ì¸ë±ìŠ¤ ë°ì´í„°ê°€ ì‚­ì œë©ë‹ˆë‹¤!")
        confirm_delete = st.checkbox("ì‚­ì œ í™•ì¸ (ëª¨ë“  ë°ì´í„° ì‚­ì œ)", key="confirm_delete_index")
        if st.button("ì¸ë±ìŠ¤ ì‚­ì œ í›„ ì¬ìƒì„±", disabled=not confirm_delete, help="ì¸ë±ìŠ¤ë¥¼ ì™„ì „íˆ ì‚­ì œí•˜ê³  ìƒˆë¡œ ìƒì„±í•©ë‹ˆë‹¤. íŒŒì¼ì„ ë‹¤ì‹œ ì—…ë¡œë“œí•´ì•¼ í•©ë‹ˆë‹¤."):
            with st.spinner("ì¸ë±ìŠ¤ ì‚­ì œ ë° ì¬ìƒì„± ì¤‘..."):
                manager = get_search_manager()
                
                # 1. ì¸ë±ìŠ¤ ì‚­ì œ
                del_success, del_msg = manager.delete_index()
                if del_success:
                    st.info(f"ğŸ—‘ï¸ {del_msg}")
                    time.sleep(2)  # ì‚­ì œ ì™„ë£Œ ëŒ€ê¸°
                    
                    # 2. ì¸ë±ìŠ¤ ì¬ìƒì„±
                    create_success, create_msg = manager.create_index()
                    if create_success:
                        st.success(f"âœ… {create_msg}")
                        st.success("""
                        **ì¸ë±ìŠ¤ê°€ ì„±ê³µì ìœ¼ë¡œ ì¬ìƒì„±ë˜ì—ˆìŠµë‹ˆë‹¤!**
                        
                        ì´ì œ íŒŒì¼ì„ ë‹¤ì‹œ ì—…ë¡œë“œí•˜ê±°ë‚˜ ì¬ë¶„ì„í•˜ì„¸ìš”.
                        """)
                    else:
                        st.error(f"âŒ ì¬ìƒì„± ì‹¤íŒ¨: {create_msg}")
                else:
                    st.error(f"âŒ ì‚­ì œ ì‹¤íŒ¨: {del_msg}")
    
    st.divider()
    
    confirm_reset = st.checkbox("ìœ„ í´ë”ë¥¼ ì¸ë±ì‹± ëŒ€ìƒìœ¼ë¡œ ì„¤ì •í•˜ê³  ì‹¶ìŠµë‹ˆë‹¤.", key="confirm_reset")
    
    if st.button("ğŸš€ í´ë” ì¸ë±ì‹± ì„¤ì • (Data Source, Indexer)", disabled=not confirm_reset):
        with st.spinner("ë¦¬ì†ŒìŠ¤ ìƒì„± ì¤‘..."):
            manager = get_search_manager()
            
            # 1. Index í™•ì¸/ìƒì„± (í•œë²ˆë§Œ í•„ìš”)
            st.write("1. Index í™•ì¸ ì¤‘...")
            success, msg = manager.create_index()
            if success:
                st.success(msg)
            else:
                st.error(msg)
                
            # 2. Data Source (í´ë”ë³„)
            st.write(f"2. Data Source ìƒì„± ì¤‘... (í´ë”: {selected_folder})")
            success, msg, datasource_name = manager.create_data_source(
                SEARCH_DATASOURCE_NAME, 
                STORAGE_CONN_STR, 
                CONTAINER_NAME, 
                query=target_folder,
                folder_name=target_folder
            )
            if success:
                st.success(msg)
            else:
                st.error(msg)
                st.stop()  # Stop execution if datasource creation fails
                
            # 2.5 Skillset (OCR) - Optional
            skillset_name = None
            enable_ocr = st.checkbox("ğŸ“¸ OCR(ì´ë¯¸ì§€ í…ìŠ¤íŠ¸ ì¶”ì¶œ) í™œì„±í™”", value=False, help="PDF ë„ë©´ì´ë‚˜ ì´ë¯¸ì§€ íŒŒì¼ì—ì„œ í…ìŠ¤íŠ¸ë¥¼ ì¶”ì¶œí•©ë‹ˆë‹¤. Azure AI Services í‚¤ê°€ í•„ìš”í•˜ë©° ë¹„ìš©ì´ ë°œìƒí•  ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
            
            if enable_ocr:
                st.write(f"2.5. Skillset (OCR) ìƒì„± ì¤‘...")
                # Use Translator Key as Cognitive Services Key (assuming it's a multi-service key)
                cog_key = st.secrets.get("AZURE_TRANSLATOR_KEY", os.environ.get("AZURE_TRANSLATOR_KEY"))
                
                if not cog_key:
                    st.warning("âš ï¸ Azure AI Services í‚¤(AZURE_TRANSLATOR_KEY)ê°€ ì„¤ì •ë˜ì§€ ì•Šì•„ OCRì„ ê±´ë„ˆëœë‹ˆë‹¤.")
                else:
                    skillset_name = f"skillset-{target_folder}" if target_folder else "skillset-all"
                    success, msg = manager.create_skillset(skillset_name, cog_key)
                    if success:
                        st.success(msg)
                    else:
                        st.error(f"Skillset ìƒì„± ì‹¤íŒ¨: {msg}")
                        skillset_name = None # Fallback to no skillset
                
            # 3. Indexer (í´ë”ë³„)
            st.write(f"3. Indexer ìƒì„± ì¤‘... (í´ë”: {selected_folder})")
            # ê¸°ì¡´ ì¸ë±ì„œ ì‚­ì œ (ê°™ì€ í´ë”ì˜ ì´ì „ ì„¤ì • ì œê±°)
            manager.delete_indexer(target_folder)
            success, msg, indexer_name = manager.create_indexer(target_folder, datasource_name, skillset_name=skillset_name)
            if success:
                st.success(msg)
                st.info(f"âœ… '{selected_folder}' í´ë”ì— ëŒ€í•œ ì¸ë±ì‹± ì„¤ì •ì´ ì™„ë£Œë˜ì—ˆìŠµë‹ˆë‹¤. ì•„ë˜ 'ì¸ë±ì„œ ìˆ˜ë™ ì‹¤í–‰'ì„ ëˆŒëŸ¬ ì¸ë±ì‹±ì„ ì‹œì‘í•˜ì„¸ìš”.")
            else:
                st.error(msg)
    
    st.divider()
    
    # ------------------------------------------------------------------
    # 4. ì¸ë±ìŠ¤ ë‚´ìš© ì¡°íšŒ (ë””ë²„ê¹…ìš©)
    # ------------------------------------------------------------------
    st.subheader("ğŸ” ì¸ë±ìŠ¤ ë‚´ìš© ì¡°íšŒ (OCR í™•ì¸ìš©)")
    with st.expander("íŠ¹ì • íŒŒì¼ì˜ ì¸ë±ì‹±ëœ ë‚´ìš© í™•ì¸í•˜ê¸°"):
        target_filename = st.text_input("í™•ì¸í•  íŒŒì¼ëª… (ì˜ˆ: drawing.pdf)", help="ì •í™•í•œ íŒŒì¼ëª…ì„ ì…ë ¥í•˜ì„¸ìš”.")
        if st.button("ë‚´ìš© ì¡°íšŒ"):
            if target_filename:
                manager = get_search_manager()
                with st.spinner("ì¡°íšŒ ì¤‘..."):
                    content = manager.get_document_content(target_filename)
                    st.text_area("ì¸ë±ì‹±ëœ ë‚´ìš© (ì•ë¶€ë¶„ 2000ì)", content[:2000], height=300)
            else:
                st.warning("íŒŒì¼ëª…ì„ ì…ë ¥í•˜ì„¸ìš”.")

    st.divider()
    
    # ìˆ˜ë™ ì‹¤í–‰ ì•ˆë‚´ ë° í™•ì¸
    st.info(f"ğŸ“‚ **í˜„ì¬ ì„ íƒëœ í´ë”**: {selected_folder}")
    st.markdown("ìˆ˜ë™ ì¸ë±ì„œ ì‹¤í–‰ì€ ì„ íƒí•œ í´ë”ì˜ ìƒˆ íŒŒì¼ ë˜ëŠ” ë³€ê²½ëœ íŒŒì¼ì„ ê²€ìƒ‰ ì—”ì§„ì— ë°˜ì˜í•©ë‹ˆë‹¤.")
    
    confirm_run = st.checkbox("ìœ„ í´ë”ë¥¼ ì¸ë±ì‹±í•˜ëŠ” ê²ƒì„ í™•ì¸í–ˆìœ¼ë©°, ì§„í–‰í•˜ê³  ì‹¶ìŠµë‹ˆë‹¤.", key="confirm_run")
    
    if st.button("â–¶ï¸ ì¸ë±ì„œ ìˆ˜ë™ ì‹¤í–‰", disabled=not confirm_run):
        manager = get_search_manager()
        success, msg = manager.run_indexer(target_folder)
        if success:
            st.success(msg)
            st.info("ì¸ë±ì‹±ì´ ì‹œì‘ë˜ì—ˆìŠµë‹ˆë‹¤. ì•„ë˜ 'ìƒíƒœ í™•ì¸' ë²„íŠ¼ì„ ëˆŒëŸ¬ ì§„í–‰ ìƒí™©ì„ ëª¨ë‹ˆí„°ë§í•˜ì„¸ìš”.")
        else:
            st.error(msg)
            
    # Add Delete Indexer Button
    if st.button("ğŸ›‘ ì¸ë±ì„œ ì‚­ì œ (ìë™ ì¸ë±ì‹± ì¤‘ì§€)", help="ìë™ìœ¼ë¡œ ì‹¤í–‰ë˜ëŠ” ì¸ë±ì„œë¥¼ ì‚­ì œí•˜ì—¬ ì¤‘ë³µ ì¸ë±ì‹±ì„ ë°©ì§€í•©ë‹ˆë‹¤."):
        manager = get_search_manager()
        indexer_name = f"indexer-{target_folder}" if target_folder else "indexer-all"
        try:
            manager.indexer_client.delete_indexer(indexer_name)
            st.success(f"ì¸ë±ì„œ '{indexer_name}'ê°€ ì‚­ì œë˜ì—ˆìŠµë‹ˆë‹¤. ì´ì œ ìë™ ì¸ë±ì‹±ì´ ì¤‘ì§€ë©ë‹ˆë‹¤.")
        except Exception as e:
            st.error(f"ì¸ë±ì„œ ì‚­ì œ ì‹¤íŒ¨: {e}")
            
    st.divider()
    
    col_status, col_refresh = st.columns([3, 1])
    with col_status:
        st.markdown("### ğŸ“Š ì¸ë±ì‹± í˜„í™© ëª¨ë‹ˆí„°ë§")
    with col_refresh:
        auto_refresh = st.checkbox("ìë™ ìƒˆë¡œê³ ì¹¨ (5ì´ˆ)", value=False)

    # ìƒíƒœ í™•ì¸ ë¡œì§ (ë²„íŠ¼ í´ë¦­ ë˜ëŠ” ìë™ ìƒˆë¡œê³ ì¹¨)
    if st.button("ìƒíƒœ ë° ì§„í–‰ë¥  í™•ì¸") or auto_refresh:
        manager = get_search_manager()
        
        # 1. ì†ŒìŠ¤ íŒŒì¼ ê°œìˆ˜ í™•ì¸ (ì§„í–‰ë¥  ê³„ì‚°ìš©)
        with st.spinner("ì†ŒìŠ¤ íŒŒì¼ ê°œìˆ˜ ê³„ì‚° ì¤‘..."):
            total_blobs = manager.get_source_blob_count(STORAGE_CONN_STR, CONTAINER_NAME, folder_path=target_folder)
        
        # 2. ì¸ë±ì„œ ìƒíƒœ í™•ì¸
        status_info = manager.get_indexer_status(target_folder)
        
        # ìƒíƒœ ì–¸íŒ©
        status = status_info.get("status")
        item_count = status_info.get("item_count", 0)
        failed_count = status_info.get("failed_item_count", 0)
        error_msg = status_info.get("error_message")
        errors = status_info.get("errors", [])
        warnings = status_info.get("warnings", [])
        
        # 3. ì¸ë±ìŠ¤ ë¬¸ì„œ ê°œìˆ˜
        doc_count = manager.get_document_count()
        
        # UI í‘œì‹œ
        st.metric(label="ì´ ì†ŒìŠ¤ íŒŒì¼ ìˆ˜", value=f"{total_blobs}ê°œ")
        
        # ì§„í–‰ë¥  ê³„ì‚° (ì‹¤ì œ ì¸ë±ìŠ¤ëœ ë¬¸ì„œ ìˆ˜ ê¸°ì¤€)
        if total_blobs > 0:
            progress = min(doc_count / total_blobs, 1.0)
        else:
            progress = 0.0
            
        st.progress(progress, text=f"ì¸ë±ì‹± ì§„í–‰ë¥ : {int(progress * 100)}% ({doc_count}/{total_blobs})")
        
        # ìƒíƒœ ë©”ì‹œì§€
        if status == "inProgress":
            st.info(f"â³ ì¸ë±ì‹± ì§„í–‰ ì¤‘... (ì²˜ë¦¬ëœ ë¬¸ì„œ: {item_count}, ì‹¤íŒ¨: {failed_count})")
            if auto_refresh:
                time.sleep(5)
                st.rerun()
        elif status == "success":
            st.success(f"âœ… ì¸ë±ì‹± ì™„ë£Œ! (ì´ ì¸ë±ìŠ¤ ë¬¸ì„œ: {doc_count}ê°œ)")
        elif status == "error":
            st.error(f"âŒ ì¸ë±ì‹± ì˜¤ë¥˜ ë°œìƒ: {error_msg}")
        elif status == "transientFailure":
            st.warning("âš ï¸ ì¼ì‹œì  ì˜¤ë¥˜ ë°œìƒ (ì¬ì‹œë„ ì¤‘...)")
        else:
            st.write(f"ìƒíƒœ: {status}")

        # ì˜¤ë¥˜ ìƒì„¸ í‘œì‹œ
        if failed_count > 0 or errors:
            st.error(f"âŒ ì‹¤íŒ¨í•œ ë¬¸ì„œ: {failed_count}ê°œ")
            with st.expander("ğŸš¨ ì˜¤ë¥˜ ìƒì„¸ ë¡œê·¸ í™•ì¸", expanded=True):
                for err in errors:
                    st.write(f"- {err}")
        
        if warnings:
            with st.expander("âš ï¸ ê²½ê³  ë¡œê·¸ í™•ì¸"):
                for warn in warnings:
                    st.warning(f"- {warn}")
    
    st.divider()
    
    # ------------------------------------------------------------------
    # ğŸ” ë””ë²„ê·¸ íˆ´ - Index Content Peek
    # ------------------------------------------------------------------
    st.subheader("ğŸ” ë””ë²„ê·¸ íˆ´")
    st.markdown("ì¸ë±ìŠ¤ì— ì €ì¥ëœ ë¬¸ì„œë¥¼ í™•ì¸í•˜ì—¬ ìŠ¤í‚¤ë§ˆ í•„ë“œê°€ ì˜¬ë°”ë¥´ê²Œ ì±„ì›Œì¡ŒëŠ”ì§€ ê²€ì¦í•©ë‹ˆë‹¤.")
    
    with st.expander("Index Content Peek", expanded=False):
        st.markdown("ì¸ë±ìŠ¤ì—ì„œ ìµœê·¼ ë¬¸ì„œë¥¼ ê°€ì ¸ì™€ í•„ë“œ ê°’ì„ í™•ì¸í•©ë‹ˆë‹¤.")
        
        # í”„ë¡œì íŠ¸ í•„í„° ì˜µì…˜
        filter_project = st.text_input(
            "í”„ë¡œì íŠ¸ í•„í„° (ì„ íƒì‚¬í•­)", 
            value="drawings_analysis",
            help="íŠ¹ì • í”„ë¡œì íŠ¸ì˜ ë¬¸ì„œë§Œ ì¡°íšŒ (ë¹„ì›Œë‘ë©´ ëª¨ë“  ë¬¸ì„œ)"
        )
        
        peek_limit = st.slider("ì¡°íšŒí•  ë¬¸ì„œ ìˆ˜", min_value=1, max_value=20, value=5)
        
        if st.button("ğŸ“„ Peek Index", key="peek_index_btn"):
            with st.spinner("ì¸ë±ìŠ¤ ì¡°íšŒ ì¤‘..."):
                try:
                    manager = get_search_manager()
                    
                    # Search with filter
                    if filter_project:
                        results = manager.search_client.search(
                            search_text="*",
                            filter=f"project eq '{filter_project}'",
                            top=peek_limit,
                            select=["id", "content", "title", "drawing_no", "page_number", "filename", "metadata_storage_name", "project"]
                        )
                    else:
                        results = manager.search_client.search(
                            search_text="*",
                            top=peek_limit,
                            select=["id", "content", "title", "drawing_no", "page_number", "filename", "metadata_storage_name", "project"]
                        )
                    
                    docs = list(results)
                    
                    if not docs:
                        st.warning("ì¸ë±ìŠ¤ì— ë¬¸ì„œê°€ ì—†ìŠµë‹ˆë‹¤.")
                    else:
                        st.success(f"ì´ {len(docs)}ê°œ ë¬¸ì„œë¥¼ ì°¾ì•˜ìŠµë‹ˆë‹¤.")
                        
                        for i, doc in enumerate(docs):
                            with st.expander(f"ğŸ“„ Document {i+1}: {doc.get('filename', 'N/A')} - Page {doc.get('page_number', 'N/A')}", expanded=(i==0)):
                                # ì¤‘ìš” í•„ë“œ í•˜ì´ë¼ì´íŠ¸
                                col1, col2 = st.columns(2)
                                with col1:
                                    st.markdown("**í•µì‹¬ ë©”íƒ€ë°ì´í„°:**")
                                    st.json({
                                        "title": doc.get("title"),
                                        "drawing_no": doc.get("drawing_no"),
                                        "page_number": doc.get("page_number"),
                                        "filename": doc.get("filename"),
                                        "project": doc.get("project")
                                    })
                                
                                with col2:
                                    st.markdown("**í•„ë“œ ìƒíƒœ ê²€ì¦:**")
                                    title_status = "âœ…" if doc.get("title") else "âŒ"
                                    drawing_status = "âœ…" if doc.get("drawing_no") else "âŒ"
                                    page_status = "âœ…" if doc.get("page_number") is not None else "âŒ"
                                    
                                    st.markdown(f"""
                                    - `title`: {title_status} {doc.get("title") or "NULL"}
                                    - `drawing_no`: {drawing_status} {doc.get("drawing_no") or "NULL"}
                                    - `page_number`: {page_status} {doc.get("page_number") if doc.get("page_number") is not None else "NULL"}
                                    """)
                                
                                # Content preview
                                st.markdown("**Content Preview (ì²˜ìŒ 500ì):**")
                                content_preview = doc.get("content", "")[:500]
                                st.text_area("", content_preview, height=150, key=f"content_{i}", disabled=True)
                                
                                # Full JSON
                                with st.expander("ì „ì²´ JSON ë³´ê¸°"):
                                    st.json(dict(doc))
                
                except Exception as e:
                    st.error(f"ì¸ë±ìŠ¤ ì¡°íšŒ ì‹¤íŒ¨: {e}")
                    import traceback
                    st.code(traceback.format_exc())


if menu == "ì‚¬ìš©ì ì„¤ì •":
    from modules.user_settings_module import render_user_settings
    render_user_settings(auth_manager)



if menu == "ë””ë²„ê·¸ (Debug)":
    st.title("ğŸ” Search Debug Tool (Cloud)")
    st.write("Debug Menu Loaded...") # Debug print
    
    # Secrets (Already loaded in app.py as global variables, but we can reuse get_search_manager)
    try:
        search_manager = get_search_manager()
        st.write("Search Manager Loaded.") # Debug print
    except Exception as e:
        st.error(f"Failed to load Search Manager: {e}")
        st.stop()
    



    # ========================================
    st.header("ğŸ¯ 2ë‹¨ê³„ ê²€ìƒ‰ í…ŒìŠ¤íŠ¸ (ì •í™•í•œ í‚¤ì›Œë“œ ìš°ì„ )")
    st.info("**ëª©ì **: ì‚¬ìš©ì ì…ë ¥ ê·¸ëŒ€ë¡œ ë¨¼ì € ê²€ìƒ‰í•˜ì—¬ ì •í™•í•œ í‚¤ì›Œë“œ ë§¤ì¹­ì„ ìš°ì„ ìˆœìœ„ë¡œ ë‘¡ë‹ˆë‹¤.")
    
    test_query = st.text_input(
        "í…ŒìŠ¤íŠ¸ ê²€ìƒ‰ì–´",
        value="piping and instrument diagram list",
        key="two_stage_query"
    )
    
    test_filename = st.text_input(
        "ëŒ€ìƒ íŒŒì¼",
        value="ì œ4ê¶Œ ë„ë©´(ì²­ì£¼).pdf",
        key="two_stage_file"
    )
    
    if st.button("ğŸš€ 2ë‹¨ê³„ ê²€ìƒ‰ ì‹¤í–‰", type="primary"):
        st.markdown("---")
        
        # Build filter
        filter_expr = None
        if test_filename and test_filename.strip():
            filter_expr = f"search.ismatch('{test_filename}', 'metadata_storage_name')"
        
        # Stage 1: Exact search
        st.subheader("ğŸ“ Stage 1: ì •í™•í•œ í‚¤ì›Œë“œ ê²€ìƒ‰ (ì¿¼ë¦¬ í™•ì¥ ì—†ìŒ)")
        st.code(f"Query: '{test_query}'")
        
        with st.spinner("Stage 1 ê²€ìƒ‰ ì¤‘..."):
            stage1_results = search_manager.search(
                test_query,  # ì›ë³¸ ê·¸ëŒ€ë¡œ
                filter_expr=filter_expr,
                search_mode="any",
                top=50
            )
        
        st.success(f"âœ… Stage 1 ê²°ê³¼: {len(stage1_results)}ê°œ")
        
        if stage1_results:
            st.markdown("**Top 10 ê²°ê³¼:**")
            for i, doc in enumerate(stage1_results[:10], 1):
                doc_name = doc.get('metadata_storage_name', 'Unknown')
                content_snippet = doc.get('content', '')[:100].replace('\n', ' ')
                
                # Check if this is page 7
                is_page_7 = "(p.7)" in doc_name
                marker = "ğŸ¯ **[TARGET PAGE]** " if is_page_7 else ""
                
                st.markdown(f"{i}. {marker}{doc_name}")
                
                # Detailed view for page 7
                if is_page_7:
                    with st.expander("ğŸ“„ 7í˜ì´ì§€ ìƒì„¸ ë‚´ìš©"):
                        st.markdown(f"**Content Preview:**")
                        st.text_area("", doc.get('content', '')[:1000], height=200, key=f"p7_content_{i}")
        else:
            st.warning("Stage 1ì—ì„œ ê²°ê³¼ë¥¼ ì°¾ì§€ ëª»í–ˆìŠµë‹ˆë‹¤.")
        
        # Stage 2 simulation
        st.markdown("---")
        st.subheader("ğŸ“ Stage 2: ì¿¼ë¦¬ í™•ì¥ ê²€ìƒ‰ (ì°¸ê³ ìš©)")
        
        THRESHOLD = 20
        if len(stage1_results) >= THRESHOLD:
            st.info(f"â„¹ï¸ Stage 1ì—ì„œ {len(stage1_results)}ê°œ ê²°ê³¼ë¥¼ ì°¾ì•˜ìœ¼ë¯€ë¡œ Stage 2ëŠ” **ì‹¤í–‰ë˜ì§€ ì•ŠìŠµë‹ˆë‹¤** (threshold: {THRESHOLD})")
        else:
            st.warning(f"âš ï¸ Stage 1ì—ì„œ {len(stage1_results)}ê°œë§Œ ì°¾ì•˜ìœ¼ë¯€ë¡œ Stage 2 ì¿¼ë¦¬ í™•ì¥ì´ í•„ìš”í•©ë‹ˆë‹¤.")
            
            # Simulate query expansion
            expanded_query = f"{test_query} PIPING INSTRUMENT DIAGRAM LIST INDEX TABLE DRAWING"
            st.code(f"Expanded Query: '{expanded_query}'")
            
            with st.spinner("Stage 2 ê²€ìƒ‰ ì¤‘..."):
                stage2_results = search_manager.search(
                    expanded_query,
                    filter_expr=filter_expr,
                    search_mode="any",
                    top=50
                )
            
            st.success(f"âœ… Stage 2 ì¶”ê°€ ê²°ê³¼: {len(stage2_results)}ê°œ")
    
    st.markdown("---")
    st.markdown("---")
    
    # ========================================
    # ì‚¬ìš©ì ì •ì˜ ê²€ìƒ‰ ì…ë ¥
    # ========================================
    st.header("ğŸ“ ì‚¬ìš©ì ì§€ì • ê²€ìƒ‰")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        custom_query = st.text_input(
            "ê²€ìƒ‰í•  í‚¤ì›Œë“œ ì…ë ¥",
            value="piping and instrument diagram list",
            help="ê²€ìƒ‰í•˜ê³  ì‹¶ì€ í‚¤ì›Œë“œë¥¼ ì…ë ¥í•˜ì„¸ìš”"
        )
    
    with col2:
        custom_top = st.number_input(
            "ê²€ìƒ‰ ê²°ê³¼ ìˆ˜",
            min_value=1,
            max_value=200,
            value=50,
            step=10
        )
    
    filename = st.text_input(
        "ëŒ€ìƒ íŒŒì¼ëª… (ì„ íƒì‚¬í•­)",
        value="ì œ4ê¶Œ ë„ë©´(ì²­ì£¼).pdf",
        help="íŠ¹ì • íŒŒì¼ë§Œ ê²€ìƒ‰í•˜ë ¤ë©´ ì…ë ¥í•˜ì„¸ìš”. ë¹„ì›Œë‘ë©´ ì „ì²´ ì¸ë±ìŠ¤ë¥¼ ê²€ìƒ‰í•©ë‹ˆë‹¤."
    )
    
    if st.button("ğŸ” ê²€ìƒ‰ ì‹¤í–‰", type="primary", use_container_width=True):
        st.markdown("---")
        st.subheader(f"ğŸ” ê²€ìƒ‰ ê²°ê³¼: '{custom_query}'")
        
        with st.spinner("ê²€ìƒ‰ ì¤‘..."):
            # Build filter
            filter_expr = None
            if filename and filename.strip():
                filter_expr = f"search.ismatch('{filename}', 'metadata_storage_name')"
            
            # Execute search
            results = search_manager.search(
                custom_query,
                filter_expr=filter_expr,
                search_mode="any",
                top=custom_top
            )
            
            st.success(f"âœ… **{len(results)}ê°œ ê²°ê³¼ ë°œê²¬**")
            
            if len(results) == 0:
                st.warning("ê²€ìƒ‰ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤. ë‹¤ë¥¸ í‚¤ì›Œë“œë¥¼ ì‹œë„í•´ë³´ì„¸ìš”.")
            else:
                # Display results
                for i, doc in enumerate(results, 1):
                    doc_name = doc.get('metadata_storage_name', 'Unknown')
                    content = doc.get('content', '')
                    title = doc.get('title', 'No title')
                    
                    with st.expander(f"**{i}. {doc_name}**", expanded=(i <= 3)):
                        st.markdown(f"**Title**: {title}")
                        
                        # Highlight if query keywords are in content
                        content_upper = content.upper()
                        query_upper = custom_query.upper()
                        
                        # Check for keyword presence
                        keywords_found = []
                        for word in query_upper.split():
                            if word in content_upper:
                                keywords_found.append(word)
                        
                        if keywords_found:
                            st.success(f"âœ… í‚¤ì›Œë“œ ë§¤ì¹­: {', '.join(keywords_found)}")
                        
                        # Content preview
                        st.markdown("**Content Preview (ì²˜ìŒ 500ì):**")
                        st.text_area("", content[:500], height=150, key=f"custom_result_{i}", disabled=True)
                        
                        # Full content
                        with st.expander("ì „ì²´ ë‚´ìš© ë³´ê¸°"):
                            st.text_area("", content, height=400, key=f"custom_full_{i}", disabled=True)

    st.markdown("---")

    st.markdown("---")

    # ---------------------------------------------------------
    # NEW: Keyword Search Debug (Stage 1 Simulation)
    # ---------------------------------------------------------
    st.header("ğŸ” í‚¤ì›Œë“œ ê²€ìƒ‰ ë° LLM ì»¨í…ìŠ¤íŠ¸ í™•ì¸ (Keyword Search Debug)")
    st.info("LLMì´ íŠ¹ì • ì •ë³´ë¥¼ ì°¾ì§€ ëª»í•  ë•Œ, ì‹¤ì œë¡œ ê²€ìƒ‰ ì—”ì§„ì´ í•´ë‹¹ ì •ë³´ë¥¼ ì°¾ì•„ë‚´ëŠ”ì§€ í™•ì¸í•˜ëŠ” ë„êµ¬ì…ë‹ˆë‹¤.")
    
    col_debug_search, col_debug_opts = st.columns([0.7, 0.3])
    
    with col_debug_search:
        debug_keyword = st.text_input("ê²€ìƒ‰í•  í‚¤ì›Œë“œ ì…ë ¥ (ì˜ˆ: ëƒ‰ê°ìˆ˜íŒí”„ ì „ê¸°ì‹¤)", value="ëƒ‰ê°ìˆ˜íŒí”„ ì „ê¸°ì‹¤")
    
    with col_debug_opts:
        debug_top_k = st.number_input("ê²€ìƒ‰ ê°œìˆ˜ (Top K)", min_value=1, max_value=50, value=20)
    
    if st.button("ğŸš€ í‚¤ì›Œë“œ ê²€ìƒ‰ ì‹¤í–‰ (Stage 1 Logic)"):
        with st.spinner(f"'{debug_keyword}' ê²€ìƒ‰ ì¤‘..."):
            # Use exact same logic as chat_manager Stage 1
            # 1. Sanitize
            import re
            sanitized_query = re.sub(r'\bAND\b', ' ', debug_keyword, flags=re.IGNORECASE)
            sanitized_query = re.sub(r'[&+\-|!(){}\[\]^"~*?:\\]', ' ', sanitized_query)
            sanitized_query = " ".join(sanitized_query.split())
            
            st.write(f"**Sanitized Query:** `{sanitized_query}`")
            
            # 2. Search
            results = search_manager.search(
                sanitized_query,
                use_semantic_ranker=False, # Stage 1 uses standard BM25
                search_mode="all",         # Stage 1 uses AND logic
                top=debug_top_k
            )
            
            if results:
                st.success(f"âœ… ì´ {len(results)}ê°œì˜ ë¬¸ì„œë¥¼ ì°¾ì•˜ìŠµë‹ˆë‹¤.")
                
                debug_data = []
                for rank, res in enumerate(results, 1):
                    name = res.get('metadata_storage_name', 'Unknown')
                    score = res.get('@search.score', 0)
                    content = res.get('content', '')
                    
                    # Highlight keywords in content preview
                    preview = content[:300].replace('\n', ' ') + "..."
                    
                    debug_data.append({
                        "Rank": rank,
                        "Score": f"{score:.4f}",
                        "File": name,
                        "Content Preview": preview
                    })
                
                st.dataframe(pd.DataFrame(debug_data), use_container_width=True)
                
                # Detailed View
                with st.expander("ğŸ“„ ìƒì„¸ ë‚´ìš© ë³´ê¸° (Top 5)"):
                    for i, res in enumerate(results[:5], 1):
                        st.markdown(f"### {i}. {res.get('metadata_storage_name')}")
                        st.text(res.get('content', '')[:1000])
                        st.markdown("---")
            else:
                st.warning("âš ï¸ ê²€ìƒ‰ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤. (No results found)")
                st.markdown("""
                **ê°€ëŠ¥í•œ ì›ì¸:**
                1. ë¬¸ì„œì— í•´ë‹¹ í‚¤ì›Œë“œê°€ ì •í™•íˆ í¬í•¨ë˜ì–´ ìˆì§€ ì•ŠìŒ (OCR ì˜¤ë¥˜ ë“±)
                2. 'AND' ì¡°ê±´ìœ¼ë¡œ ì¸í•´ ëª¨ë“  ë‹¨ì–´ê°€ í¬í•¨ëœ ë¬¸ì„œë§Œ ê²€ìƒ‰ë¨
                """)
    
    st.markdown("---")

    # ---------------------------------------------------------
    # NEW: Target Page Debug (Why is this page missing?)
    # ---------------------------------------------------------
    st.header("ğŸ¯ íŠ¹ì • í˜ì´ì§€ ê²€ìƒ‰ ëˆ„ë½ ì›ì¸ ë¶„ì„ (Target Page Debug)")
    st.info("íŠ¹ì • í˜ì´ì§€ê°€ ê²€ìƒ‰ ê²°ê³¼ì— ë‚˜ì˜¤ì§€ ì•Šì„ ë•Œ, í•´ë‹¹ í˜ì´ì§€ê°€ ì¸ë±ìŠ¤ì— ì¡´ì¬í•˜ëŠ”ì§€, í‚¤ì›Œë“œê°€ í¬í•¨ë˜ì–´ ìˆëŠ”ì§€ ë¶„ì„í•©ë‹ˆë‹¤.")
    
    col_target_1, col_target_2, col_target_3 = st.columns([2, 2, 1])
    
    with col_target_1:
        target_query = st.text_input("ê²€ìƒ‰ ì¿¼ë¦¬", value="ëƒ‰ê°ìˆ˜íŒí”„ ì „ê¸°ì‹¤", key="target_debug_query")
    with col_target_2:
        target_filename = st.text_input("íŒŒì¼ëª… (ì •í™•íˆ ì…ë ¥)", value="ì œ5ê¶Œ ë¬¼ëŸ‰ë‚´ì—­ì„œ(ì²­ì£¼).pdf", key="target_debug_file")
    with col_target_3:
        target_page = st.number_input("í˜ì´ì§€ ë²ˆí˜¸", value=82, key="target_debug_page")
        
    if st.button("ğŸ•µï¸â€â™‚ï¸ í˜ì´ì§€ ë¶„ì„ ì‹¤í–‰"):
        target_doc_name = f"{target_filename} (p.{target_page})"
        st.write(f"**Target Document Name:** `{target_doc_name}`")
        
        # 1. Check if page exists in index
        with st.spinner("ì¸ë±ìŠ¤ì—ì„œ í˜ì´ì§€ ì¡°íšŒ ì¤‘..."):
            # Escape single quotes for OData
            safe_doc_name = target_doc_name.replace("'", "''")
            direct_check = search_manager.search(
                "*",
                filter_expr=f"metadata_storage_name eq '{safe_doc_name}'",
                top=1
            )
            
        if not direct_check:
            st.error(f"âŒ **í˜ì´ì§€ê°€ ì¸ë±ìŠ¤ì— ì—†ìŠµë‹ˆë‹¤!** (`{target_doc_name}`)")
            st.warning("íŒŒì¼ëª…ì´ë‚˜ í˜ì´ì§€ ë²ˆí˜¸ë¥¼ í™•ì¸í•´ì£¼ì„¸ìš”. ë˜ëŠ” í•´ë‹¹ íŒŒì¼ì´ ì¸ë±ì‹±ë˜ì§€ ì•Šì•˜ì„ ìˆ˜ ìˆìŠµë‹ˆë‹¤.")
        else:
            doc = direct_check[0]
            st.success(f"âœ… **í˜ì´ì§€ê°€ ì¸ë±ìŠ¤ì— ì¡´ì¬í•©ë‹ˆë‹¤.** (ID: `{doc.get('id', 'N/A')}`)")
            
            raw_content = doc.get('content', '')
            
            # Apply same cleaning logic as Chat Manager
            import re
            cleaned_content = raw_content
            
            # 1. Remove XML comments
            cleaned_content = re.sub(r'<!--.*?-->', '', cleaned_content, flags=re.DOTALL)
            
            # 2. Mark intended line breaks
            LINE_BREAK = "___LB___"
            cleaned_content = re.sub(r'</tr>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'<br\s*/?>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'</p>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'</div>', LINE_BREAK, cleaned_content, flags=re.IGNORECASE)
            
            # 3. Replace cell endings with pipe
            cleaned_content = re.sub(r'</td>', ' | ', cleaned_content, flags=re.IGNORECASE)
            cleaned_content = re.sub(r'</th>', ' | ', cleaned_content, flags=re.IGNORECASE)
            
            # 4. Remove all original newlines
            cleaned_content = cleaned_content.replace('\n', ' ').replace('\r', ' ')
            
            # 5. Remove remaining tags
            cleaned_content = re.sub(r'<[^>]+>', '', cleaned_content)
            
            # 6. Restore intended line breaks
            cleaned_content = cleaned_content.replace(LINE_BREAK, '\n')
            
            # 7. Noise
            cleaned_content = cleaned_content.replace("AutoCAD SHX Text", "").replace("%%C", "Ã˜")
            
            # 8. Collapse whitespace
            cleaned_content = re.sub(r'[ \t]+', ' ', cleaned_content)
            cleaned_content = re.sub(r'\n\s*\n', '\n\n', cleaned_content)
            cleaned_content = cleaned_content.strip()
            
            st.markdown("### ğŸ“„ í˜ì´ì§€ ë‚´ìš© (Content Preview)")
            
            tab_clean, tab_raw = st.tabs(["âœ¨ Cleaned (AIê°€ ë³´ëŠ” í™”ë©´)", "ğŸ“ Raw (ì›ë³¸ ë°ì´í„°)"])
            
            with tab_clean:
                st.info("AIì—ê²ŒëŠ” ì•„ë˜ì™€ ê°™ì´ **í‘œ êµ¬ì¡°ê°€ ì •ë¦¬ëœ í…ìŠ¤íŠ¸**ê°€ ì „ë‹¬ë©ë‹ˆë‹¤.")
                st.text_area("Cleaned Content", cleaned_content, height=400)
                
            with tab_raw:
                st.warning("ì¸ë±ìŠ¤ì— ì €ì¥ëœ ì›ë³¸ ë°ì´í„°ì…ë‹ˆë‹¤ (HTML íƒœê·¸ í¬í•¨).")
                st.text_area("Raw Content", raw_content, height=400)

            
            # 2. Analyze Keyword Matching (Check against CLEANED content)
            st.markdown("### ğŸ” í‚¤ì›Œë“œ ë§¤ì¹­ ë¶„ì„ (Cleaned Content ê¸°ì¤€)")
            keywords = target_query.split()
            
            match_data = []
            content_upper = cleaned_content.upper()
            
            all_matched = True
            for kw in keywords:
                kw_upper = kw.upper()
                count = content_upper.count(kw_upper)
                matched = count > 0
                if not matched:
                    all_matched = False
                
                match_data.append({
                    "Keyword": kw,
                    "Found": "âœ… Yes" if matched else "âŒ No",
                    "Count": count
                })
            
            st.dataframe(pd.DataFrame(match_data), use_container_width=True)
            
            if all_matched:
                st.success("âœ… ëª¨ë“  í‚¤ì›Œë“œê°€ ë³¸ë¬¸ì— í¬í•¨ë˜ì–´ ìˆìŠµë‹ˆë‹¤. ê²€ìƒ‰ ë­í‚¹ ë¬¸ì œì¼ ê°€ëŠ¥ì„±ì´ ë†’ìŠµë‹ˆë‹¤.")
            else:
                st.error("âŒ ì¼ë¶€ í‚¤ì›Œë“œê°€ ë³¸ë¬¸ì— ì—†ìŠµë‹ˆë‹¤! ì´ë˜ì„œ ê²€ìƒ‰ì´ ì•ˆ ë˜ëŠ” ê²ƒì…ë‹ˆë‹¤.")
                st.markdown("""
                **í•´ê²° ë°©ë²•:**
                1. **OCR ì˜¤ë¥˜ í™•ì¸**: ë³¸ë¬¸ í…ìŠ¤íŠ¸ë¥¼ ìì„¸íˆ ì½ì–´ë³´ì„¸ìš”. ì˜¤íƒ€ê°€ ìˆë‚˜ìš”? (ì˜ˆ: `ì „ê¸°ì‹¤` -> `ì „ ê¸° ì‹¤` or `ì „ê¸°ìŠ¬`)
                2. **ë™ì˜ì–´ í™•ì¥**: ì‚¬ìš©ìê°€ ì…ë ¥í•œ ë‹¨ì–´ì™€ ë¬¸ì„œì— ìˆëŠ” ë‹¨ì–´ê°€ ë‹¤ë¥¼ ìˆ˜ ìˆìŠµë‹ˆë‹¤.
                """)
                
            # 3. Run actual search to see Rank
            st.markdown("### ğŸ“Š ì‹¤ì œ ê²€ìƒ‰ ë­í‚¹ í™•ì¸")
            with st.spinner("ì‹¤ì œ ê²€ìƒ‰ ìˆ˜í–‰ ì¤‘..."):
                # Use same logic as Stage 1
                import re
                sanitized_query = re.sub(r'\bAND\b', ' ', target_query, flags=re.IGNORECASE)
                sanitized_query = re.sub(r'[&+\-|!(){}\[\]^"~*?:\\]', ' ', sanitized_query)
                sanitized_query = " ".join(sanitized_query.split())
                
                # Filter by filename to narrow down
                safe_filename = target_filename.replace("'", "''")
                escaped_filename = re.sub(r'([+\-&|!(){}\[\]^"~*?:\\])', r'\\\1', safe_filename)
                file_filter = f"search.ismatch('\"{escaped_filename}\"', 'metadata_storage_name')"
                
                search_results = search_manager.search(
                    sanitized_query,
                    filter_expr=file_filter,
                    use_semantic_ranker=False,
                    search_mode="all",
                    top=200
                )
                
                found_rank = None
                for i, res in enumerate(search_results, 1):
                    if res.get('metadata_storage_name') == target_doc_name:
                        found_rank = i
                        break
                
                if found_rank:
                    st.info(f"â„¹ï¸ ì´ í˜ì´ì§€ëŠ” í˜„ì¬ ê²€ìƒ‰ ê²°ê³¼ **{found_rank}ìœ„**ì— ìˆìŠµë‹ˆë‹¤.")
                else:
                    st.warning("âš ï¸ ì´ í˜ì´ì§€ëŠ” Top 200 ê²€ìƒ‰ ê²°ê³¼ì— í¬í•¨ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")

    st.markdown("---")

elif menu == "ë””ë²„ê·¸ (Debug)":
    # ========================================
    # ğŸ” ì‹¬ì¸µ ë­í‚¹ ë¶„ì„ (Deep Ranking Analysis)
    # ========================================
    st.header("ğŸ” ì‹¬ì¸µ ë­í‚¹ ë¶„ì„ (Deep Ranking Analysis)")
    st.info("ê²€ìƒ‰ì–´ì— ëŒ€í•´ ê° í˜ì´ì§€ê°€ ì™œ ê·¸ ì ìˆ˜ë¥¼ ë°›ì•˜ëŠ”ì§€ ìƒì„¸íˆ ë¶„ì„í•©ë‹ˆë‹¤.")

    col_deep_1, col_deep_2 = st.columns([2, 1])
    with col_deep_1:
        deep_query = st.text_input("ë¶„ì„í•  ê²€ìƒ‰ì–´", value="ëƒ‰ê°ìˆ˜íŒí”„ ì „ê¸°ì‹¤", key="deep_query_app")
    with col_deep_2:
        # Reuse target_filename as default if possible, otherwise generic default
        default_deep_file = ""
        deep_file = st.text_input("ëŒ€ìƒ íŒŒì¼ (í•„í„° - ì„ íƒì‚¬í•­)", value=default_deep_file, key="deep_file_input_app")

    if st.button("ğŸ”¬ ë­í‚¹ ë¶„ì„ ì‹¤í–‰", type="primary", use_container_width=True):
        st.markdown("### 1. ì¿¼ë¦¬ ë¶„ì„ (Query Analysis)")
        
        # 1. Sanitization Logic
        import re
        sanitized_query = re.sub(r'\bAND\b', ' ', deep_query, flags=re.IGNORECASE)
        sanitized_query = re.sub(r'[&+\-|!(){}\[\]^"~*?:\\]', ' ', sanitized_query)
        sanitized_query = " ".join(sanitized_query.split())
        
        st.code(f"Original: '{deep_query}'\nSanitized: '{sanitized_query}'", language="text")
        
        keywords = sanitized_query.split()
        st.write(f"**Keywords extracted:** {keywords}")
        
        # 2. Execute Search
        st.markdown("### 2. ê²€ìƒ‰ ê²°ê³¼ ë­í‚¹ (Top 20)")
        
        filter_expr = None
        if deep_file and deep_file.strip():
            filter_expr = f"search.ismatch('{deep_file}', 'metadata_storage_name')"
            
        with st.spinner("ë­í‚¹ ë¶„ì„ ì¤‘..."):
            # Ensure search_manager is available (it's initialized at top level)
            manager = get_search_manager()
            results = manager.search(
                sanitized_query,
                filter_expr=filter_expr,
                search_mode="all", # Strict mode
                top=20,
                use_semantic_ranker=False # Raw score analysis
            )
            
        if not results:
            st.warning("ê²€ìƒ‰ ê²°ê³¼ê°€ ì—†ìŠµë‹ˆë‹¤.")
        else:
            rank_data = []
            for i, doc in enumerate(results, 1):
                name = doc.get('metadata_storage_name', 'Unknown')
                content = doc.get('content', '')
                score = doc.get('@search.score', 0)
                
                # Keyword Matching Analysis
                content_upper = content.upper()
                matched_kws = []
                for kw in keywords:
                    if kw.upper() in content_upper:
                        matched_kws.append(kw)
                
                match_status = "âœ… All" if len(matched_kws) == len(keywords) else f"âš ï¸ {len(matched_kws)}/{len(keywords)}"
                
                # Highlight specific pages
                highlight = ""
                if "(p.17)" in name: highlight = "ğŸ”´ (Issue)"
                if "(p.82)" in name: highlight = "ğŸŸ¢ (Target)"
                
                rank_data.append({
                    "Rank": i,
                    "Score": f"{score:.4f}",
                    "Page": f"{name} {highlight}",
                    "Match": match_status,
                    "Matched Keywords": ", ".join(matched_kws),
                    "Snippet": content[:100].replace("\n", " ") + "..."
                })
                
            st.dataframe(pd.DataFrame(rank_data), use_container_width=True)
            
            # Detailed Comparison
            st.markdown("### 3. ì£¼ìš” í˜ì´ì§€ ìƒì„¸ ë¹„êµ")
            target_pages = [d for d in results if "(p.82)" in d.get('metadata_storage_name', '') or "(p.17)" in d.get('metadata_storage_name', '')]
            
            if target_pages:
                for doc in target_pages:
                    name = doc.get('metadata_storage_name')
                    score = doc.get('@search.score')
                    st.markdown(f"#### ğŸ“„ {name} (Score: {score:.4f})")
                    st.text_area(f"Content of {name}", doc.get('content', ''), height=200)
            else:
                st.info("ë¹„êµí•  ì£¼ìš” í˜ì´ì§€(p.17, p.82)ê°€ Top 20 ë‚´ì— ì—†ìŠµë‹ˆë‹¤.")
